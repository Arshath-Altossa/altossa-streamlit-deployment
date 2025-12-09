import os
import unicodedata
import requests
import pandas as pd
import streamlit as st
from bs4 import BeautifulSoup
from io import BytesIO
from typing import Optional, Any
import re
from sqlalchemy import create_engine, text
#from concurrent.futures import ThreadPoolExecutor, as_completed
from dotenv import load_dotenv, find_dotenv
import time
from datetime import datetime
import hashlib
import uuid
# NEW: persistent storage for attachments (disk-backed)
import pathlib, shutil
import threading
import bcrypt

# --- NEW: persistent cart autosave/restore (disk-backed; per-app) ---
import json
from pathlib import Path
from hashlib import sha1

from PIL import Image, ImageDraw, ImageFont
import textwrap
import io

# [ADD THESE IMPORTS]
from reportlab.lib import colors
from reportlab.lib.pagesizes import A4, landscape
# [UPDATED] - Ensure these imports are present
from reportlab.platypus import (
    SimpleDocTemplate, PageTemplate, Frame, Paragraph, Spacer, 
    Image as RLImage, Table, TableStyle, NextPageTemplate, PageBreak
)
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import cm
from reportlab.lib.utils import ImageReader
from reportlab.lib.colors import black, HexColor

st.set_page_config(page_title="Altossa Selector", layout="wide")

# Load .env (walks up from current working dir)
load_dotenv(dotenv_path=find_dotenv(usecwd=True), override=True)

# [NEW] Helper to read secrets from Streamlit Cloud or local env
def get_secret(name: str, default: str = "") -> str:
    """Try st.secrets first, fallback to os.getenv."""
    try:
        if name in st.secrets:
            return str(st.secrets[name])
    except Exception:
        pass
    return os.getenv(name, default)

# Remember app-auth across Streamlit sessions for the current APP_PASSWORD
_APP_AUTH_CACHE: dict[str, bool] = {}

@st.cache_resource
def db_engine(database_url: str):
    """Reuse a single SQLAlchemy engine across reruns (faster + fewer connections)."""
    return create_engine(database_url, pool_pre_ping=True)

def get_database_url() -> str:
    """Secrets-first, then Env-first."""
    return get_secret("DATABASE_URL")

# ---------------------------
# Defaults: Room → Product Types (editable in UI)
# ---------------------------
DEFAULT_ROOM_TYPES = {
    "Drawing Room": ["Sofa", "Armchair", "Coffee Table","Carpets"],
    "Living Room": ["Sofa", "Armchair", "Table", "Coffee Table", "Tv Holder"],
    "Dining Room": ["Chair", "Table"],
    "Kitchen": ["Modular Kitchen", "Trolly/Slider", "Hob-Hood-Sink-Tap"],
    "Bedroom 1": ["Bed", "Night Stand", "Bed - Side Table", "Chest of Drawers", "Wardrobes"],
    "Bedroom 2": ["Bed", "Night Stand", "Bed - Side Table", "Chest of Drawers", "Wardrobes"],
    "Bedroom 3": ["Bed", "Night Stand", "Bed - Side Table", "Chest of Drawers", "Wardrobes"],
    "Others": ["Armchair - Cushion", "Pillow"],
}

ROOM_CHOICES = [
    "Drawing Room",
    "Living Room",
    "Dining Room",
    "Kitchen",
    "Bedroom 1",
    "Bedroom 2",
    "Bedroom 3",
    "Others",
]

# Placeholder label used for single-quantity room selector
ROOM_PLACEHOLDER = "Select your room"

# Placeholder label used for Drawing Room visual-role selector
VISUAL_ROLE_PLACEHOLDER = "Select the label"

IMG_CACHE_MAX = 300  # keep last ~300 image entries in session cache
IMG_CACHE_TTL_SECONDS = 24 * 60 * 60  # 24 hours

@st.cache_resource
def _get_memory_image_store() -> dict:
    """
    Returns a persistent dictionary {key: (bytes, timestamp)} shared across sessions/reruns.
    Replaces the SQLite file for Streamlit Cloud compatibility.
    """
    return {}

# ---------------------------
# Utils
# ---------------------------
def norm_txt(s: str) -> str:
    if s is None:
        return ""
    s = unicodedata.normalize("NFKD", str(s)).encode("ascii", "ignore").decode("ascii")
    s = s.strip()
    return s

def to_num(x):
    """
    Dimensions parser:
      - Pure number -> float
      - Any value that looks like a range or multi-choice
        (e.g. '75-97', '75/97', '75 to 97', '25/30/40/50', '30/40/50')
        -> keep original string (unchanged)
      - Otherwise -> None
    """
    import re
    import unicodedata
    import pandas as pd

    if pd.isna(x):
        return None

    s_orig = str(x).strip()
    if s_orig == "":
        return None

    # ASCII-normalize only for pattern checks
    s = unicodedata.normalize("NFKD", s_orig).encode("ascii", "ignore").decode("ascii").strip()

    # If it contains separators between digits ( - / en/em dash or "to" ),
    # treat it as a text range / multi-size and KEEP the original string.
    # This catches: '75-97', '260/330', '25/30/40/50', '30/40/50', '75 to 97', etc.
    if re.search(r"\d\s*(?:-|/|–|—|\bto\b)\s*\d", s, flags=re.IGNORECASE):
        return s_orig

    # Pure number (allow thousands commas or decimal comma/dot)
    num_candidate = s.replace(",", ".")
    if re.match(r"^\s*\d+(?:\.\d+)?\s*$", num_candidate):
        try:
            return float(num_candidate)
        except Exception:
            return None

    return None

def calculate_final_price(brand, base_price):
    """Brand-based final price.
    Returns None if price is missing/NaN so callers can render '—' safely.
    """
    try:
        p = float(str(base_price).strip())
    except Exception:
        return None
    # guard against NaN (float('nan') != float('nan'))
    if p != p:  # NaN check
        return None

    b = (brand or "").strip().lower()
    if "cattelan italia" in b:
        mult = 120
    elif "ditre italia" in b:
        mult = 100
    elif "pianca" in b:
        mult = 100
    else:
        mult = 100
    try:
        return int(round(p * mult))
    except Exception:
        return None

def _get_item_price(item: dict) -> float:
    """
    Return the effective unit price for an item.
    Priority: 1. Session override, 2. Calculated DB price.
    """
    if "override_price" in item:
        try:
            val = float(item["override_price"])
            if val >= 0: return val
        except Exception:
            pass
            
    val = calculate_final_price(item.get("Brand", ""), item.get("Base Price"))
    return float(val) if val is not None else 0.0

# NEW: safe string for values that may be None/NaN/<NA>
def nz_str(x: Any) -> str:
    """Return empty string for None/NaN/<NA>, else str(x)."""
    try:
        return "" if x is None or pd.isna(x) else str(x)
    except Exception:
        # fallback if pd.isna() can't handle the type
        return "" if x is None else str(x)

def _cart_to_dataframe(cart):
    """Build a tidy DataFrame from cart with computed unit/subtotals + composition fields.
    Also normalize to a single 'Category' column for export.
    """
    rows = []
    for it in cart:
        unit_price = _get_item_price(it)  # [UPDATED]
        qty = int(it.get("qty", 1) or 1)
        subtotal = unit_price * qty       # [UPDATED]

        # Single 'Category' export field (prefer Fabric/Leather label else Other Category)
        cat = (it.get("Material") or "").strip() or (it.get("Other Category") or "").strip()

        # NEW: human-readable room allocation text
        room_text = _room_allocation_text(it)

        rows.append({
            "Brand": it.get("Brand", ""),
            "Product type": it.get("Product type", ""),
            "Product name": it.get("Product name", ""),
            "Name of subtype": it.get("Name of subtype", ""),
            "W": it.get("W", ""),
            "D": it.get("D", ""),
            "H": it.get("H", ""),
            "Category": cat,
            "Rooms": room_text,                 # NEW column
            "Qty": qty,
            "Unit Price (INR)": unit_price,     # no currency symbol, header indicates INR
            "Subtotal (INR)": subtotal,         # no currency symbol, header indicates INR
            "Link": it.get("Link", ""),         # product page (hero fallback)
            "Composition parts": it.get("Composition parts", ""),
            "Composition link": it.get("Composition link", ""),  # preferred hero if present
        })
    df = pd.DataFrame(rows)
    if not df.empty:
        df = df[
            [
                "Brand","Product type","Product name","Name of subtype","W","D","H",
                "Category","Rooms","Qty","Unit Price (INR)","Subtotal (INR)",
                "Composition parts","Composition link","Link",
            ]
        ]
    return df

def _room_qty(val) -> int:
    """
    Normalize a stored room-allocation value into an int quantity.
    Supports both plain ints and dicts like {'qty': 2, 'visual_role': 'Sofa'}.
    """
    if isinstance(val, dict):
        base = val.get("qty", val.get("quantity", 0))
    else:
        base = val
    try:
        return int(base or 0)
    except Exception:
        return 0

def _room_allocation_text(item: dict) -> str:
    """
    Convert item['room_alloc'] into a human-readable text like:
      'Drawing Room × 2 ; Bedroom × 1'
    Falls back to 'Others × <qty>' if nothing is set.
    """
    qty = 1
    try:
        qty = int(item.get("qty", 1) or 1)
    except Exception:
        qty = 1

    raw_alloc = item.get("room_alloc") or {}
    alloc: dict[str, int] = {}

    # Normalize to int counts > 0, supporting both int and dict values
    if isinstance(raw_alloc, dict):
        for k, v in raw_alloc.items():
            c = _room_qty(v)
            if c > 0:
                alloc[str(k)] = c

    if not alloc:
        return f"Others × {qty}"

    parts: list[str] = []

    # Stable order: known rooms first
    for room in ROOM_CHOICES:
        c = int(alloc.get(room, 0) or 0)
        if c > 0:
            parts.append(f"{room} × {c}")

    # Then any custom room labels (if you ever add them later)
    for room, c in alloc.items():
        if room not in ROOM_CHOICES and c > 0:
            parts.append(f"{room} × {c}")

    return " ; ".join(parts) if parts else f"Others × {qty}"

def get_room_allocation_violations(cart: list[dict]) -> list[str]:
    """
    For multi-quantity items, check that the sum of room allocations
    equals the item qty.

    Returns a list of human-readable violation messages, e.g.:
      'Alta Armchair · Armchair · Ditre Italia: Qty 2, assigned 3 across rooms'
    """
    violations: list[str] = []

    for it in cart:
        # Skip items without quantity or with qty <= 1 (single piece is always fine)
        try:
            qty = int(it.get("qty", 1) or 1)
        except Exception:
            qty = 1

        if qty <= 1:
            continue

        raw_alloc = it.get("room_alloc") or {}
        total_assigned = 0

        if isinstance(raw_alloc, dict):
            for _, v in raw_alloc.items():
                total_assigned += _room_qty(v)

        if total_assigned != qty:
            pname = str(it.get("Product name") or "").strip() or "Unnamed"
            ptype = str(it.get("Product type") or "").strip()
            brand = str(it.get("Brand") or "").strip()
            label = " · ".join([x for x in [pname, ptype, brand] if x])

            violations.append(
                f"{label}: Qty {qty}, assigned {total_assigned if total_assigned >= 0 else '?'} across rooms"
            )

    return violations

def _excel_bytes_from_df(df: pd.DataFrame) -> bytes:
    """
    Return an .xlsx file in-memory.

    Changes:
    - Adds a new visible column **Product Image** (hero image from Product page `Link`).
    - Keeps **Composition Image** as-is (from `Composition link`), but sizes both image columns
      so thumbnails fit the column without overlapping cells.
    - Auto-sets row heights where images are inserted.
    - Totals row remains.
    """
    output = BytesIO()

    # ----- Prep dataframe (add "Product Image" & keep "Composition Image" if comp cols exist) -----
    df_out = df.copy()

    has_link_col = "Link" in df_out.columns
    has_comp_col = "Composition link" in df_out.columns

    # Insert "Product Image" immediately after "Link" (if present), otherwise append at end.
    if has_link_col and "Product Image" not in df_out.columns:
        df_out.insert(df_out.columns.get_loc("Link") + 1, "Product Image", "")
    elif "Product Image" not in df_out.columns:
        df_out["Product Image"] = ""

    # Keep composition image column next to comp link if available.
    if has_comp_col and "Composition Image" not in df_out.columns:
        df_out.insert(df_out.columns.get_loc("Composition link") + 1, "Composition Image", "")

    # ---------- Preferred path: xlsxwriter (with proper sizing) ----------
    try:
        with pd.ExcelWriter(output, engine="xlsxwriter") as writer:
            df_out.to_excel(writer, index=False, sheet_name="Selection")
            ws = writer.sheets["Selection"]

            # Column sizing: make both image columns wider so thumbnails are clearly visible
            # NOTE: xlsxwriter width units ~= Excel "character" width; ~1 unit ≈ 7px
            def _set_col_width(col_label: str, width_chars: float = 24.0):
                if col_label in df_out.columns:
                    col_idx = df_out.columns.get_loc(col_label)
                    ws.set_column(col_idx, col_idx, width_chars)

            _set_col_width("Product Image", 24.0)
            _set_col_width("Composition Image", 24.0)

            # Helper to place a scaled thumbnail that fits within ~150×110 px cell box
            def _insert_thumb(r_idx_1based: int, c_idx_0based: int, img_bytes: Optional[bytes]) -> None:
                if not img_bytes:
                    return
                try:
                    # set row height tall enough for thumbnails (in points; ~1pt = 1/72")
                    ws.set_row(r_idx_1based - 1, 110)
                    ws.insert_image(
                        r_idx_1based - 1,
                        c_idx_0based,
                        "ignored.png",
                        {
                            "image_data": BytesIO(img_bytes),
                            "x_offset": 2,
                            "y_offset": 2,
                            # scale chosen to fit comfortably in the 24-char column and 110pt row
                            "x_scale": 0.55,
                            "y_scale": 0.55,
                        },
                    )
                except Exception:
                    pass

            # Insert Product Image thumbnails (prefer product hero from Link)
            if has_link_col and "Product Image" in df_out.columns:
                prod_img_col = df_out.columns.get_loc("Product Image")
                for r_idx, link in enumerate(df_out["Link"], start=2):
                    url = str(link or "").strip()
                    if not url:
                        continue
                    img = get_image_if_cached(url, "") or get_image_cached(url, "")
                    _insert_thumb(r_idx, prod_img_col, img)

            # Insert Composition Image thumbnails
            if has_comp_col and "Composition Image" in df_out.columns:
                comp_img_col = df_out.columns.get_loc("Composition Image")
                for r_idx, link in enumerate(df_out["Composition link"], start=2):
                    url = str(link or "").strip()
                    if not url:
                        continue
                    img = get_image_if_cached("", url) or get_image_cached("", url)
                    _insert_thumb(r_idx, comp_img_col, img)

            # Totals row (unchanged)
            if not df_out.empty and "Subtotal (INR)" in df_out.columns:
                last_row = len(df_out) + 1  # 1-indexed with header
                name_col = df_out.columns.get_loc("Product name")
                ws.write(last_row, name_col, "TOTAL (INR)")
                total_col = df_out.columns.get_loc("Subtotal (INR)")
                col_letter = chr(65 + total_col)
                ws.write_formula(last_row, total_col, f"=SUM({col_letter}2:{col_letter}{len(df_out)+1})")

        output.seek(0)
        return output.read()

    except ImportError:
        pass  # fallback to openpyxl path below

    # ---------- Fallback path: openpyxl (with optional images if Pillow is present) ----------
    try:
        from openpyxl import Workbook
        from openpyxl.utils.dataframe import dataframe_to_rows
        try:
            from openpyxl.drawing.image import Image as XLImage  # type: ignore
            _can_embed_images = True
        except Exception:
            _can_embed_images = False

        wb = Workbook()
        ws = wb.active
        ws.title = "Selection"

        # Write headers + data
        for r_idx, row in enumerate(dataframe_to_rows(df_out, index=False, header=True), start=1):
            ws.append(row)
            if r_idx == 1:
                continue
            # make rows taller when we plan to add thumbnails
            if _can_embed_images:
                ws.row_dimensions[r_idx].height = 84  # ~84pt ~= 112px

        # Place images if possible
        def _add_pic(row_1based: int, col_1based: int, img_bytes: Optional[bytes]):
            if not (_can_embed_images and img_bytes):
                return
            try:
                bio = BytesIO(img_bytes)
                img = XLImage(bio)
                # openpyxl does not support scaling perfectly to cell; rely on row height we set
                ws.add_image(img, ws.cell(row=row_1based, column=col_1based).coordinate)
            except Exception:
                pass

        if has_link_col and "Product Image" in df_out.columns:
            c_idx = df_out.columns.get_loc("Product Image") + 1
            for r_idx, link in enumerate(df_out["Link"], start=2):
                url = str(link or "").strip()
                if not url:
                    continue
                img = get_image_if_cached(url, "") or get_image_cached(url, "")
                _add_pic(r_idx, c_idx, img)

        if has_comp_col and "Composition Image" in df_out.columns:
            c_idx = df_out.columns.get_loc("Composition Image") + 1
            for r_idx, link in enumerate(df_out["Composition link"], start=2):
                url = str(link or "").strip()
                if not url:
                    continue
                img = get_image_if_cached("", url) or get_image_cached("", url)
                _add_pic(r_idx, c_idx, img)

        # Totals (plain value)
        if not df_out.empty and "Subtotal (INR)" in df_out.columns:
            last_row = ws.max_row + 1
            name_col = df_out.columns.get_loc("Product name") + 1
            sub_col  = df_out.columns.get_loc("Subtotal (INR)") + 1
            ws.cell(row=last_row, column=name_col, value="TOTAL (INR)")
            try:
                total_val = float(pd.to_numeric(df_out["Subtotal (INR)"], errors="coerce").fillna(0).sum())
            except Exception:
                total_val = 0.0
            ws.cell(row=last_row, column=sub_col, value=total_val)

        tmp = BytesIO()
        wb.save(tmp)
        tmp.seek(0)
        return tmp.read()

    except Exception:
        # Last-resort: no images
        with pd.ExcelWriter(output, engine="openpyxl") as writer:
            df_out.to_excel(writer, index=False, sheet_name="Selection")
        output.seek(0)
        return output.read()

def _pdf_bytes_from_df(df: pd.DataFrame, client_name: str = "", approved_by: str = "") -> Optional[bytes]:
    """
    Landscape A4 pro-forma PDF:
      • Altossa logo ONLY on the first page.
      • A single header box below the logo (first page only).
      • Each Product-type section is kept together (title + table + subtotal row INSIDE the table).
      • Table widths guaranteed to fit the printable area (no clipping).
      • Row image is *contained* and *scaled* to the image-cell box (no overflow / no crop).
      • Clean grand-total block that aligns with the tables.
    """
    try:
        from reportlab.lib import colors
        from reportlab.lib.pagesizes import A4, landscape
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.enums import TA_RIGHT
        from reportlab.platypus import (
            SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer, Image, KeepTogether
        )
        from reportlab.lib.utils import ImageReader
        from datetime import datetime
    except Exception:
        return None

    # ---- Prepare the logo once
    ALTOSSA_LOGO = "https://drive.google.com/file/d/1eK3VDb7gemhJTxenmIxYrl8lauw5HBLH/view?usp=sharing"
    _logo_bytes = get_image_if_cached("", ALTOSSA_LOGO) or get_image_cached("", ALTOSSA_LOGO)
    _logo_reader = ImageReader(BytesIO(_logo_bytes)) if _logo_bytes else None

    # ---- Page setup
    buffer = BytesIO()
    page_size = landscape(A4)  # 842 × 595 pt
    doc = SimpleDocTemplate(
        buffer,
        pagesize=page_size,
        leftMargin=18, rightMargin=18,   # content width ≈ 842 - 36 = 806 pt
        topMargin=42, bottomMargin=18,
    )
    content_width = page_size[0] - doc.leftMargin - doc.rightMargin  # ~806

    styles = getSampleStyleSheet()
    base = styles["Normal"]; base.fontSize = 8; base.leading = 10
    wrap_desc  = ParagraphStyle("wrap_desc",  parent=base, fontSize=8, leading=10, wordWrap="CJK")
    wrap_parts = ParagraphStyle("wrap_parts", parent=base, fontSize=8, leading=10, wordWrap="CJK")
    wrap_cat   = ParagraphStyle("wrap_cat",   parent=base, fontSize=8, leading=10, wordWrap="CJK")
    right_num  = ParagraphStyle("right_num",  parent=base, alignment=TA_RIGHT)

    # ---- First page header (logo only on page 1)
    def _draw_first_header(canvas, _doc):
        if _logo_reader:
            pw, ph = page_size
            lw, lh = _logo_reader.getSize()
            max_w, max_h = 240, 40
            scale = min(max_w / float(lw), max_h / float(lh))
            dw, dh = lw * scale, lh * scale
            x = (pw - dw) / 2.0
            y = ph - dh - 8
            canvas.drawImage(_logo_reader, x, y, width=dw, height=dh,
                             preserveAspectRatio=True, mask='auto')

    def _no_header(canvas, _doc):
        return

    elems = []

    # ---- Pro-forma header box (front-page flow)
    hdr_data = [
        ["Client", Paragraph(f"<b>{(client_name or '').strip() or '-'}</b>", base),
         "Date & Time", Paragraph(datetime.now().strftime("%d %b %Y, %I:%M %p"), base),
         "Approved by", Paragraph(f"<b>{(approved_by or '').strip() or '-'}</b>", base)],
    ]
    hdr = Table(hdr_data, colWidths=[60, 180, 80, 160, 80, 160])
    hdr.setStyle(TableStyle([
        ("BOX", (0,0), (-1,-1), 0.8, colors.HexColor("#0f1b3d")),
        ("INNERGRID", (0,0), (-1,-1), 0.3, colors.HexColor("#0f1b3d")),
        ("BACKGROUND", (0,0), (-1,-1), colors.HexColor("#f6f7fb")),
        ("FONTNAME", (0,0), (-1,-1), "Helvetica"),
        ("FONTSIZE", (0,0), (-1,-1), 9),
        ("VALIGN", (0,0), (-1,-1), "MIDDLE"),
        ("LEFTPADDING", (0,0), (-1,-1), 6),
        ("RIGHTPADDING",(0,0), (-1,-1), 6),
        ("TOPPADDING",  (0,0), (-1,-1), 4),
        ("BOTTOMPADDING",(0,0), (-1,-1), 4),
    ]))
    elems.append(hdr)
    elems.append(Spacer(1, 10))

    if df.empty:
        elems.append(Paragraph("No items selected.", base))
        doc.build(elems, onFirstPage=_draw_first_header, onLaterPages=_no_header)
        buffer.seek(0);  return buffer.read()

    # ---- Helpers (tightened image handling)
    CELL_PAD = 3.5  # must match table paddings below

    # NOTE: Total column widths MUST fit content_width (~806 pt).
    # We now have 12 columns:
    #   idx:   0     1     2   3   4    5      6      7      8    9     10     11
    #         Br   Desc    W   D   H  Cat   Rooms  Comp   Qty  Unit  Sub    Img
    # Sum ≈ 780 pt → still safely inside content_width.
    _COL_WIDTHS = [50, 170, 26, 26, 26, 78, 60, 118, 24, 66, 76, 70]

    IMG_CELL_W = _COL_WIDTHS[10] - 2 * CELL_PAD
    IMG_CELL_H = 58  # keep rows compact & consistent

    def _image_for_row(comp_url: str, prod_link: str) -> Optional[bytes]:
        comp_url = (comp_url or "").strip()
        prod_link = (prod_link or "").strip()
        img = None
        if comp_url:
            img = get_image_if_cached("", comp_url) or get_image_cached("", comp_url)
        if not img and prod_link:
            img = get_image_if_cached(prod_link, "") or get_image_cached(prod_link, "")
        return img

    def _scaled_image_flowable(img_bytes: Optional[bytes]) -> Any:
        """
        Return an Image scaled to *fit* the fixed image-box (IMG_CELL_W × IMG_CELL_H).
        This guarantees no overflow and no cropping.
        """
        if not img_bytes:
            return ""
        try:
            rdr = ImageReader(BytesIO(img_bytes))
            iw, ih = rdr.getSize()
            if iw <= 0 or ih <= 0:
                return ""
            scale = min(IMG_CELL_W / float(iw), IMG_CELL_H / float(ih))
            w, h = iw * scale, ih * scale
            return Image(BytesIO(img_bytes), width=w, height=h)  # already fitted, no KeepInFrame needed
        except Exception:
            return ""

    def _render_type_section(type_name: str, gdf: pd.DataFrame) -> int:
        gdf = gdf.copy()
        gdf["Description"] = gdf["Product name"].astype(str) + " — " + gdf["Name of subtype"].astype(str)

        headers = [
            "Brand","Description","W","D","H",
            "Category","Rooms","Composition parts",
            "Qty","Unit Price (INR)","Subtotal (INR)","Image",
        ]

        data = [headers]
        for _, r in gdf.iterrows():
            desc_p  = Paragraph(str(r.get("Description","")), wrap_desc)
            cat_p   = Paragraph(str(r.get("Category","")),   wrap_cat)
            rooms_p = Paragraph(str(r.get("Rooms","")),      wrap_cat)
            parts_p = Paragraph(str(r.get("Composition parts","")), wrap_parts)
            qty_p   = Paragraph(str(int(r.get("Qty", 1) or 1)), right_num)
            unit_p  = Paragraph(str(int(r.get("Unit Price (INR)") or 0)), right_num)
            sub_p   = Paragraph(str(int(r.get("Subtotal (INR)") or 0)),   right_num)

            img_bytes = _image_for_row(str(r.get("Composition link") or ""), str(r.get("Link") or ""))
            img_flow  = _scaled_image_flowable(img_bytes)

            data.append([
                r.get("Brand",""),
                desc_p,
                r.get("W",""),
                r.get("D",""),
                r.get("H",""),
                cat_p,
                rooms_p,
                parts_p,
                qty_p,
                unit_p,
                sub_p,
                img_flow,
            ])

        # Subtotal (inside the SAME table as a styled last row)
        subtotal_val = int(gdf["Subtotal (INR)"].fillna(0).sum())
        subtotal_label = Paragraph(f"<b>Subtotal – {type_name}:</b>", base)
        subtotal_amt   = Paragraph(f"<b>{subtotal_val:,} INR</b>", right_num)

        # Append a row and span label across first 10 columns (0..9),
        # amount in col 10, image col 11 blank.
        data.append([""] * len(headers))
        subtotal_row_idx = len(data) - 1

        tbl = Table(
            data,
            colWidths=_COL_WIDTHS,
            repeatRows=1,
            splitByRow=1,
            hAlign="LEFT",
        )

        # Base styling
        style_cmds = [
            ("BACKGROUND", (0,0), (-1,0), colors.HexColor("#0f1b3d")),
            ("TEXTCOLOR", (0,0), (-1,0), colors.white),
            ("FONTNAME", (0,0), (-1,0), "Helvetica"),
            ("FONTSIZE", (0,0), (-1,-1), 8),
            ("GRID", (0,0), (-1,-1), 0.25, colors.HexColor("#bfbfbf")),
            ("VALIGN", (0,1), (-1,-1), "MIDDLE"),
            ("ALIGN", (8,1), (10,-1), "RIGHT"),   # Qty, Unit, Subtotal
            ("LEFTPADDING", (0,0), (-1,-1), CELL_PAD),
            ("RIGHTPADDING",(0,0), (-1,-1), CELL_PAD),
            ("TOPPADDING",  (0,0), (-1,-1), 2.5),
            ("BOTTOMPADDING",(0,0), (-1,-1), 2.5),
        ]

        # Subtotal row style: span label, right-align amount, subtle background
        style_cmds += [
            ("SPAN", (0, subtotal_row_idx), (9, subtotal_row_idx)),  # label span
            ("BACKGROUND", (0, subtotal_row_idx), (-1, subtotal_row_idx), colors.HexColor("#f1f2f6")),
            ("FONTNAME", (0, subtotal_row_idx), (-1, subtotal_row_idx), "Helvetica-Bold"),
            ("ALIGN", (10, subtotal_row_idx), (10, subtotal_row_idx), "RIGHT"),
        ]
        tbl.setStyle(TableStyle(style_cmds))

        # Put the actual subtotal label & value into the spanned cells
        tbl._cellvalues[subtotal_row_idx][0] = subtotal_label
        tbl._cellvalues[subtotal_row_idx][10] = subtotal_amt
        # image col (11) remains blank for subtotal row

        section = [
            Paragraph(f"<b>{type_name}</b>", styles["Heading3"]),
            tbl,
            Spacer(1, 10),
        ]
        elems.append(KeepTogether(section))
        return subtotal_val

    # ---- Render sections grouped by Product type
    grand_total = 0
    for ptype, g in df.groupby("Product type", sort=True):
        grand_total += _render_type_section(str(ptype), g)

    # ---- Grand total block (exactly aligns to the same content_width)
    elems.append(Spacer(1, 6))
    gt = Table(
        [[Paragraph("<b>Grand Total:</b>", styles["Heading3"]),
          Paragraph(f"{grand_total:,} INR", ParagraphStyle('gt_right', parent=styles["Heading3"], alignment=TA_RIGHT))]],
        colWidths=[content_width - 220, 220],  # total exactly equals content_width
        hAlign="LEFT",
        style=TableStyle([
            ("BOX", (0,0), (-1,-1), 0.6, colors.HexColor("#0f1b3d")),
            ("BACKGROUND", (0,0), (-1,-1), colors.HexColor("#f6f7fb")),
            ("LEFTPADDING", (0,0), (-1,-1), 8),
            ("RIGHTPADDING",(0,0), (-1,-1), 8),
            ("TOPPADDING",  (0,0), (-1,-1), 6),
            ("BOTTOMPADDING",(0,0), (-1,-1), 6),
        ])
    )
    elems.append(gt)

    # Build: logo only on first page
    doc.build(elems, onFirstPage=_draw_first_header, onLaterPages=_no_header)
    buffer.seek(0)
    return buffer.read()

# --- NEW: Internal PPT export (one slide per selected product) ---
def _ppt_bytes_from_cart(cart: list[dict], client_name: str = "") -> bytes | None:
    """
    16:9 clean slide per item:
      • Title:  Product type · Product name · Brand
      • Left:   main image (composition preferred; else product hero)
      • Right-top:  Client Notes / Adjustments (in the open whitespace)
      • Bottom: table with visible header row (Subtype, Size, Material/Other, Qty, Unit, Subtotal, Brand, Product Page)
      • Bottom-right BELOW table: up to 3 attached images (so they never overlap the table)
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor
    except Exception:
        return None

    def _unit_and_subtotal(it: dict) -> tuple[int | None, int]:
        unit = calculate_final_price(it.get("Brand", ""), it.get("Base Price"))
        qty = int(it.get("qty", 1) or 1)
        return unit, (int(unit) * qty if unit is not None else 0)

    def _img_bytes_for_item(it: dict) -> bytes | None:
        comp = (it.get("Composition link") or "").strip()
        link = (it.get("Link") or "").strip()
        img = None
        if comp:
            img = get_image_if_cached("", comp) or get_image_cached("", comp)
        if not img and link:
            img = get_image_if_cached(link, "") or get_image_cached(link, "")
        return img

    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    for it in cart:
        sld = prs.slides.add_slide(blank)

        # Title: "Type · Name · Brand"
        title = f"{it.get('Product type','')} · {it.get('Product name','')} · {it.get('Brand','')}"
        tx = sld.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12.3), Inches(0.7))
        tf = tx.text_frame; tf.clear()
        p = tf.paragraphs[0]; p.text = title; p.font.size = Pt(28); p.font.bold = True

        # Left: main image (height-first fit)
        img = _img_bytes_for_item(it)
        if img:
            from io import BytesIO
            try:
                sld.shapes.add_picture(BytesIO(img), Inches(0.5), Inches(1.2), height=Inches(4.6))
            except Exception:
                pass

        # Right-top: notes in the open whitespace (no collision with image)
        comment = (it.get("internal_comment") or "").strip()
        notes = sld.shapes.add_textbox(Inches(7.3), Inches(1.2), Inches(5.45), Inches(2.2))
        nt = notes.text_frame; nt.word_wrap = True; nt.clear()
        h = nt.paragraphs[0]; h.text = "Client Notes / Adjustments"; h.font.bold = True; h.font.size = Pt(16)
        b = nt.add_paragraph(); b.text = comment or "—"; b.font.size = Pt(13)

        # Bottom table: visible header row
        subtype = it.get("Name of subtype") or ""
        dims = " × ".join([v for v in [str(it.get("W") or ""), str(it.get("D") or ""), str(it.get("H") or "")] if v])
        material = (it.get("Material") or it.get("Other Category") or "")
        qty = int(it.get("qty", 1) or 1)
        unit_price, subtotal = _unit_and_subtotal(it)
        link  = it.get("Link") or ""

        # NEW: text summary of room allocation
        rooms_txt = _room_allocation_text(it)

        tbl = sld.shapes.add_table(
            rows=2, cols=9,
            left=Inches(0.5), top=Inches(5.2),
            width=Inches(12.3), height=Inches(1.4)
        ).table

        headers = [
            "Subtype",
            "Size",
            "Material/Other",
            "Rooms",
            "Qty",
            "Unit Price",
            "Subtotal",
            "Brand",
            "Product Page",
        ]
        vals = [
            subtype,
            dims,
            material,
            rooms_txt,
            str(qty),
            (f"₹{unit_price:,}" if unit_price is not None else "—"),
            f"₹{subtotal:,}",
            (it.get("Brand") or ""),
            link,
        ]

        # Column widths sum to ≈ 12.3"
        col_w = [1.4, 1.7, 1.9, 1.4, 0.7, 1.0, 1.2, 1.0, 2.0]
        for i, w in enumerate(col_w):
            tbl.columns[i].width = Inches(w)

        # Header styling (blue background, white text, bold)
        for c, text in enumerate(headers):
            hdr_p = tbl.cell(0, c).text_frame.paragraphs[0]
            hdr_p.text = text
            hdr_p.font.bold = True
            hdr_p.font.size = Pt(12)
            cell = tbl.cell(0, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(16, 46, 92)  # deep blue
            hdr_p.font.color.rgb = RGBColor(255, 255, 255)

        # Body row
        # Header styling (unchanged)
        for c, text in enumerate(headers):
            hdr_p = tbl.cell(0, c).text_frame.paragraphs[0]
            hdr_p.text = text
            hdr_p.font.bold = True
            hdr_p.font.size = Pt(12)
            cell = tbl.cell(0, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(16, 46, 92)  # deep blue
            hdr_p.font.color.rgb = RGBColor(255, 255, 255)

        # Body row
        for c, text in enumerate(vals):
            par = tbl.cell(1, c).text_frame.paragraphs[0]
            par.text = str(text)
            par.font.size = Pt(11)
            # Qty, Unit Price, Subtotal are now indices 4, 5, 6
            if c in (4, 5, 6):
                par.alignment = PP_ALIGN.RIGHT

        # Bottom-right: attached images placed *below* the table to avoid overlap
        assets: list[bytes] = it.get("internal_assets_bytes") or []
        if isinstance(assets, bytes):
            assets = [assets]
        from io import BytesIO
        y = Inches(6.65)        # clearly below the table
        x = Inches(9.4)
        h_img = Inches(0.75)    # small thumbnails so they never reach the table
        pad = Inches(0.25)
        for ab in assets[:3]:
            try:
                sld.shapes.add_picture(BytesIO(ab), x, y, height=h_img)
                x += Inches(1.6) + pad
            except Exception:
                pass

        # Footer
        fbox = sld.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.35))
        ftf = fbox.text_frame; ftf.clear()
        fp = ftf.paragraphs[0]; fp.text = "Altossa – Internal Review"; fp.font.size = Pt(10)

    from io import BytesIO
    bio = BytesIO(); prs.save(bio); bio.seek(0)
    return bio.read()

# ---------- /Export helpers ----------
def safe_multiselect(label, options, default=None, **kwargs):
    """Multiselect wrapper that drops defaults not present in options (prevents StreamlitAPIException)."""
    default = default or []
    valid_default = [v for v in default if v in options]
    return st.multiselect(label, options=options, default=valid_default, **kwargs)

def _num_or_none(v):
    """
    Used when adding rows via the dashboard:
      - Pure number -> float
      - Range like '75-97' / '75/97' / '75 to 97' -> return original string
      - Else -> None
    """
    import re
    import unicodedata

    if v is None:
        return None
    s_orig = str(v).strip()
    if s_orig == "":
        return None

    s = unicodedata.normalize("NFKD", s_orig).encode("ascii", "ignore").decode("ascii").strip()

    range_pattern = re.compile(
        r"""^\s*
            \d+(?:[.,]\d+)?\s*
            (?:-|/|to|–|—)\s*
            \d+(?:[.,]\d+)?\s*$
        """,
        re.IGNORECASE | re.VERBOSE,
    )
    if range_pattern.match(s):
        return s_orig

    num_candidate = s.replace(",", ".")
    if re.match(r"^\s*\d+(?:\.\d+)?\s*$", num_candidate):
        try:
            return float(num_candidate)
        except Exception:
            return None

    return None

@st.cache_data(show_spinner=False)
def load_data(path):
    df = pd.read_excel(path)

    # numeric columns
    for c in ["W", "D", "H"]:
        if c in df.columns:
            df[c] = df[c].map(to_num)

    # string columns (add Composition link)
    str_cols = [
        "Brand", "Product name", "Link", "Product type", "Name of subtype",
        "Fabric Category", "Leather Category", "Other Categories", "Composition link"
    ]
    for c in str_cols:
        if c in df.columns:
            df[c] = df[c].fillna("").astype(str).map(norm_txt)
        else:
            df[c] = ""

    # price columns
    price_cols = ["Fabric Category Price", "Leather Category Price", "Other Price"]
    for c in price_cols:
        if c in df.columns:
            df[c] = pd.to_numeric(df[c], errors="coerce")
        else:
            df[c] = pd.NA

    # clean placeholders
    for col in ["Fabric Category", "Leather Category", "Other Categories"]:
        if col in df.columns:
            df[col] = df[col].replace({".": "", "-": "", "–": "", "—": ""})

    return df

def save_data(path: str, df: pd.DataFrame):
    num_cols = ["W", "D", "H", "Fabric Category Price", "Leather Category Price", "Other Price"]
    for c in num_cols:
        if c in df.columns:
            df[c] = pd.to_numeric(df[c], errors="coerce")
    df.to_excel(path, index=False)
    load_data.clear()

# --- ADD: centralized catalog cache invalidation ---
def _invalidate_catalog_cache() -> None:
    """Force-refresh the cached catalog after any DB mutation."""
    try:
        load_data_from_db.clear()
    except Exception:
        pass

@st.cache_data(show_spinner=False)
def load_data_from_db(database_url: str) -> pd.DataFrame:
    """
    Load catalog from Postgres and return a DataFrame shaped like the Excel:
      Brand, Product name, Link, Product type, Name of subtype, W, D, H,
      Fabric Category, Leather Category, Fabric Category Price, Leather Category Price,
      Other Categories, Other Price,
      Composition link, Composition parts, Composition part image URLs

    IMPORTANT:
    - Composition fields are pulled ONLY from variants.* (NOT products.*),
      so variants without compositions remain blank and won't "inherit" from other subtypes.
    - One row per pricing row (fabric/leather/other).
    """
    engine = db_engine(database_url)

    sql = """
    SELECT
      b.id   AS brand_id,
      pt.id  AS product_type_id,
      p.id   AS product_id,
      v.id   AS variant_id,
      pr.id  AS price_id,

      b.name  AS brand,
      p.name  AS product_name,
      p.link,
      pt.name AS product_type,
      v.subtype,
      v.width_raw, v.depth_raw, v.height_raw,
      v.width_mm,  v.depth_mm,  v.height_mm,
      pr.kind, pr.category, pr.base_price,

      -- Pull composition ONLY from the variant (not the product)
      v.composition_link,
      v.composition_parts,
      v.composition_part_image_urls
    FROM products p
    JOIN brands b           ON b.id  = p.brand_id
    JOIN product_types pt   ON pt.id = p.product_type_id
    LEFT JOIN variants v    ON v.product_id = p.id
    LEFT JOIN pricing_options pr ON pr.variant_id = v.id
    """
    raw = pd.read_sql(sql, engine)

    rows = []
    for _, r in raw.iterrows():
        base = {
            "brand_id": r.get("brand_id"),
            "product_type_id": r.get("product_type_id"),
            "product_id": r.get("product_id"),
            "variant_id": r.get("variant_id"),
            "price_id": r.get("price_id"),

            "Brand":        norm_txt(r.get("brand")),
            "Product name": norm_txt(r.get("product_name")),
            "Link":         norm_txt(r.get("link")),
            "Product type": norm_txt(r.get("product_type")),
            "Name of subtype": norm_txt(r.get("subtype")) if pd.notna(r.get("subtype")) else "",

            # Prefer raw W/D/H if provided, else mm values
            "W": r.get("width_raw")  if pd.notna(r.get("width_raw"))  and str(r.get("width_raw")).strip()  != "" else r.get("width_mm"),
            "D": r.get("depth_raw")  if pd.notna(r.get("depth_raw"))  and str(r.get("depth_raw")).strip()  != "" else r.get("depth_mm"),
            "H": r.get("height_raw") if pd.notna(r.get("height_raw")) and str(r.get("height_raw")).strip() != "" else r.get("height_mm"),

            "Fabric Category": "",
            "Leather Category": "",
            "Other Categories": "",
            "Fabric Category Price": pd.NA,
            "Leather Category Price": pd.NA,
            "Other Price": pd.NA,

            # Variant-scoped composition fields
            "Composition link":              norm_txt(r.get("composition_link")) if pd.notna(r.get("composition_link")) else "",
            "Composition parts":             norm_txt(r.get("composition_parts")) if pd.notna(r.get("composition_parts")) else "",
            "Composition part image URLs":   norm_txt(r.get("composition_part_image_urls")) if pd.notna(r.get("composition_part_image_urls")) else "",
        }

        kind = str(r.get("kind")).lower() if pd.notna(r.get("kind")) else None
        if kind == "fabric":
            base["Fabric Category"] = norm_txt(r.get("category"))
            base["Fabric Category Price"] = pd.to_numeric(r.get("base_price"), errors="coerce")
        elif kind == "leather":
            base["Leather Category"] = norm_txt(r.get("category"))
            base["Leather Category Price"] = pd.to_numeric(r.get("base_price"), errors="coerce")
        elif kind == "other":
            base["Other Categories"] = norm_txt(r.get("category"))
            base["Other Price"] = pd.to_numeric(r.get("base_price"), errors="coerce")

        rows.append(base)

    df = pd.DataFrame(rows)

    # Normalize numeric dims (no change)
    for c in ["W", "D", "H"]:
        if c in df.columns:
            df[c] = df[c].map(to_num)

    # Normalize strings – ALWAYS create these columns (even if df is empty)
    str_cols = [
        "Brand", "Product name", "Link", "Product type", "Name of subtype",
        "Fabric Category", "Leather Category", "Other Categories",
        "Composition link", "Composition parts", "Composition part image URLs",
    ]
    for c in str_cols:
        if c in df.columns:
            df[c] = df[c].fillna("").astype(str).map(norm_txt)
        else:
            # create empty string column so code like df["Product type"] never KeyErrors
            df[c] = ""

    # Numeric prices – ALWAYS make sure columns exist
    num_price_cols = ["Fabric Category Price", "Leather Category Price", "Other Price"]
    for c in num_price_cols:
        if c in df.columns:
            df[c] = pd.to_numeric(df[c], errors="coerce")
        else:
            # create empty float column for consistency
            df[c] = pd.Series(dtype="float64")

    # Clean placeholder tokens on category text columns
    for col in ["Fabric Category", "Leather Category", "Other Categories"]:
        if col in df.columns:
            df[col] = df[col].replace({".": "", "-": "", "–": "", "—": ""})

    return df

def _dim_raw_and_mm(v):
    s = None if pd.isna(v) else str(v).strip()
    if not s:
        return None, None
    # if looks like a range or non-numeric token → raw only
    if re.search(r"(-|–|—|/|\bto\b)", s, flags=re.I):
        return s, None
    try:
        mm = float(s.replace(",", "."))
        return s, mm
    except Exception:
        return s, None

def _kind_and_price(row: pd.Series):
    """Detect which price kind this row represents and return (kind, category, price)."""
    if pd.notna(row.get("Fabric Category Price")):
        return ("fabric", row.get("Fabric Category"), float(pd.to_numeric(row.get("Fabric Category Price"), errors="coerce")))
    if pd.notna(row.get("Leather Category Price")):
        return ("leather", row.get("Leather Category"), float(pd.to_numeric(row.get("Leather Category Price"), errors="coerce")))
    if pd.notna(row.get("Other Price")):
        return ("other", row.get("Other Categories"), float(pd.to_numeric(row.get("Other Price"), errors="coerce")))
    return (None, None, None)

def upsert_row_to_db(database_url: str, row: dict | pd.Series) -> None:
    """Upsert ONE Excel-like row into Postgres and invalidate the cached catalog."""
    if isinstance(row, dict):
        row = pd.Series(row)
    engine = db_engine(database_url)
    with engine.begin() as conn:
        # 1) brand
        brand = norm_txt(row.get("Brand"))
        bid = conn.execute(text("""
            INSERT INTO brands(name) VALUES (:n)
            ON CONFLICT (name) DO NOTHING
            RETURNING id;
        """), {"n": brand}).scalar_one_or_none()
        if bid is None:
            bid = conn.execute(text("SELECT id FROM brands WHERE name=:n"), {"n": brand}).scalar_one()

        # 2) product type
        ptype = norm_txt(row.get("Product type"))
        ptid = conn.execute(text("""
            INSERT INTO product_types(name) VALUES (:n)
            ON CONFLICT (name) DO NOTHING
            RETURNING id;
        """), {"n": ptype}).scalar_one_or_none()
        if ptid is None:
            ptid = conn.execute(text("SELECT id FROM product_types WHERE name=:n"), {"n": ptype}).scalar_one()

        # 3) product
        pname = norm_txt(row.get("Product name"))
        link = norm_txt(row.get("Link"))
        comp = norm_txt(row.get("Composition link"))
        comp_parts = norm_txt(row.get("Composition parts"))
        comp_part_urls = norm_txt(row.get("Composition part image URLs"))

        pid = conn.execute(text("""
            INSERT INTO products(brand_id, product_type_id, name, link)
            VALUES (:b,:t,:n,:l)
            ON CONFLICT (brand_id, product_type_id, name) DO UPDATE
              SET link = COALESCE(EXCLUDED.link, products.link)
            RETURNING id;
        """), {"b": bid, "t": ptid, "n": pname, "l": link}).scalar_one()

        # 4) variant
        subtype = norm_txt(row.get("Name of subtype"))
        wr, wmm = _dim_raw_and_mm(row.get("W"))
        dr, dmm = _dim_raw_and_mm(row.get("D"))
        hr, hmm = _dim_raw_and_mm(row.get("H"))

        vid = conn.execute(text("""
            INSERT INTO variants(
                product_id, subtype,
                width_raw, depth_raw, height_raw,
                width_mm, depth_mm, height_mm,
                composition_link,
                composition_parts,
                composition_part_image_urls
            )
            VALUES (:pid, :sub, :wr, :dr, :hr, :wmm, :dmm, :hmm, :comp, :comp_parts, :comp_part_urls)
            ON CONFLICT (product_id, subtype_norm, width_raw_norm, depth_raw_norm, height_raw_norm) DO UPDATE
              SET width_mm   = COALESCE(EXCLUDED.width_mm,  variants.width_mm),
                  depth_mm   = COALESCE(EXCLUDED.depth_mm,  variants.depth_mm),
                  height_mm  = COALESCE(EXCLUDED.height_mm, variants.height_mm),
                  composition_link = COALESCE(NULLIF(EXCLUDED.composition_link,''), variants.composition_link),
                  composition_parts = COALESCE(NULLIF(EXCLUDED.composition_parts,''), variants.composition_parts),
                  composition_part_image_urls = COALESCE(NULLIF(EXCLUDED.composition_part_image_urls,''), variants.composition_part_image_urls)
            RETURNING id;
        """), {
            "pid": pid, "sub": subtype or None,
            "wr": wr, "dr": dr, "hr": hr,
            "wmm": wmm, "dmm": dmm, "hmm": hmm,
            "comp": comp or None,
            "comp_parts": comp_parts or None,
            "comp_part_urls": comp_part_urls or None
        }).scalar_one()

        # 5) pricing (only one kind per row)
        kind, category, price = _kind_and_price(row)
        if kind is not None and price is not None:
            conn.execute(text("""
                INSERT INTO pricing_options(variant_id, kind, category, base_price)
                VALUES (:v, :k, :c, :p)
                ON CONFLICT (variant_id, kind, category_norm) DO UPDATE
                  SET base_price = EXCLUDED.base_price;
            """), {"v": vid, "k": kind, "c": norm_txt(category) or None, "p": float(price)})

    _invalidate_catalog_cache()

def delete_rows_in_db(database_url: str, rows_df: pd.DataFrame) -> int:
    """
    Delete by selected rows from the edited view.
    Priority:
      1) If price_id exists -> delete that pricing_options row.
      2) Else if variant_id exists -> delete that variant (cascades prices).
      3) Else if product_id exists -> delete the entire product and all its variants/prices.
    Returns number of delete statements executed (best-effort count).
    Also invalidates the cached catalog.
    """
    engine = db_engine(database_url)
    count = 0
    with engine.begin() as conn:
        for _, r in rows_df.iterrows():
            pid = r.get("price_id")
            vid = r.get("variant_id")
            prod_id = r.get("product_id")

            if pd.notna(pid):
                conn.execute(text("DELETE FROM pricing_options WHERE id=:id"), {"id": int(pid)})
                count += 1
                continue

            if pd.notna(vid):
                conn.execute(text("DELETE FROM variants WHERE id=:id"), {"id": int(vid)})
                count += 1
                continue

            if pd.notna(prod_id):
                # Delete all pricing rows for all variants of this product, then variants, then the product itself
                conn.execute(text("""
                    DELETE FROM pricing_options
                    WHERE variant_id IN (SELECT id FROM variants WHERE product_id=:pid)
                """), {"pid": int(prod_id)})
                conn.execute(text("DELETE FROM variants WHERE product_id=:pid"), {"pid": int(prod_id)})
                conn.execute(text("DELETE FROM products WHERE id=:pid"), {"pid": int(prod_id)})
                count += 1

    _invalidate_catalog_cache()
    return count

def prune_orphan_products(database_url: str) -> int:
    """
    Delete products that have no variants left. Returns number of products deleted.
    """
    engine = db_engine(database_url)
    with engine.begin() as conn:
        rows = conn.execute(text("""
            SELECT p.id
            FROM products p
            LEFT JOIN variants v ON v.product_id = p.id
            GROUP BY p.id
            HAVING COUNT(v.id) = 0
        """)).fetchall()
        ids = [int(r[0]) for r in rows]
        for pid in ids:
            conn.execute(text("DELETE FROM products WHERE id=:pid"), {"pid": pid})
    if rows:
        _invalidate_catalog_cache()
    return len(rows)

def update_product_details(
    database_url: str,
    product_id: int,
    brand_name: str,
    product_type_name: str,
    product_name: str,
    link: Optional[str],
    composition_link: Optional[str],
) -> None:
    """
    Update/move a PRODUCT row (not variants). Composition is NOT stored on products anymore.
      - Ensures brand & product type exist.
      - If moving to an existing (brand,type,name), merge by reparenting variants and delete the old product.
    """
    engine = db_engine(database_url)

    bname = norm_txt(brand_name)
    tname = norm_txt(product_type_name)
    pname = norm_txt(product_name)
    link  = norm_txt(link)

    with engine.begin() as conn:
        bid = conn.execute(text("""
            INSERT INTO brands(name) VALUES (:n)
            ON CONFLICT (name) DO NOTHING
            RETURNING id;
        """), {"n": bname}).scalar_one_or_none()
        if bid is None:
            bid = conn.execute(text("SELECT id FROM brands WHERE name=:n"), {"n": bname}).scalar_one()

        ptid = conn.execute(text("""
            INSERT INTO product_types(name) VALUES (:n)
            ON CONFLICT (name) DO NOTHING
            RETURNING id;
        """), {"n": tname}).scalar_one_or_none()
        if ptid is None:
            ptid = conn.execute(text("SELECT id FROM product_types WHERE name=:n"), {"n": tname}).scalar_one()

        target_pid = conn.execute(text("""
            SELECT id FROM products
            WHERE brand_id = :b AND product_type_id = :t AND name = :n
            LIMIT 1
        """), {"b": bid, "t": ptid, "n": pname}).scalar_one_or_none()

        if target_pid is None or int(target_pid) == int(product_id):
            conn.execute(text("""
                UPDATE products
                SET brand_id = :b,
                    product_type_id = :t,
                    name = :n,
                    link = :l
                WHERE id = :pid
            """), {"b": bid, "t": ptid, "n": pname, "l": (link or None), "pid": int(product_id)})

        else:
            conn.execute(text("""
                DELETE FROM variants v
                USING variants vt
                WHERE v.product_id = :old_pid
                  AND vt.product_id = :new_pid
                  AND COALESCE(v.subtype_norm,'')     = COALESCE(vt.subtype_norm,'')
                  AND COALESCE(v.width_raw_norm,'')   = COALESCE(vt.width_raw_norm,'')
                  AND COALESCE(v.depth_raw_norm,'')   = COALESCE(vt.depth_raw_norm,'')
                  AND COALESCE(v.height_raw_norm,'')  = COALESCE(vt.height_raw_norm,'')
            """), {"old_pid": int(product_id), "new_pid": int(target_pid)})

            conn.execute(text("""
                UPDATE variants
                SET product_id = :new_pid
                WHERE product_id = :old_pid
            """), {"new_pid": int(target_pid), "old_pid": int(product_id)})

            conn.execute(text("""
                UPDATE products
                SET link = COALESCE(NULLIF(:l, ''), link)
                WHERE id = :new_pid
            """), {"l": (link or None), "new_pid": int(target_pid)})

            conn.execute(text("DELETE FROM products WHERE id = :pid"), {"pid": int(product_id)})

    _invalidate_catalog_cache()

def update_variant_and_prices(
    database_url: str,
    variant_id: int,
    subtype: Optional[str],
    w: Optional[str],
    d: Optional[str],
    h: Optional[str],
    fab_cat: Optional[str],
    fab_price: Optional[str | float],
    lea_cat: Optional[str],
    lea_price: Optional[str | float],
    oth_cat: Optional[str],
    oth_price: Optional[str | float],
    comp_link: Optional[str],
    comp_parts: Optional[str],
    comp_part_urls: Optional[str],
) -> None:
    """
    Update one existing VARIANT (subtype + W/D/H) and its three price kinds,
    and also the variant-scoped composition fields:
      - composition_link
      - composition_parts
      - composition_part_image_urls

    Conflict-safe merge remains identical to before.
    """
    engine = db_engine(database_url)

    wr, wmm = _dim_raw_and_mm(w)
    dr, dmm = _dim_raw_and_mm(d)
    hr, hmm = _dim_raw_and_mm(h)

    def _num(v):
        try:
            x = pd.to_numeric(v, errors="coerce")
            return None if pd.isna(x) else float(x)
        except Exception:
            return None

    def _norm_txt_db(s: Optional[str]) -> Optional[str]:
        if s is None:
            return None
        s2 = str(s).strip()
        return s2 if s2 != "" else None

    with engine.begin() as conn:
        product_id = conn.execute(
            text("SELECT product_id FROM variants WHERE id=:vid"),
            {"vid": int(variant_id)}
        ).scalar_one()

        target_variant_id = conn.execute(text("""
            SELECT v2.id
            FROM variants v2
            WHERE v2.product_id = :pid
              AND v2.id <> :vid
              AND COALESCE(lower(btrim(v2.subtype)), '')      = COALESCE(lower(btrim(:sub)), '')
              AND COALESCE(lower(btrim(v2.width_raw)), '')    = COALESCE(lower(btrim(:wr)),  '')
              AND COALESCE(lower(btrim(v2.depth_raw)), '')    = COALESCE(lower(btrim(:dr)),  '')
              AND COALESCE(lower(btrim(v2.height_raw)), '')   = COALESCE(lower(btrim(:hr)),  '')
            LIMIT 1
        """), {
            "pid": int(product_id),
            "vid": int(variant_id),
            "sub": _norm_txt_db(subtype),
            "wr": wr, "dr": dr, "hr": hr
        }).scalar_one_or_none()

        # --- no collision → update this variant in place (incl. compositions)
        if target_variant_id is None:
            conn.execute(
                text("""
                    UPDATE variants
                    SET subtype   = :sub,
                        width_raw = :wr,  depth_raw = :dr,  height_raw = :hr,
                        width_mm  = :wmm, depth_mm  = :dmm, height_mm  = :hmm,
                        composition_link = COALESCE(NULLIF(:comp,''),
                                                    composition_link),
                        composition_parts = COALESCE(NULLIF(:comp_parts,''),
                                                     composition_parts),
                        composition_part_image_urls = COALESCE(NULLIF(:comp_urls,''),
                                                               composition_part_image_urls)
                    WHERE id = :vid
                """),
                {"sub": _norm_txt_db(subtype), "wr": wr, "dr": dr, "hr": hr,
                 "wmm": wmm, "dmm": dmm, "hmm": hmm, "vid": int(variant_id),
                 "comp": (comp_link or None),
                 "comp_parts": (comp_parts or None),
                 "comp_urls": (comp_part_urls or None)}
            )
            survivor_vid = int(variant_id)

        # --- collision → merge into survivor, then update survivor’s mm + compositions
        else:
            survivor_vid = int(target_variant_id)

            price_rows = conn.execute(text("""
                SELECT kind, category, base_price
                FROM pricing_options
                WHERE variant_id=:v
            """), {"v": int(variant_id)}).fetchall()

            for kind, category, base_price in price_rows:
                conn.execute(text("""
                    INSERT INTO pricing_options(variant_id, kind, category, base_price)
                    VALUES (:v, :k, :c, :p)
                    ON CONFLICT (variant_id, kind, category_norm) DO UPDATE
                      SET base_price = EXCLUDED.base_price,
                          category   = EXCLUDED.category
                """), {
                    "v": survivor_vid,
                    "k": kind,
                    "c": norm_txt(category) or None,
                    "p": float(base_price) if base_price is not None else None
                })

            conn.execute(text("DELETE FROM pricing_options WHERE variant_id=:v"), {"v": int(variant_id)})
            conn.execute(text("DELETE FROM variants WHERE id=:v"), {"v": int(variant_id)})

            conn.execute(text("""
                UPDATE variants
                SET width_mm  = COALESCE(:wmm, width_mm),
                    depth_mm  = COALESCE(:dmm, depth_mm),
                    height_mm = COALESCE(:hmm, height_mm),
                    composition_link = COALESCE(NULLIF(:comp,''),
                                                composition_link),
                    composition_parts = COALESCE(NULLIF(:comp_parts,''),
                                                 composition_parts),
                    composition_part_image_urls = COALESCE(NULLIF(:comp_urls,''),
                                                           composition_part_image_urls)
                WHERE id = :vid
            """), {"wmm": wmm, "dmm": dmm, "hmm": hmm, "vid": survivor_vid,
                   "comp": (comp_link or None),
                   "comp_parts": (comp_parts or None),
                   "comp_urls": (comp_part_urls or None)})

        # prices (unchanged)
        def _upsert_or_delete(kind: str, cat: Optional[str], price):
            price_val = _num(price)
            cat_norm  = norm_txt(cat) or None
            if price_val is None:
                conn.execute(text("DELETE FROM pricing_options WHERE variant_id=:v AND kind=:k"),
                             {"v": survivor_vid, "k": kind})
            else:
                conn.execute(text("""
                    INSERT INTO pricing_options(variant_id, kind, category, base_price)
                    VALUES (:v, :k, :c, :p)
                    ON CONFLICT (variant_id, kind, category_norm) DO UPDATE
                      SET base_price = EXCLUDED.base_price,
                          category   = EXCLUDED.category
                """), {"v": survivor_vid, "k": kind, "c": cat_norm, "p": float(price_val)})

        _upsert_or_delete("fabric",  fab_cat, fab_price)
        _upsert_or_delete("leather", lea_cat, lea_price)
        _upsert_or_delete("other",   oth_cat, oth_price)

    _invalidate_catalog_cache()

@st.cache_resource
def _http_session() -> requests.Session:
    s = requests.Session()
    s.headers.update({
        "User-Agent": "Altossa/1.0 (+streamlit)",
        "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8",
    })
    return s

def _img_key(product_link: str, composition_link: str | None) -> str:
    return f"{(product_link or '').strip()}|{(composition_link or '').strip()}"

def _persisted_get_image(key: str) -> Optional[bytes]:
    """Return bytes if in memory store and not expired."""
    store = _get_memory_image_store()
    if key not in store:
        return None
    
    img, ts = store[key]
    
    # TTL check
    if int(time.time()) - int(ts) > IMG_CACHE_TTL_SECONDS:
        try:
            del store[key]
        except KeyError:
            pass
        return None
        
    return img

def _persisted_put_image(key: str, img: Optional[bytes]) -> None:
    """Upsert image bytes into memory store with current timestamp."""
    if not img:
        return
    store = _get_memory_image_store()
    store[key] = (img, int(time.time()))

def _persisted_cleanup(limit_rows: int = 5000) -> None:
    """Keep memory store small: delete oldest entries beyond limit."""
    store = _get_memory_image_store()
    if len(store) <= limit_rows:
        return
    
    # Sort by timestamp (oldest first)
    # items() returns (key, (img, ts))
    items = sorted(store.items(), key=lambda x: x[1][1])
    
    # Calculate how many to remove
    to_remove_count = len(store) - limit_rows
    
    for i in range(to_remove_count):
        k = items[i][0]
        try:
            del store[k]
        except KeyError:
            pass

def get_image_if_cached(product_link: str, composition_link: Optional[str]) -> Optional[bytes]:
    """
    Return image bytes ONLY if already cached (session or SQLite). Never hits the network.
    """
    if "img_cache" not in st.session_state:
        st.session_state.img_cache = {}

    key = _img_key(product_link, composition_link)

    # session cache
    if key in st.session_state.img_cache:
        return st.session_state.img_cache[key]

    # persisted cache
    img = _persisted_get_image(key)
    if img:
        st.session_state.img_cache[key] = img
        return img

    return None

def get_image_cached(product_link: str, composition_link: Optional[str]) -> Optional[bytes]:
    """
    3-layer cache for image bytes:
      1) session (fastest)
      2) SQLite (24h persistent)
      3) network fetch (then persist to 1 & 2)
    """
    if "img_cache" not in st.session_state:
        st.session_state.img_cache = {}

    key = _img_key(product_link, composition_link)

    # 1) session cache
    if key in st.session_state.img_cache:
        return st.session_state.img_cache[key]

    # 2) persisted (24h)
    img = _persisted_get_image(key)
    if img:
        st.session_state.img_cache[key] = img
        return img

    # 3) fetch
    img = _resolve_product_image(product_link, composition_link)
    st.session_state.img_cache[key] = img
    _persisted_put_image(key, img)

    # trim session cache a bit
    if len(st.session_state.img_cache) > IMG_CACHE_MAX:
        for old in list(st.session_state.img_cache.keys())[: max(1, IMG_CACHE_MAX // 10)]:
            st.session_state.img_cache.pop(old, None)

    _persisted_cleanup(limit_rows=5000)
    return img

# === NEW: build a big "hero" URL for Google Drive links (1600px wide) ===
_DRIVE_ID = re.compile(r"""
    (?:/d/|id=)                  # /d/<ID> or id=<ID>
    ([A-Za-z0-9_-]{20,})         # the file id
""", re.X)

def _drive_hero_url(url: str, width: int = 1600) -> str:
    """Return a large, direct image URL for a Google Drive share link."""
    if not url:
        return url
    m = _DRIVE_ID.search(url)
    if not m:
        return url
    fid = m.group(1)
    # thumbnail endpoint returns an actual image at requested width
    return f"https://drive.google.com/thumbnail?id={fid}&sz=w{int(width)}"

def _split_pipe_list(s: str) -> list[str]:
    """Split a pipe-delimited string into a trimmed list; tolerates surrounding spaces."""
    if not s or str(s).strip() == "":
        return []
    return [p.strip() for p in str(s).split("|") if str(p).strip()]

def render_compositions_for_variant(g: pd.DataFrame, product_link: str) -> None:
    """
    Render composition parts in a responsive grid with larger tiles.
    - Shows cached images immediately (no blocking HTTP).
    - Triggers background prefetch for missing ones.
    - Uses 1..3 columns adaptively per row so images are larger.
    """
    if g is None or g.empty:
        st.info("No composition data for this selection.")
        return

    # Pull the first non-empty composition strings
    parts_str = ""
    urls_str  = ""
    for _, r in g.iterrows():
        ps = str(r.get("Composition parts") or "").strip()
        us = str(r.get("Composition part image URLs") or "").strip()
        if ps or us:
            parts_str, urls_str = ps, us
            break

    parts = _split_pipe_list(parts_str)
    urls  = _split_pipe_list(urls_str)

    if not parts and not urls:
        st.info("No composition data for this selection.")
        return

    # Pair up by index; truncate to shortest length
    n = min(len(parts), len(urls)) if urls else len(parts)
    items = []
    for i in range(n):
        label = parts[i] if i < len(parts) else ""
        url   = urls[i]  if i < len(urls)  else ""
        items.append((label, url))

    # Kick off a background prefetch for all images up front (non-blocking)
    prefetch_product_images_async([("", u) for _, u in items if u], max_workers=8)

    # Determine columns per row: 1 for 1 item, 2 for 2 items, else 3
    def _cols_for_row(count: int):
        if count == 1: return st.columns(1)
        if count == 2: return st.columns(2)
        return st.columns(3)

    # Render grid
    for i in range(0, len(items), 3):
        row_slice   = items[i:i+3]
        cols        = _cols_for_row(len(row_slice))
        for j, (label, url) in enumerate(row_slice):
            with cols[j].container(border=True):
                if label:
                    st.markdown(f"**{label}**")
                img_bytes = get_image_if_cached("", url)  # non-blocking cache check
                if img_bytes:
                    # new API: width='stretch' instead of use_container_width=True
                    st.image(img_bytes, width="stretch")
                else:
                    # Placeholder while background thread fetches
                    ph = st.empty()
                    ph.markdown(
                        "<div style='height:140px; background:rgba(255,255,255,0.05); "
                        "border-radius:8px; display:flex; align-items:center; justify-content:center;'>"
                        "Loading…</div>",
                        unsafe_allow_html=True,
                    )

def prefetch_product_images(pairs: list[tuple[str, Optional[str]]], max_workers: int = 8) -> None:
    """
    Preload many images in parallel.
    Uses session cache, then SQLite, then fetches missing in threads.
    """
    if "img_cache" not in st.session_state:
        st.session_state.img_cache = {}

    # Normalize + split into (already_cached) and (need_fetch)
    need_fetch_keys = []
    for pl, cl in pairs:
        key = _img_key(pl, cl)
        if key in st.session_state.img_cache:
            continue
        img = _persisted_get_image(key)
        if img:
            st.session_state.img_cache[key] = img
        else:
            need_fetch_keys.append(key)

    if not need_fetch_keys:
        return

    def _task(key: str):
        pl, cl = key.split("|", 1)
        try:
            img = _resolve_product_image(pl, cl or None)
            return key, img
        except Exception:
            return key, None

    from concurrent.futures import ThreadPoolExecutor, as_completed
    with ThreadPoolExecutor(max_workers=max_workers) as ex:
        futures = [ex.submit(_task, k) for k in need_fetch_keys]
        for fut in as_completed(futures):
            key, img = fut.result()
            st.session_state.img_cache[key] = img
            _persisted_put_image(key, img)

    _persisted_cleanup(limit_rows=5000)

def _start_daemon_thread(name: str, target, args=()):
    """
    Start a daemon thread and attach Streamlit's ScriptRunContext if available
    to avoid 'missing ScriptRunContext' warnings.
    """
    import threading
    t = threading.Thread(name=name, target=target, args=args, daemon=True)
    try:
        # Streamlit ≥1.31 has this util; safe to no-op if not present
        from streamlit.runtime.scriptrunner import add_script_run_ctx
        add_script_run_ctx(t)
    except Exception:
        pass
    t.start()
    return t

def prefetch_product_images_async(pairs: list[tuple[str, Optional[str]]], max_workers: int = 8) -> None:
    """
    Fire-and-forget prefetch. Returns immediately while images are loaded in background.
    """
    if not pairs:
        return
    _start_daemon_thread(
        name="prefetch_product_images",
        target=prefetch_product_images,
        args=(pairs, max_workers),
    )

def _drive_candidates(url: str, size: int = 1600) -> list[str]:
    """
    Given any Google Drive URL, return a list of fast, direct image candidates (thumbnail + uc).
    Falls back to the original URL.
    """
    u = (url or "").strip()
    if not u:
        return []
    m = re.search(r"drive\.google\.com/(?:file/d/|open\?id=|uc\?.*?id=|thumbnail\?id=)([^/&?]+)", u)
    if not m:
        return [u]
    fid = m.group(1)
    return [
        f"https://drive.google.com/thumbnail?id={fid}&sz=w{size}",
        f"https://drive.google.com/uc?export=view&id={fid}",
        f"https://drive.google.com/uc?export=download&id={fid}",
        u,
    ]

def start_global_image_warmup(df: pd.DataFrame, max_workers: int = 12) -> None:
    """
    Warm the 24h on-disk cache ONCE per session:
      • Product heroes (from product links)
      • All composition images (from 'Composition link' and 'Composition part image URLs')
    Runs in a daemon thread and never blocks the UI.

    Safe when the dataframe is empty or missing columns (e.g., empty database).
    """
    # If there is no data at all, or no Link column, just skip warmup gracefully.
    if df is None or df.empty or "Link" not in df.columns:
        return

    hero_pairs: set[tuple[str, str]] = set()
    comp_pairs: set[tuple[str, str]] = set()

    # Product page heroes (force composition empty so we always fetch hero)
    for link in df["Link"].dropna().astype(str):
        link = link.strip()
        if link:
            hero_pairs.add((link, ""))

    # Composition link per row (variant-scoped)
    if "Composition link" in df.columns:
        for c in df["Composition link"].dropna().astype(str):
            c = c.strip()
            if c:
                # choose best Drive candidate but keep key stable by using original in key
                for cand in _drive_candidates(c):
                    comp_pairs.add(("", cand))

    # Composition part image URLs (pipe-separated)
    if "Composition part image URLs" in df.columns:
        for s in df["Composition part image URLs"].dropna().astype(str):
            for piece in (p.strip() for p in s.split("|") if p.strip()):
                for cand in _drive_candidates(piece):
                    comp_pairs.add(("", cand))

    # Fire and forget
    prefetch_product_images_async(list(hero_pairs), max_workers=max_workers)
    prefetch_product_images_async(list(comp_pairs), max_workers=max_workers)

@st.cache_data(show_spinner=False, ttl=21600)  # 6 hours
def fetch_image_from_page(url: str):
    """Extract og:image or distinct product image using a shared HTTP session."""
    from urllib.parse import urljoin
    try:
        # [UPDATED] Robust fetch with absolute URL conversion
        r = _http_session().get(url, timeout=8)
        if not r.ok:
            return None
            
        soup = BeautifulSoup(r.text, "html.parser")
        
        # [NEW] Cattelan Italia specific handler
        # Targets the specific slider structure used on their product pages
        if "cattelanitalia.com" in url.lower():
            # Use CSS selector to find the landscape image wrapper and the slide image inside it
            cattelan_img = soup.select_one("div.image-wrapper.landscape img.background-slide")
            if cattelan_img:
                # Prefer data-src (often used for lazy loading high-res), fallback to src
                img_url = cattelan_img.get("data-src") or cattelan_img.get("src")
                if img_url:
                    return urljoin(url, img_url)

        # Priority 1: OpenGraph (Standard for most other sites)
        candidates = [
            soup.find("meta", property="og:image"),
            soup.find("meta", name="twitter:image"),
            soup.find("link", rel="image_src"),
        ]
        
        for c in candidates:
            if c and c.get("content"):
                return urljoin(url, c.get("content"))
            if c and c.get("href"):
                return urljoin(url, c.get("href"))

        # Priority 2: Generic fallback classes (Hero/Main/Product)
        img = soup.find("img", class_=re.compile(r"product|hero|main|detail", re.I))
        if img and img.get("src"):
            return urljoin(url, img.get("src"))

    except Exception:
        pass
    return None

@st.cache_data(show_spinner=False, ttl=21600)  # 6 hours
def _resolve_product_image(product_link: str, composition_link: str | None) -> Optional[bytes]:
    """
    Return raw image bytes ready for st.image(...).
    Preference:
      1) og:image from product_link (hero)
      2) composition_link (Drive-optimized sequence)
    """
    session = _http_session()

    def _get_bytes(u: str, timeout: float = 6.0) -> Optional[bytes]:
        try:
            r = session.get(u, timeout=timeout, allow_redirects=True, stream=True)
            if not r.ok:
                return None
            # Some Drive endpoints serve octet-stream; accept both image/* and octet-stream
            ct = (r.headers.get("Content-Type") or "").lower()
            if (ct.startswith("image/") or "octet-stream" in ct):
                # read at most ~5 MB to avoid runaway downloads
                chunk = r.raw.read(5 * 1024 * 1024)
                return chunk if chunk else None
        except Exception:
            return None
        return None

    # 1) Try product hero
    if product_link:
        og = fetch_image_from_page(product_link)
        if og:
            b = _get_bytes(og, timeout=6.0)
            if b:
                return b

    # 2) Fallback to composition link (Drive candidates first)
    comp = (composition_link or "").strip()
    if comp:
        for cand in _drive_candidates(comp, size=1600):
            b = _get_bytes(cand, timeout=6.0)
            if b:
                return b

    return None

# === NEW: product-type tiles (animated crossfade, hover title) ===

@st.cache_data(show_spinner=False, ttl=21600)
def representative_image_urls_for_type(df: pd.DataFrame, type_name: str, max_images: int = 4) -> list[str]:
    """
    Pick up to 'max_images' representative images for a Product type:
      1) Prefer variant 'Composition link' (Drive hero URL at 1600px)
      2) Else use og:image from the product 'Link'
    Ensures each comes from a distinct product_id.
    """
    if not type_name:
        return []

    sub = df[df["Product type"].astype(str).str.lower() == str(type_name).lower()].copy()
    if sub.empty:
        return []

    urls: list[str] = []
    used_pids: set[int] = set()
    for _, r in sub.sort_values("product_id").iterrows():
        try:
            pid = int(r.get("product_id"))
        except Exception:
            continue
        if pid in used_pids:
            continue

        # 1) composition hero
        comp = str(r.get("Composition link") or "").strip()
        picked: str | None = None
        if comp:
            picked = _drive_hero_url(comp, width=1600)

        # 2) fallback to og:image of product page
        if not picked:
            link = str(r.get("Link") or "").strip()
            if link:
                og = fetch_image_from_page(link)  # returns URL string or None
                if og:
                    picked = og

        if picked:
            urls.append(picked)
            used_pids.add(pid)

        if len(urls) >= max_images:
            break

    return urls

def first_product_links_for_type(df: pd.DataFrame, type_name: str) -> tuple[str, str]:
    """
    Return (product_link, composition_link) for the FIRST product in this type
    that actually produces a valid image.

    Behaviour:
    - Iterate products in this Product type in the same order as the Products page.
    - For each product, try get_image_cached(product_link, composition_link).
      If it returns bytes, use that product for the type tile.
    - If none of the products yield an image, fall back to the very first
      product's (Link, "") so behaviour is deterministic.
    """
    if not type_name:
        return "", ""

    sub = df[df["Product type"].astype(str).str.lower() == str(type_name).lower()]
    if sub.empty:
        return "", ""

    fallback_links: tuple[str, str] | None = None

    # Walk products in the same order as elsewhere (by product_id)
    for _, g0 in sub.groupby("product_id", sort=False):
        row0 = g0.iloc[0]
        product_link = str(row0.get("Link") or "").strip()
        comp_link = str(row0.get("Composition link") or "").strip()

        if fallback_links is None:
            # Remember the very first product in case no images work at all
            fallback_links = (product_link, "")

        # If we have at least one URL, check whether it can actually render an image
        if product_link or comp_link:
            img_bytes = get_image_cached(product_link, comp_link or "")
            if img_bytes:
                # Use this product for the tile
                return product_link, (comp_link or "")

    # No product yielded an image – fall back to the very first product
    if fallback_links is not None:
        return fallback_links

    return "", ""

def render_type_tile(ptype: str, product_link: str, composition_link: str, key: str) -> None:
    """
    Render a 3:2 tile for a Product type.

    • Image source is EXACTLY the same as the first product card inside that type:
      get_image_cached(product_link, composition_link).
    • The tile is purely visual; navigation happens via a small "View" button
      rendered just below the tile, which calls open_product_type(ptype) and
      keeps us inside the same Streamlit session (no password prompt).
    • The product type name is ALWAYS visible (no hover needed).
    """
    import base64, re, uuid

    # Try to fetch bytes using the same logic as product cards
    data_uri = ""
    if product_link or composition_link:
        b = (
            get_image_if_cached(product_link, composition_link)
            or get_image_cached(product_link, composition_link)
        )
        if b:
            try:
                data_uri = "data:image/*;base64," + base64.b64encode(b).decode("ascii")
            except Exception:
                data_uri = ""

    safe_key = re.sub(r"[^a-zA-Z0-9_-]+", "", key) or ("k" + uuid.uuid4().hex[:6])
    cls = f"tile_{safe_key}"

    style_bg = (
        f"background-image:url('{data_uri}');"
        if data_uri
        else "background:#0b0f1a;"
    )

    css = f"""
    <style>
      .{cls}-wrap {{
        margin: 10px 8px 8px 8px;
      }}
      .{cls} {{
        position: relative;
        width: 100%;
        aspect-ratio: 3 / 2;
        border-radius: 12px;
        overflow: hidden;
        background: #0b0f1a;
        box-shadow: 0 8px 24px rgba(0,0,0,0.25);
      }}
      .{cls}__img {{
        position: absolute;
        inset: 0;
        background-size: cover;
        background-position: center;
        {style_bg}
      }}
      /* Overlay is ALWAYS visible now (no hover needed) */
      .{cls}__overlay {{
        position: absolute;
        inset: 0;
        display: flex;
        align-items: flex-end;
        justify-content: flex-start;
        padding: 12px 14px;
        background: linear-gradient(
          180deg,
          rgba(0, 0, 0, 0.0) 45%,
          rgba(0, 0, 0, 0.55) 100%
        );
        opacity: 1;
      }}
      .{cls}__title {{
        font-weight: 600;
        font-size: 1.1rem;
        color: #fff;
        text-shadow: 0 1px 2px rgba(0,0,0,0.5);
      }}
    </style>
    """

    # Pure visual tile (no link, no page reload)
    html = f"""
      {css}
      <div class="{cls}-wrap">
        <div class="{cls}">
          <div class="{cls}__img"></div>
          <div class="{cls}__overlay">
            <div class="{cls}__title">{ptype}</div>
          </div>
        </div>
      </div>
    """
    st.markdown(html, unsafe_allow_html=True)

    # Small "View" button under the tile – uses in-session navigation
    if st.button("View", key=f"view_{safe_key}", use_container_width=True):
        open_product_type(ptype)

def _dedupe_size_cards(g: pd.DataFrame):
    """
    From a product rows DataFrame (g), return unique size cards by (W, D, H, Name of subtype).
    Only include rows that have full dimensions (all of W, D, H notna()).
    If multiple rows share the same dimensions but different subtypes, they are kept separately.
    """
    cards = []
    seen = set()
    for idx, r in g.iterrows():
        W, D, H = r.get("W"), r.get("D"), r.get("H")
        subtype = r.get("Name of subtype", "").strip()
        if pd.notna(W) and pd.notna(D) and pd.notna(H):
            # Uniqueness key: dimensions + subtype
            key = (W, D, H, subtype)
            if key not in seen:
                seen.add(key)
                cards.append({
                    "idx": idx,
                    "W": W, "D": D, "H": H,
                    "subtype": subtype,
                    "row": r
                })
    return cards

def _ensure_item_uid(idx: int) -> str:
    """
    Ensure st.session_state.cart[idx] has a *stable* per-line UID that never
    depends on product fields (brand/type/dimensions/etc.). This prevents
    attachments from 'disappearing' when any field representation changes.
    """
    item = st.session_state.cart[idx]
    cur = str(item.get("uid") or "").strip()
    if cur:
        return cur
    new_uid = uuid.uuid4().hex
    item["uid"] = new_uid
    return new_uid

def _uploader_key_for(uid: str) -> str:
    return f"upload_{uid}"

def _stable_key(item: dict) -> str:
    """
    Build a stable identity for a selected item based on its semantics, not index.
    Keys on: Brand | Product type | Product name | Name of subtype | W | D | H
    """
    def _nz(v):
        return "" if v is None else str(v).strip()
    return "|".join([
        _nz(item.get("Brand")),
        _nz(item.get("Product type")),
        _nz(item.get("Product name")),
        _nz(item.get("Name of subtype")),
        _nz(item.get("W")),
        _nz(item.get("D")),
        _nz(item.get("H")),
    ])

def _uid_index() -> dict:
    """
    Central mapping: stable_key -> uid. Lives in session_state.
    Ensures we reuse the same uid for an item even if the cart reorders or re-renders.
    """
    if "uid_index" not in st.session_state:
        st.session_state.uid_index = {}
    return st.session_state.uid_index

# >>> ADD THIS BLOCK (right after _uploader_key_for)
def _attach_store() -> dict[str, list[bytes]]:
    """Central store for attachments keyed by item uid."""
    if "attach_store" not in st.session_state:
        st.session_state.attach_store = {}
    return st.session_state.attach_store

def _attachments_dir() -> pathlib.Path:
    """Directory to persist per-item attachment bytes across reruns."""
    d = pathlib.Path(".altossa_attachments")
    d.mkdir(exist_ok=True)
    return d

def _asset_paths_for(uid: str) -> list[pathlib.Path]:
    """List all saved files for a uid."""
    d = _attachments_dir()
    return sorted(d.glob(f"{uid}_*.bin"))

def _save_assets_to_disk(uid: str, assets: list[bytes]) -> None:
    """Rewrite all files for uid (truncate + write fresh)."""
    d = _attachments_dir()
    # wipe old files first
    for p in _asset_paths_for(uid):
        try: p.unlink()
        except Exception: pass
    # write fresh
    for i, b in enumerate(assets):
        try:
            (d / f"{uid}_{i}.bin").write_bytes(b)
        except Exception:
            pass

def _load_assets_from_disk(uid: str, cap: int = 3) -> list[bytes]:
    """Load assets for uid from disk (up to cap)."""
    out: list[bytes] = []
    for p in _asset_paths_for(uid)[:cap]:
        try:
            out.append(p.read_bytes())
        except Exception:
            pass
    return out

def _delete_assets_for(uid: str) -> None:
    """Delete both in-memory and disk-backed assets for a uid."""
    # in-memory
    _attach_store().pop(uid, None)
    # on disk
    for p in _asset_paths_for(uid):
        try: p.unlink()
        except Exception: pass

def _get_assets_for(uid: str) -> list[bytes]:
    """Read attachments for a uid from memory, falling back to disk."""
    store = _attach_store()
    assets = store.get(uid)
    if assets is None:
        # try disk
        assets = _load_assets_from_disk(uid, cap=3)
        store[uid] = assets or []
    if isinstance(assets, bytes):
        assets = [assets]
        store[uid] = assets
    return assets

def _set_assets_for(uid: str, assets: list[bytes], cap: int = 3) -> None:
    """Write attachments for uid into memory + disk (dedup + cap)."""
    from hashlib import sha1
    uniq, seen = [], set()
    for b in assets[:cap]:
        h = sha1(b).hexdigest()
        if h not in seen:
            uniq.append(b); seen.add(h)
    # persist in memory
    _attach_store()[uid] = uniq
    # persist on disk
    _save_assets_to_disk(uid, uniq)

def _cart_state_file() -> Path:
    d = Path(".altossa_state")
    d.mkdir(exist_ok=True)
    return d / "cart_state.json"

def _snapshot_cart_dict() -> dict:
    """Minimal state we want to persist across sleeps/restarts."""
    return {
        "cart": st.session_state.get("cart", []),
        "client_name": st.session_state.get("client_name_sidebar", ""),
        "approved_by": st.session_state.get("approved_by_sidebar", ""),
        "last_selection_ctx": st.session_state.get("last_selection_ctx", None),
    }

def _restore_cart_from_disk() -> None:
    """One-time best-effort restore of a previous session."""
    try:
        p = _cart_state_file()
        if not p.exists():
            return
        data = json.loads(p.read_text(encoding="utf-8"))
        if isinstance(data, dict):
            if isinstance(data.get("cart"), list) and not st.session_state.get("cart"):
                st.session_state.cart = data["cart"]
            if isinstance(data.get("client_name"), str) and not st.session_state.get("client_name_sidebar"):
                st.session_state.client_name_sidebar = data["client_name"]
            if isinstance(data.get("approved_by"), str) and not st.session_state.get("approved_by_sidebar"):
                st.session_state.approved_by_sidebar = data["approved_by"]
            if data.get("last_selection_ctx") and not st.session_state.get("last_selection_ctx"):
                st.session_state.last_selection_ctx = data["last_selection_ctx"]
    except Exception:
        pass  # never block UI

def _autosave_cart_if_changed() -> None:
    """Write snapshot to disk only when the payload changes."""
    try:
        snap = _snapshot_cart_dict()
        blob = json.dumps(snap, ensure_ascii=False, separators=(",", ":"), sort_keys=True)
        sig = sha1(blob.encode("utf-8")).hexdigest()
        if st.session_state.get("__cart_sig__") != sig:
            _cart_state_file().write_text(blob, encoding="utf-8")
            st.session_state["__cart_sig__"] = sig
    except Exception:
        pass

# ---------------------------
# Session state
# ---------------------------
if "page" not in st.session_state: st.session_state.page = "home"   # rooms-first
if "sel_type" not in st.session_state: st.session_state.sel_type = None
if "sel_room" not in st.session_state: st.session_state.sel_room = None

# NEW – persist filters
if "last_brand_filter" not in st.session_state:st.session_state.last_brand_filter = None
if "last_search_text" not in st.session_state:st.session_state.last_search_text = ""

if "sel_product" not in st.session_state: st.session_state.sel_product = None
if "sel_size_idx" not in st.session_state: st.session_state.sel_size_idx = None
if "cart" not in st.session_state: st.session_state.cart = []
if "nav" not in st.session_state: st.session_state.nav = []
if "fwd" not in st.session_state: st.session_state.fwd = []
if "room_types" not in st.session_state: st.session_state.room_types = DEFAULT_ROOM_TYPES.copy()
if "to_delete_defaults" not in st.session_state: st.session_state.to_delete_defaults = []
if "sel_product_id" not in st.session_state: st.session_state.sel_product_id = None

# NEW: persist product page filters (Brand, Price, Search) across navigation
if "product_filters" not in st.session_state: st.session_state.product_filters = {}

# NEW: compositions-only path (used by the "Available Compositions" entry on Home)
if "compositions_only" not in st.session_state: st.session_state.compositions_only = False

# NEW: one-time cart restore guard
if "cart_restored" not in st.session_state: st.session_state.cart_restored = False

# [ADD THIS LINE]
if "last_deleted" not in st.session_state: st.session_state.last_deleted = None

if not st.session_state.cart_restored:
    _restore_cart_from_disk()
    st.session_state.cart_restored = True

# --- Simple password gates (env-driven) ---

def _env_password(name: str, fallback: str = "") -> str:
    """
    Backward compatible secret lookup using get_secret:
    - Prefer <NAME>_HASH (bcrypt hash from secrets)
    - Otherwise fall back to <NAME> (plain text, legacy)
    """
    hashed = get_secret(f"{name}_HASH")
    if hashed:
        return hashed.strip()
    return (get_secret(name) or fallback).strip()

def _verify_password(plain_password: str, expected_secret: str) -> bool:
    """
    Verify a password against either:
      - a bcrypt hash (preferred, from APP_PASSWORD_HASH / DASHBOARD_PASSWORD_HASH), or
      - a legacy plain-text secret (fallback).
    """
    if not expected_secret:
        return False

    expected_secret = expected_secret.strip()
    plain_password = plain_password.strip()

    # If it looks like a bcrypt hash ($2a$, $2b$, $2y$), use bcrypt
    if expected_secret.startswith("$2a$") or expected_secret.startswith("$2b$") or expected_secret.startswith("$2y$"):
        try:
            return bcrypt.checkpw(
                plain_password.encode("utf-8"),
                expected_secret.encode("utf-8"),
            )
        except Exception:
            return False

    # Fallback: constant-time-ish comparison for legacy plain text
    import hmac
    return hmac.compare_digest(
        hashlib.sha256(plain_password.encode("utf-8")).hexdigest(),
        hashlib.sha256(expected_secret.encode("utf-8")).hexdigest(),
    )


def _login_form(title: str, key_prefix: str, expected_secret: str) -> bool:
    """
    Render a tiny password-only login form.
    Uses _verify_password() so expected_secret can be a bcrypt hash or plain text.
    """
    st.markdown("<div style='height:10px'></div>", unsafe_allow_html=True)
    st.markdown(f"<h3 style='text-align:center'>{title}</h3>", unsafe_allow_html=True)

    with st.form(key=f"{key_prefix}_form", clear_on_submit=False, border=True):
        st.text_input("User", value="admin", disabled=True, key=f"{key_prefix}_user")
        pw = st.text_input("Password", type="password", key=f"{key_prefix}_pw")
        ok = st.form_submit_button("Sign in", use_container_width=True)

    if ok:
        if _verify_password(pw, expected_secret):
            st.session_state[f"__authed__{key_prefix}"] = True
            st.success("Signed in")
            st.rerun()
        else:
            st.error("Incorrect password")

    return bool(st.session_state.get(f"__authed__{key_prefix}", False))

def gate_app_auth() -> bool:
    """
    Gate the entire app.
    - Within a single browser session, we use session_state.
    - Across Streamlit sessions (e.g., when query params change), we also
      remember a successful login in a small process-wide cache keyed by
      the current APP_PASSWORD value.
    """
    exp = _env_password("APP_PASSWORD", "").strip()

    # If no password configured → always allow
    if not exp:
        st.session_state["__authed__app"] = True
        return True

    # Derive a stable token for this password value
    token = sha1(exp.encode("utf-8")).hexdigest()

    # If any prior session has already authenticated for this token,
    # immediately mark this session as authed too.
    if _APP_AUTH_CACHE.get(token):
        st.session_state["__authed__app"] = True
        return True

    # Normal session-level flag (same tab / same session reruns)
    if st.session_state.get("__authed__app", False):
        _APP_AUTH_CACHE[token] = True
        return True

    # Not yet authed → show login form once
    authed = _login_form("Altossa – Sign In", "app", exp)
    if authed:
        st.session_state["__authed__app"] = True
        _APP_AUTH_CACHE[token] = True

    return bool(st.session_state.get("__authed__app", False))

def gate_dashboard_auth() -> bool:
    """
    Gate only the Dashboard page (separate password).
    Uses env DASHBOARD_PASSWORD (from .env). Username is fixed: admin.
    """
    if st.session_state.get("__authed__dashboard", False):
        return True
    exp = _env_password("DASHBOARD_PASSWORD", "")
    if not exp:
        # If no dashboard password configured, allow access
        st.session_state["__authed__dashboard"] = True
        return True
    authed = _login_form("Dashboard – Admin Sign In", "dashboard", exp)
    if authed:
        st.session_state["__authed__dashboard"] = True
    return bool(st.session_state.get("__authed__dashboard", False))

# --- NAV CORE: back/forward + home ---
def go(page: str) -> None:
    """Navigate to a page; clears forward history like a browser."""
    st.session_state.nav.append(st.session_state.page)
    st.session_state.page = page
    # any new navigation clears the forward stack
    st.session_state.fwd = []

def back() -> None:
    """Go to previous page; pushes current page onto forward stack."""
    if st.session_state.nav:
        cur = st.session_state.page
        st.session_state.page = st.session_state.nav.pop()
        # keep forward history
        st.session_state.fwd.append(cur)

def forward() -> None:
    """Go to the next page in history (if any)."""
    fwd = st.session_state.get("fwd", [])
    if fwd:
        cur = st.session_state.page
        st.session_state.page = fwd.pop()          # consume one step forward
        st.session_state.nav.append(cur)           # current becomes a back step
        st.session_state.fwd = fwd                 # persist (already popped)

def go_home() -> None:
    """Jump straight to Home; clears both back and forward history."""
    # if you were on review, callers should have flushed uploads first (see wrappers below)
    st.session_state.nav.clear()
    st.session_state.fwd = []
    st.session_state.page = "home"

def go_to_current_selection() -> None:
    """Jump to the last selection page (e.g., the size chooser for the last product you picked).
    Falls back to Products if no context exists.
    """
    # If leaving Review, persist any uploads first
    if st.session_state.get("page") == "review":
        _flush_review_uploads_to_cart()

    ctx = st.session_state.get("last_selection_ctx")
    if not ctx:
        # No stored context → fall back to product browser
        go("products")
        return

    # Restore the minimal state needed for the target page to render correctly
    st.session_state.sel_type = ctx.get("sel_type")
    st.session_state.sel_product_id = ctx.get("sel_product_id")
    st.session_state.sel_product = ctx.get("sel_product")
    st.session_state.sel_size_idx = None          # ensure we land on the chooser, not a specific size
    st.session_state.compositions_only = bool(ctx.get("compositions_only", False))

    go(ctx.get("page") or "products")

def open_product_type(ptype: str) -> None:
    """
    In-session navigation to a product type without reloading the app.
    Prevents auth prompts by keeping the same Streamlit session.

    NOTE:
    We now force a rerun after changing st.session_state.page so that
    clicking a "View" button on the Home → Product Type tiles navigates
    immediately (no need for a second click).
    """
    # clear any compositions shortcut; set chosen type; go to products
    st.session_state.update(
        sel_room=None,
        sel_type=str(ptype),
        compositions_only=False,
    )
    go("products")
    # Immediately rerun so the new page ("products") is rendered in this click
    st.rerun()

def _consume_type_query_and_jump():
    """
    If the URL has ?type=<Product Type>, set selection and jump to Products once.
    Works whether the query param is a string or a list (Streamlit versions differ).
    Clears the param to avoid loops.
    """
    try:
        q = st.query_params
    except Exception:
        return

    if not q:
        return

    raw = q.get("type")
    # Streamlit may give "Sofa" or ["Sofa"]
    if isinstance(raw, (list, tuple)):
        t = raw[0] if raw else ""
    else:
        t = str(raw) if raw is not None else ""

    t = (t or "").strip()
    if not t:
        return

    # prepare state then clear param and rerun into Products
    st.session_state.update(
        sel_type=t,
        sel_room=None,
        compositions_only=False,
    )

    try:
        st.query_params.clear()
    except Exception:
        # best-effort fallback for older versions
        try:
            st.query_params["type"] = []
        except Exception:
            pass

    go("products")
    st.rerun()

def inject_scroll_restoration(page_key: str) -> None:
    """
    Remember and restore scroll position for a given logical page key.
    We use window.sessionStorage so when you come back to the Products
    page (via Back), it returns to the same scroll offset.
    """
    import json as _json

    js = f"""
    <script>
      const storageKey = "scroll_" + {_json.dumps(page_key)};

      function saveScrollPos() {{
        try {{
          const y = window.scrollY || window.pageYOffset || 0;
          window.sessionStorage.setItem(storageKey, String(y));
        }} catch (e) {{}}
      }}

      function restoreScrollPos() {{
        try {{
          const raw = window.sessionStorage.getItem(storageKey);
          if (!raw) return;
          const y = parseInt(raw, 10);
          if (!isNaN(y)) {{
            window.scrollTo(0, y);
          }}
        }} catch (e) {{}}
      }}

      // On initial load + after Streamlit re-render
      window.addEventListener("load", restoreScrollPos);
      setTimeout(restoreScrollPos, 80);

      // Periodically save while the user scrolls
      setInterval(saveScrollPos, 400);

      // And also right before page unload
      window.addEventListener("beforeunload", saveScrollPos);
    </script>
    """
    st.markdown(js, unsafe_allow_html=True)

def render_nav(prefix: str, on_back, on_forward, on_home, on_products) -> None:
    """
    Navigation row:
      [ Back ] ........................................ [ Home ] [ Go to Products ] [ → Forward ]
    - Updated: Increased width of "Go to Products" + "→ Forward" buttons.
    - No UI changes elsewhere.
    """
    # Three zones: left(back) | spacer | right(cluster)
    left, spacer, right = st.columns([0.16, 0.46, 0.38])   # <-- wider right cluster

    with left:
        st.button("⬅ Back", on_click=on_back,
                  key=f"{prefix}_back",
                  help="Go to previous step",
                  use_container_width=True)

    # Right cluster
    with right:
        # Increase the width ratios: Home smaller, Products & Forward bigger
        c_home, c_prod, c_fwd = st.columns([0.8, 1.3, 1.3], gap="small")

        with c_home:
            st.button("Home",
                      on_click=on_home,
                      key=f"{prefix}_home",
                      use_container_width=True)

        with c_prod:
            st.button("Go to Products",
                      on_click=on_products,
                      key=f"{prefix}_to_products",
                      use_container_width=True)

        with c_fwd:
            st.button("→ Forward",
                      on_click=on_forward,
                      key=f"{prefix}_fwd",
                      use_container_width=True)

# NEW: route depending on whether the product has any full dimensions
def on_select_product(product_id: int, product_name: str):
    """Select a specific product by ID (not just name) and route based on whether it has sizes.
    Also records a 'last selection context' so the sidebar 'Current Selection' button can jump back.
    """
    st.session_state.sel_product_id = int(product_id)
    st.session_state.sel_product = product_name  # keep for display
    sel_pid = int(product_id)

    # Filter rows for the chosen product id
    g = df[df["product_id"] == sel_pid]

    # 'Has full dimensions' means at least one row has W, D, H all present
    has_full_dims = False
    if not g.empty:
        has_full_dims = (g[["W", "D", "H"]].notna().all(axis=1)).any()

    st.session_state.sel_size_idx = None
    target_page = "sizes" if has_full_dims else "materials"

    # --- Record the "last selection" context so we can return here later
    st.session_state.last_selection_ctx = {
        "page": target_page,                           # where to go back to
        "sel_type": st.session_state.get("sel_type"),  # e.g., Bookcase
        "sel_product_id": sel_pid,                     # product id
        "sel_product": product_name,                   # product name (for header)
        "compositions_only": False,                    # be explicit
    }

    go(target_page)

# NEW: safely clear a file_uploader on the *next* run (avoids StreamlitAPIException)
def _mark_uploader_for_clear(uploader_key: str) -> None:
    st.session_state[f"__clear__{uploader_key}"] = True

def _ensure_uploader_cleared(uploader_key: str) -> None:
    """Call this *before* creating the file_uploader widget."""
    flag_key = f"__clear__{uploader_key}"
    if st.session_state.get(flag_key):
        # Remove both the flag and the old uploader state (if any) before instantiating the widget.
        st.session_state.pop(flag_key, None)
        st.session_state.pop(uploader_key, None)

# ---------------------------
# Data (PostgreSQL)
# ---------------------------
DATABASE_URL = get_database_url()
if not DATABASE_URL:
    st.error("Missing DATABASE_URL.")
    st.stop()

df = load_data_from_db(DATABASE_URL)

# Warm image cache once per session (heroes + all composition images)
if not st.session_state.get("warmup_started"):
    start_global_image_warmup(df, max_workers=12)
    st.session_state.warmup_started = True

# --- Global App Auth (one-time, per session)
if not gate_app_auth():
    st.stop()

# ---------------------------
# Sidebar: Selected Products with Quantity  (instant refresh)
# ---------------------------
with st.sidebar:
# [UPDATED] Sidebar Navigation & Undo Menu
    def _nav(to):
        if st.session_state.get("page") == "review":
            _flush_review_uploads_to_cart()
        go(to)

    # Dropdown Menu Container
    with st.expander("Menu", expanded=False):
        st.button("Home", use_container_width=True, on_click=lambda: (_flush_review_uploads_to_cart() if st.session_state.get("page") == "review" else None, go_home()))
        st.button("Dashboard", use_container_width=True, on_click=lambda: _nav("dashboard"))
        st.button("Review Selection", use_container_width=True, on_click=lambda: _nav("review"))
        st.button("Current Selection", use_container_width=True, on_click=go_to_current_selection)
        
        # Undo Button (Only appears if something was deleted)
        if st.session_state.get("last_deleted"):
            st.markdown("---")
            if st.button("↩ Undo Remove", key="undo_btn", use_container_width=True):
                restore_data = st.session_state.last_deleted
                if restore_data:
                    # Restore item to original position or append
                    idx = restore_data["index"]
                    item = restore_data["item"]
                    
                    # Safety check for index
                    if 0 <= idx <= len(st.session_state.cart):
                        st.session_state.cart.insert(idx, item)
                    else:
                        st.session_state.cart.append(item)
                    
                    # Clear undo buffer
                    st.session_state.last_deleted = None
                    st.success("Product Restored!")
                    time.sleep(0.5) # Visual feedback
                    st.rerun()

    st.markdown("---")

    st.title("Selected Products")

    # NEW: remove all items in cart (and clear any saved attachments)
    if st.button(
        "Remove all",
        key="cart_clear_all",
        use_container_width=True,
    ):
        # Clear uploader state + disk assets for every item
        for it in st.session_state.cart:
            uid = it.get("uid")
            if uid:
                _mark_uploader_for_clear(_uploader_key_for(uid))
                _delete_assets_for(uid)
        st.session_state.cart = []
        st.success("All items removed")
        st.rerun()

    total = 0
    keep = []
    
    # [UPDATED] Use stable UIDs for keys instead of index 'i' to prevent data shifting
    for i, item in enumerate(st.session_state.cart):
        # Ensure item has a persistent ID
        uid = _ensure_item_uid(i)
        
        with st.container(border=True):
            st.markdown(f"**{item['Product name']}** · {item['Product type']}")
            st.caption(item.get("Brand", ""))
            if item.get("W") and item.get("D") and item.get("H"):
                st.write(f"**Size** {item['W']} × {item['D']} × {item['H']}")
            if item.get("Material"):
                st.write(f"**Material** {item['Material']}")
            if oc := item.get("Other Category"):
                st.write(f"**Other Category** {oc}")
            if item.get("Name of subtype"):
                st.write(f"**Subtype** {item['Name of subtype']}")
            if link := item.get("Link"):
                st.write(f"[Product link]({link})")

            # --- Preview image in sidebar:
            comp_link = (item.get("Composition link") or "").strip()
            prod_link = (item.get("Link") or "").strip()

            img = None
            if comp_link:
                img = get_image_if_cached("", comp_link) or get_image_cached("", comp_link)
            if not img and prod_link:
                img = get_image_if_cached(prod_link, "") or get_image_cached(prod_link, "")

            if img:
                st.image(img, use_container_width=True)

            comp_parts = (item.get("Composition parts") or "").strip()
            if comp_parts:
                st.caption(f"**Composition:** {comp_parts}")

            # --- Quantity (Keyed by UID)
            qty_key = f"qty_{uid}" 
            current_qty = int(item.get("qty", 1))
            qty = st.number_input(
                "Qty",
                min_value=1,
                max_value=20,
                step=1,
                value=current_qty,
                key=qty_key,
            )
            if qty != current_qty:
                item["qty"] = int(qty)
                st.rerun()
            qty_int = int(item.get("qty", 1) or 1)

            # --- Room allocation per item (Keys updated to use UID)
            st.markdown("**Room allocation**")
            existing_alloc = item.get("room_alloc") or {}

            if qty_int <= 1:
                # Single piece logic
                current_room = None
                try:
                    non_zero: list[tuple[str, int]] = []
                    if isinstance(existing_alloc, dict):
                        for r, v in existing_alloc.items():
                            c = _room_qty(v)
                            if c > 0:
                                non_zero.append((str(r), c))
                except Exception:
                    non_zero = []

                if (
                    len(non_zero) == 1
                    and non_zero[0][1] == 1
                    and non_zero[0][0] in ROOM_CHOICES
                ):
                    current_room = non_zero[0][0]

                options = [ROOM_PLACEHOLDER] + ROOM_CHOICES

                if current_room and current_room in ROOM_CHOICES:
                    default_index = options.index(current_room)
                else:
                    default_index = 0

                # [UPDATED KEY]
                room_sel = st.selectbox(
                    "Room",
                    options,
                    index=default_index,
                    key=f"room_single_{uid}",
                )

                if room_sel == ROOM_PLACEHOLDER:
                    item["room_alloc"] = {}
                    st.caption("Assigned: —")
                else:
                    if room_sel == "Drawing Room":
                        prev_role = None
                        dr_prev = existing_alloc.get("Drawing Room")
                        if isinstance(dr_prev, dict):
                            prev_role = dr_prev.get("visual_role")

                        vr_options = [
                            VISUAL_ROLE_PLACEHOLDER,
                            "Sofa", "Armchair", "Centre Table",
                            "Side Table", "Insert Table", "Carpet", "Others",
                        ]

                        selected_vr = st.selectbox(
                            "Drawing Room Visual Role",
                            vr_options,
                            index=vr_options.index(prev_role) if prev_role in vr_options else 0,
                            key=f"vr_single_{uid}",
                        )

                        if selected_vr == VISUAL_ROLE_PLACEHOLDER:
                            dr_alloc = {"qty": qty_int}
                        else:
                            dr_alloc = {"qty": qty_int, "visual_role": selected_vr}

                        item["room_alloc"] = {"Drawing Room": dr_alloc}

                    elif room_sel == "Living Room":
                        prev_role = None
                        lr_prev = existing_alloc.get("Living Room")
                        if isinstance(lr_prev, dict):
                            prev_role = lr_prev.get("visual_role")

                        vr_options = [
                            VISUAL_ROLE_PLACEHOLDER,
                            "Sofa", "Armchair", "Chest of Tables", 
                            "Side Tables", "TV Unit", "Carpets", "Others",
                        ]

                        selected_vr = st.selectbox(
                            "Living Room Visual Role",
                            vr_options,
                            index=vr_options.index(prev_role) if prev_role in vr_options else 0,
                            key=f"vr_single_lr_{uid}",
                        )

                        if selected_vr == VISUAL_ROLE_PLACEHOLDER:
                            lr_alloc = {"qty": qty_int}
                        else:
                            lr_alloc = {"qty": qty_int, "visual_role": selected_vr}

                        item["room_alloc"] = {"Living Room": lr_alloc}

                    # [UPDATED] Generic Bedroom Logic (Handles 1, 2, 3)
                    elif room_sel in ["Bedroom 1", "Bedroom 2", "Bedroom 3"]:
                        prev_role = None
                        br_prev = existing_alloc.get(room_sel)
                        if isinstance(br_prev, dict):
                            prev_role = br_prev.get("visual_role")

                        vr_options = [
                            VISUAL_ROLE_PLACEHOLDER,
                            "Bed", "Night Tables", "Chest of Drawers", "Wardrobe", "Others",
                        ]

                        selected_vr = st.selectbox(
                            f"{room_sel} Visual Role",
                            vr_options,
                            index=vr_options.index(prev_role) if prev_role in vr_options else 0,
                            key=f"vr_single_{room_sel.replace(' ', '')}_{uid}",
                        )

                        if selected_vr == VISUAL_ROLE_PLACEHOLDER:
                            br_alloc = {"qty": qty_int}
                        else:
                            br_alloc = {"qty": qty_int, "visual_role": selected_vr}

                        item["room_alloc"] = {room_sel: br_alloc}

                    # [NEW] Dining Room Logic
                    elif room_sel == "Dining Room":
                        prev_role = None
                        dr_prev = existing_alloc.get("Dining Room")
                        if isinstance(dr_prev, dict):
                            prev_role = dr_prev.get("visual_role")

                        vr_options = [
                            VISUAL_ROLE_PLACEHOLDER,
                            "Dining Table", "Chairs", "Others",
                        ]

                        selected_vr = st.selectbox(
                            "Dining Room Visual Role",
                            vr_options,
                            index=vr_options.index(prev_role) if prev_role in vr_options else 0,
                            key=f"vr_single_din_{uid}",
                        )

                        if selected_vr == VISUAL_ROLE_PLACEHOLDER:
                            din_alloc = {"qty": qty_int}
                        else:
                            din_alloc = {"qty": qty_int, "visual_role": selected_vr}

                        item["room_alloc"] = {"Dining Room": din_alloc}

                    else:
                        item["room_alloc"] = {room_sel: qty_int}

                    st.caption(f"Assigned: {room_sel} × {qty_int}")

            else:
                # Multi-quantity logic
                with st.expander(
                    "Room allocation (optional)",
                    expanded=bool(existing_alloc),
                ):
                    st.caption("Set how many units go to each room.")
                    slider_vals: dict[str, int] = {}
                    col_left, col_right = st.columns(2)

                    for idx, room in enumerate(ROOM_CHOICES):
                        col = col_left if idx % 2 == 0 else col_right
                        with col:
                            default_val = 0
                            if isinstance(existing_alloc, dict):
                                default_val = _room_qty(existing_alloc.get(room, 0))

                            if default_val < 0: default_val = 0
                            if default_val > qty_int: default_val = qty_int

                            # [UPDATED KEY]
                            slider_vals[room] = st.number_input(
                                room,
                                min_value=0,
                                max_value=qty_int,
                                step=1,
                                value=default_val,
                                key=f"room_{uid}_{room}",
                            )

                    eff_alloc_ints = {r: int(v) for r, v in slider_vals.items() if int(v) > 0}
                    total_assigned = sum(eff_alloc_ints.values()) if eff_alloc_ints else 0
                    
                    st.caption(f"Total: {total_assigned} / {qty_int}")
                    if total_assigned != qty_int:
                        st.warning("Total mismatch", icon="⚠️")

                    prev_role = None
                    dr_prev_raw = None
                    if isinstance(existing_alloc, dict):
                        dr_prev_raw = existing_alloc.get("Drawing Room")
                        if isinstance(dr_prev_raw, dict):
                            prev_role = dr_prev_raw.get("visual_role")

                    # Determine previous roles for pre-selection
                    dr_prev_role = None
                    lr_prev_role = None
                    br_prev_role = None
                    din_prev_role = None
                    
                    if isinstance(existing_alloc, dict):
                        # Drawing Room prev
                        dr_prev_raw = existing_alloc.get("Drawing Room")
                        if isinstance(dr_prev_raw, dict):
                            dr_prev_role = dr_prev_raw.get("visual_role")
                        # Living Room prev
                        lr_prev_raw = existing_alloc.get("Living Room")
                        if isinstance(lr_prev_raw, dict):
                            lr_prev_role = lr_prev_raw.get("visual_role")
                        # [NEW] Bedroom prev
                        br_prev_raw = existing_alloc.get("Bedroom")
                        if isinstance(br_prev_raw, dict):
                            br_prev_role = br_prev_raw.get("visual_role")
                        # [NEW] Dining Room prev
                        din_prev_raw = existing_alloc.get("Dining Room")
                        if isinstance(din_prev_raw, dict):
                            din_prev_role = din_prev_raw.get("visual_role")

                    new_alloc: dict[str, object] = {}

                    # 1. Handle Drawing Room
                    dr_qty = eff_alloc_ints.get("Drawing Room", 0)
                    if dr_qty > 0:
                        vr_options = [
                            VISUAL_ROLE_PLACEHOLDER,
                            "Sofa", "Armchair", "Centre Table",
                            "Side Table", "Insert Table", "Carpet", "Others",
                        ]
                        selected_vr = st.selectbox(
                            "Drawing Room Visual Role",
                            vr_options,
                            index=vr_options.index(dr_prev_role) if dr_prev_role in vr_options else 0,
                            key=f"vr_multi_dr_{uid}",
                        )
                        if selected_vr == VISUAL_ROLE_PLACEHOLDER:
                            new_alloc["Drawing Room"] = {"qty": dr_qty}
                        else:
                            new_alloc["Drawing Room"] = {"qty": dr_qty, "visual_role": selected_vr}

                    # 2. Handle Living Room
                    lr_qty = eff_alloc_ints.get("Living Room", 0)
                    if lr_qty > 0:
                        vr_options = [
                            VISUAL_ROLE_PLACEHOLDER,
                            "Sofa", "Armchair", "Chest of Tables", 
                            "Side Tables", "TV Unit", "Carpets", "Others",
                        ]
                        selected_vr = st.selectbox(
                            "Living Room Visual Role",
                            vr_options,
                            index=vr_options.index(lr_prev_role) if lr_prev_role in vr_options else 0,
                            key=f"vr_multi_lr_{uid}",
                        )
                        if selected_vr == VISUAL_ROLE_PLACEHOLDER:
                            new_alloc["Living Room"] = {"qty": lr_qty}
                        else:
                            new_alloc["Living Room"] = {"qty": lr_qty, "visual_role": selected_vr}

                    # 3. [UPDATED] Handle Bedrooms (1, 2, 3)
                    for br_key in ["Bedroom 1", "Bedroom 2", "Bedroom 3"]:
                        br_qty = eff_alloc_ints.get(br_key, 0)
                        if br_qty > 0:
                            # Retrieve previous role dynamically
                            prev_r = None
                            if isinstance(existing_alloc, dict):
                                raw = existing_alloc.get(br_key)
                                if isinstance(raw, dict):
                                    prev_r = raw.get("visual_role")

                            vr_options = [
                                VISUAL_ROLE_PLACEHOLDER,
                                "Bed", "Night Tables", "Chest of Drawers", "Wardrobe", "Others",
                            ]
                            selected_vr = st.selectbox(
                                f"{br_key} Visual Role",
                                vr_options,
                                index=vr_options.index(prev_r) if prev_r in vr_options else 0,
                                key=f"vr_multi_{br_key.replace(' ', '')}_{uid}",
                            )
                            if selected_vr == VISUAL_ROLE_PLACEHOLDER:
                                new_alloc[br_key] = {"qty": br_qty}
                            else:
                                new_alloc[br_key] = {"qty": br_qty, "visual_role": selected_vr}

                    # 4. [NEW] Handle Dining Room
                    din_qty = eff_alloc_ints.get("Dining Room", 0)
                    if din_qty > 0:
                        vr_options = [
                            VISUAL_ROLE_PLACEHOLDER,
                            "Dining Table", "Chairs", "Others",
                        ]
                        selected_vr = st.selectbox(
                            "Dining Room Visual Role",
                            vr_options,
                            index=vr_options.index(din_prev_role) if din_prev_role in vr_options else 0,
                            key=f"vr_multi_din_{uid}",
                        )
                        if selected_vr == VISUAL_ROLE_PLACEHOLDER:
                            new_alloc["Dining Room"] = {"qty": din_qty}
                        else:
                            new_alloc["Dining Room"] = {"qty": din_qty, "visual_role": selected_vr}

                    # 5. Handle Other Rooms
                    for room_name, q in eff_alloc_ints.items():
                        if room_name in ["Drawing Room", "Living Room", "Dining Room", "Bedroom 1", "Bedroom 2", "Bedroom 3"]: continue
                        if q > 0: new_alloc[room_name] = int(q)

                    item["room_alloc"] = new_alloc

            # --- Pricing summary
            fp = calculate_final_price(item.get("Brand", ""), item.get("Base Price"))
            if fp is not None:
                st.metric("Unit Price", f"₹{fp:,}")
                total += fp * item["qty"]
                st.caption(f"Subtotal: ₹{(fp * item['qty']):,}")

            c1, c2 = st.columns(2)
            # [UPDATED KEYS]
            if c1.button("Edit", key=f"edit_{uid}"):
                st.session_state.sel_type = item["Product type"]
                st.session_state.sel_product = item["Product name"]
                go("products")
                st.rerun()

            # [UPDATED] Remove button with Undo support
            if c2.button("Remove", key=f"rm_{uid}"):
                # 1. Save for Undo (Store index and item data)
                st.session_state.last_deleted = {
                    "index": i,
                    "item": item  # Store the dictionary reference
                }
                
                # 2. Clear uploader widget state, BUT DO NOT delete actual asset files 
                #    (_delete_assets_for) so they reappear if Undone.
                if uid:
                    _mark_uploader_for_clear(_uploader_key_for(uid))
                    # _delete_assets_for(uid) <--- DISABLED TO ALLOW UNDO RESTORE
                
                st.success("Removed")
                st.session_state.cart.pop(i)
                st.rerun()

            keep.append(item)

    st.session_state.cart = keep
    st.divider()
    st.metric("Grand Total", f"₹{total:,}")

    # --- Client name (required for exports) ---
    client_name = st.text_input("Client name", key="client_name_sidebar")
    approved_by = st.text_input("Approved by (team member)", key="approved_by_sidebar")

    export_df = _cart_to_dataframe(st.session_state.cart)

    # NEW: validate room allocations for multi-quantity items
    allocation_violations = get_room_allocation_violations(st.session_state.cart)
    if allocation_violations:
        st.error(
            "Room allocation mismatch for some multi-quantity products. "
            "Please adjust the room splits so that the total assigned per product "
            "matches its quantity before exporting:\n\n"
            + "\n".join(f"- {msg}" for msg in allocation_violations)
        )

    # Export is controlled only by client name + non-empty selection.
    # Room allocation issues are shown as an error but do not 'lock' the buttons,
    # so you won't have the situation where the UI looks correct but the button
    # stays disabled until another random action.
    allow_export = (
        (client_name.strip() != "")
        and (not export_df.empty)
        and (not allocation_violations) 
    )

    # [UPDATED] Export Menu Dropdown
    # We replace the static "### Export" header with a collapsible menu
    with st.expander("Export", expanded=False):
        
        # Excel
        if allow_export:
            xlsx_bytes = _excel_bytes_from_df(export_df)
            st.download_button(
                "⬇ Excel",
                data=xlsx_bytes,
                file_name="altossa_selection.xlsx",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                width="stretch",
            )
        else:
            st.button("⬇ Excel", disabled=True, width="stretch")

        # PDF
        if allow_export:
            pdf_bytes = _pdf_bytes_from_df(
                export_df,
                client_name=client_name.strip(),
                approved_by=approved_by.strip(),
            )
            if pdf_bytes is not None:
                st.download_button(
                    "⬇ PDF",
                    data=pdf_bytes,
                    file_name="altossa_selection.pdf",
                    mime="application/pdf",
                    width="stretch",
                )
            else:
                st.info("Install `reportlab` to enable PDF export: pip install reportlab")
                st.button("⬇ PDF", disabled=True, width="stretch")
        else:
            st.button("⬇ PDF", disabled=True, width="stretch")

        # PPT
        if allow_export:
            ppt_bytes = _ppt_bytes_from_cart(
                st.session_state.cart,
                client_name=client_name.strip(),
            )
            if ppt_bytes is not None:
                st.download_button(
                    "⬇ PPT (internal)",
                    data=ppt_bytes,
                    file_name="altossa_internal_notes.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    width="stretch",
                )
            else:
                st.info("Install `python-pptx` to enable PPT export: pip install python-pptx")
                st.button("⬇ PPT (internal)", disabled=True, width="stretch")
        else:
            st.button("⬇ PPT (internal)", disabled=True, width="stretch")

    # --- NEW: autosave selection & header info to disk so it survives sleep/restarts
    _autosave_cart_if_changed()

def _filter_variant_rows(df: pd.DataFrame, base_row: pd.Series) -> pd.DataFrame:
    """
    Return only the rows that match the selected variant exactly:
      Brand, Product type, Product name, W, D, H, and Name of subtype.
    (If Name of subtype is empty, we fall back to dimensions-only.)
    """
    # Safe getters
    brand   = str(base_row.get("Brand", "")).strip()
    ptype   = str(base_row.get("Product type", "")).strip()
    pname   = str(base_row.get("Product name", "")).strip()
    subtype = str(base_row.get("Name of subtype", "")).strip()

    W = base_row.get("W", None)
    D = base_row.get("D", None)
    H = base_row.get("H", None)

    # Build the common mask
    m = (
        (df["Brand"].astype(str).str.strip() == brand) &
        (df["Product type"].astype(str).str.strip() == ptype) &
        (df["Product name"].astype(str).str.strip() == pname) &
        (df["W"] == W) & (df["D"] == D) & (df["H"] == H)
    )

    # If subtype is present, add it to the filter; otherwise use only dimensions
    if subtype:
        m = m & (df["Name of subtype"].astype(str).str.strip() == subtype)

    return df.loc[m].copy()

def _sha1(b: bytes) -> str:
    return hashlib.sha1(b).hexdigest()

def _append_unique_assets(item: dict, new_files: list[bytes], cap: int = 3) -> None:
    """
    Keep at most 'cap' images. Avoid duplicates by SHA1. Mutates 'item'.
    Stores as list of raw bytes in key 'internal_assets_bytes'.
    """
    existing = item.get("internal_assets_bytes") or []
    if isinstance(existing, bytes):  # backward compat
        existing = [existing]
    # current hashes
    seen = { _sha1(x) for x in existing }
    for nb in new_files:
        h = _sha1(nb)
        if h not in seen:
            existing.append(nb); seen.add(h)
        if len(existing) >= cap:
            break
    item["internal_assets_bytes"] = existing[:cap]

# NEW: keep attached images in sync with the file_uploader (so ✕ removes instantly)
def _sync_uploader_assets(idx: int, uploader_key: str, cap: int = 3) -> bool:
    """
    Keep the central attachment store (uid → bytes[]) and cart[idx] in sync with the
    visible chips of the file_uploader. Returns True if changed on THIS run.
    Uses a *per-line UID* only (no derived stable key), so attachments never move/disappear.
    """
    files = st.session_state.get(uploader_key)  # None, [] or list[UploadedFile]
    if files is None:
        return False  # widget not instantiated yet

    # Read up to 'cap' files (as bytes)
    new_bytes: list[bytes] = []
    for f in files[:cap]:
        try:
            new_bytes.append(f.getvalue())
        except Exception:
            pass

    # Compute signature of uploader payload
    import hashlib
    def _sig(lst: list[bytes]) -> str:
        return "|".join(hashlib.sha1(b).hexdigest() for b in lst)

    new_sig = _sig(new_bytes)
    memo_key = f"__uploader_sig__{uploader_key}"
    old_sig = st.session_state.get(memo_key, "")

    if new_sig != old_sig:
        # Map row index to a *stable per-line uid*
        uid = _ensure_item_uid(idx)

        # 1) write into central store (and disk)
        _set_assets_for(uid, new_bytes, cap=cap)

        # 2) mirror into cart row (for exports / PPT / PDF without re-reading chips)
        st.session_state.cart[idx]["internal_assets_bytes"] = _get_assets_for(uid)

        # 3) memo + refresh now so preview appears immediately
        st.session_state[memo_key] = new_sig
        st.rerun()
        return True

    return False

def _snapshot_uploader_to_cart(idx: int, uploader_key: str, cap: int = 3) -> None:
    """
    Snapshot current uploader chips into the central store and mirror into cart
    WITHOUT calling st.rerun(). Used before navigation.
    """
    files = st.session_state.get(uploader_key)
    if files is None:
        return

    new_bytes: list[bytes] = []
    for f in files[:cap]:
        try:
            new_bytes.append(f.getvalue())
        except Exception:
            pass

    uid = _ensure_item_uid(idx)
    _set_assets_for(uid, new_bytes, cap=cap)
    st.session_state.cart[idx]["internal_assets_bytes"] = _get_assets_for(uid)

# NEW: when navigating away from the Review page (e.g., Back), make sure
# any uploader chips currently visible are written into the cart so they
# survive page changes.
def _flush_review_uploads_to_cart() -> None:
    """
    Snapshot all visible uploader chips into cart so they survive navigation.
    IMPORTANT: never calls st.rerun() here (safe for on_click navigation).
    """
    if st.session_state.get("page") != "review" or not st.session_state.get("cart"):
        return

    for idx, it in enumerate(st.session_state.cart):
        uid = it.get("uid") or _ensure_item_uid(idx)
        uploader_key = _uploader_key_for(uid)
        if uploader_key in st.session_state:
            try:
                _snapshot_uploader_to_cart(idx, uploader_key, cap=3)
            except Exception:
                pass  # don't block navigation

# -------------------------------------------------------------------
#  DRAWING ROOM VISUAL EXPORT MODULE (Finalized v9 - Long Vertical Lines)
# -------------------------------------------------------------------

DRAWING_TEMPLATE_PATH = os.path.join(os.getcwd(), "template", "drawing_room_base.png")
LIVING_TEMPLATE_PATH = os.path.join(os.getcwd(), "template", "living_room_base.png")
BEDROOM_TEMPLATE_PATH = os.path.join(os.getcwd(), "template", "bedroom_base.png")
DINING_TEMPLATE_PATH = os.path.join(os.getcwd(), "template", "dining_room_base.png")

def _get_item_pdf_data(item: dict, visual_role: str) -> list[list[str]]:
    """
    Prepares data rows for a Table.
    Format: [Label, ":", Value]
    """
    brand = nz_str(item.get("Brand"))
    ptype = nz_str(item.get("Product type"))
    pname = nz_str(item.get("Product name"))
    subtype = nz_str(item.get("Name of subtype"))
    
    # 1. Brand
    brand_val = f"{brand} - {visual_role}" if brand else visual_role

    # 2. Model
    model_parts = [p for p in [pname, ptype, subtype] if p]
    model_val = " - ".join(model_parts) if model_parts else "—"

    # 3. Dimensions
    dims = " × ".join(
        nz_str(v) for v in [item.get("W"), item.get("D"), item.get("H")] if nz_str(v)
    )
    dim_val = f"{dims} cm" if dims else "—"

    # 4. Quantity
    qty = int(item.get("qty", 1) or 1)

    # 5. Material
    mat = nz_str(item.get("Material"))
    other_cat = nz_str(item.get("Other Category"))
    
    mat_val = "—"
    if mat:
        mat_val = mat 
    elif other_cat:
        mat_val = f"Other - {other_cat}"

    # 6. Price (Use "Rs." text to avoid font issues)
    unit_price = _get_item_price(item)  # [UPDATED]
    total_price = unit_price * qty
    price_val = f"Rs. {total_price:,.0f}"

    # Return raw strings; they will be wrapped in Paragraphs later
    return [
        ["Brand", ":", brand_val],
        ["Model", ":", model_val],
        ["Dimension", ":", dim_val],
        ["Quantity", ":", str(qty)],
        ["Material", ":", mat_val],
        ["Price", ":", price_val],
    ]

def export_drawing_room_pdf(cart: list[dict], client_name: str = "", approved_by: str = "") -> Optional[bytes]:
    """
    Generates a Drawing Room visual layout PDF.
    - Page 1: Diagram + Header + Border.
    - Page 2+: Grid layout 4x2 with consistent borders on ALL pages.
    """
    from datetime import datetime

    # --- Configuration ---
    page_size = landscape(A4)
    WIDTH, HEIGHT = page_size 
    
    # --- LAYOUT SHIFT ---
    Y_SHIFT = 20  

    # --- LAYOUT COORDINATES ---
    BOX_W, BOX_H = 225, 125  
    IMG_W, IMG_H = 95, 95    

    # Slot Definitions
    SLOT_CONFIG = {
        "Sofa":         {"box": (575, 390 - Y_SHIFT), "anchor": (515, 330 - Y_SHIFT), "align": "right"},
        "Carpet":       {"box": (575, 240 - Y_SHIFT), "anchor": (515, 295 - Y_SHIFT), "align": "right"},
        "Insert Table": {"box": (575, 90 - Y_SHIFT),  "anchor": (515, 260 - Y_SHIFT), "align": "right"},
        "Armchair":     {"box": (40, 390 - Y_SHIFT),  "anchor": (325, 330 - Y_SHIFT), "align": "left"},
        "Centre Table": {"box": (40, 240 - Y_SHIFT),  "anchor": (325, 295 - Y_SHIFT), "align": "left"},
        "Side Table":   {"box": (40, 90 - Y_SHIFT),   "anchor": (325, 260 - Y_SHIFT), "align": "left"},
    }
    
    # --- Helpers ---
    def _get_img_bytes(item):
        c = nz_str(item.get("Composition link"))
        l = nz_str(item.get("Link"))
        img = None
        if c: img = get_image_if_cached("", c) or get_image_cached("", c)
        if not img and l: img = get_image_if_cached(l, "") or get_image_cached(l, "")
        return img

    def _draw_elbow_line(canvas, start_x, start_y, box_x, box_y, box_w, box_h, align):
        canvas.saveState()
        canvas.setStrokeColor(black)
        canvas.setLineWidth(1)
        canvas.setDash([3, 2])
        box_mid_y = box_y + (box_h / 2)
        p = canvas.beginPath()
        p.moveTo(start_x, start_y)
        if align == "right":
            end_x = box_x
            mid_x = start_x + 30 if abs(start_y - box_mid_y) > 50 else start_x + 15
            p.lineTo(mid_x, start_y)
            p.lineTo(mid_x, box_mid_y)
            p.lineTo(end_x, box_mid_y)
        else: 
            end_x = box_x + box_w
            mid_x = start_x - 30 if abs(start_y - box_mid_y) > 50 else start_x - 15
            p.lineTo(mid_x, start_y)
            p.lineTo(mid_x, box_mid_y)
            p.lineTo(end_x, box_mid_y)
        canvas.drawPath(p, stroke=1, fill=0)
        canvas.setDash([]) 
        canvas.setFillColor(black)
        canvas.circle(start_x, start_y, 2.5, fill=1, stroke=0)
        canvas.restoreState()

    # --- Header & Border Renderer ---
    def draw_border_and_header(canvas, doc, is_page_1=False):
        canvas.saveState()
        
        # 1. Global Page Border
        canvas.setStrokeColor(HexColor("#000000"))
        canvas.setLineWidth(1.5)
        canvas.rect(10, 10, WIDTH - 20, HEIGHT - 20)

        # 2. Page 1 Header Table
        if is_page_1:
            date_str = datetime.now().strftime("%d-%m-%Y")
            c_name = client_name.strip() or ""
            app_by = approved_by.strip() or "Altossa"
            
            s_style = ParagraphStyle('Hdr', fontName='Helvetica', fontSize=10, leading=12)
            
            row_data = [
                Paragraph(f"<b>Client :</b> {c_name}", s_style),
                Paragraph(f"<b>Approved by :</b> {app_by}", s_style),
                Paragraph(f"<b>Room :</b> Drawing room", s_style),
                Paragraph(f"<b>Date :</b> {date_str}", s_style)
            ]
            
            t_width = WIDTH - 40
            col_w = t_width / 4.0
            
            t = Table([row_data], colWidths=[col_w]*4, rowHeights=[35])
            t.setStyle(TableStyle([
                ('GRID', (0,0), (-1,-1), 1.5, colors.black),
                ('VALIGN', (0,0), (-1,-1), 'MIDDLE'),
                ('LEFTPADDING', (0,0), (-1,-1), 8),
                ('RIGHTPADDING', (0,0), (-1,-1), 8),
            ]))
            
            w, h = t.wrap(t_width, 100)
            t.drawOn(canvas, 20, HEIGHT - 10 - h - 10)

        canvas.restoreState()

    # --- Filter Content ---
    fixed_slots = {k: [] for k in SLOT_CONFIG.keys()}
    others_items = []
    
    for item in cart:
        # 1. Check allocation specifically for Drawing Room
        raw_alloc = item.get("room_alloc", {})
        dr_data = raw_alloc.get("Drawing Room")
        
        if not dr_data:
            continue

        dr_qty = 0
        visual_role = "Others"

        if isinstance(dr_data, dict):
            dr_qty = int(dr_data.get("qty", 0))
            visual_role = dr_data.get("visual_role") or "Others"
        else:
            try: dr_qty = int(dr_data)
            except: dr_qty = 0
        
        if dr_qty <= 0:
            continue

        # Create a PDF-specific item copy with room-specific qty
        pdf_item = item.copy()
        pdf_item["qty"] = dr_qty

        if visual_role in fixed_slots:
            if not fixed_slots[visual_role]:
                fixed_slots[visual_role].append(pdf_item)
            else:
                others_items.append(pdf_item)
        else:
            others_items.append(pdf_item)

    if fixed_slots.get("Others"):
        others_items.extend(fixed_slots["Others"])
        fixed_slots["Others"] = [] 

    # --- Setup ReportLab ---
    buffer = BytesIO()
    doc = SimpleDocTemplate(
        buffer, pagesize=page_size,
        leftMargin=0, rightMargin=0, topMargin=0, bottomMargin=0
    )
    styles = getSampleStyleSheet()
    
    tbl_txt_style = ParagraphStyle('Compact', parent=styles['Normal'], fontName='Helvetica', fontSize=6.5, leading=7)
    tbl_lbl_style = ParagraphStyle('CompactBold', parent=styles['Normal'], fontName='Helvetica-Bold', fontSize=6.5, leading=7)
    tbl_style = TableStyle([
        ('VALIGN', (0,0), (-1,-1), 'TOP'),
        ('LEFTPADDING', (0,0), (-1,-1), 0),
        ('RIGHTPADDING', (0,0), (-1,-1), 0),
        ('BOTTOMPADDING', (0,0), (-1,-1), 0),
        ('TOPPADDING', (0,0), (-1,-1), 0),
    ])

    # --- Page 1: Visualizer Renderer ---
    def draw_visualizer_bg(canvas, doc):
        canvas.saveState()
        if os.path.exists(DRAWING_TEMPLATE_PATH):
            try:
                canvas.drawImage(DRAWING_TEMPLATE_PATH, 0, -Y_SHIFT, width=WIDTH, height=HEIGHT, preserveAspectRatio=True, anchor='c')
            except: pass
        canvas.restoreState()

        draw_border_and_header(canvas, doc, is_page_1=True)
        
        canvas.saveState()
        for role, items in fixed_slots.items():
            if items:
                item = items[0]
                config = SLOT_CONFIG[role]
                box_x, box_y = config["box"]
                _draw_elbow_line(canvas, config["anchor"][0], config["anchor"][1], box_x, box_y, BOX_W, BOX_H, config["align"])

                canvas.setStrokeColor(HexColor("#333333"))
                canvas.setFillColor(colors.white)
                canvas.setLineWidth(1)
                canvas.roundRect(box_x, box_y, BOX_W, BOX_H, 8, fill=1, stroke=1)

                img_bytes = _get_img_bytes(item)
                img_margin = 8
                if img_bytes:
                    try:
                        rdr = ImageReader(BytesIO(img_bytes))
                        iw, ih = rdr.getSize()
                        scale = min(IMG_W / float(iw), IMG_H / float(ih))
                        w, h = iw * scale, ih * scale
                        img_y_pos = box_y + (BOX_H - h) / 2
                        canvas.drawImage(rdr, box_x + img_margin, img_y_pos, width=w, height=h, mask='auto')
                    except: pass

                table_data = _get_item_pdf_data(item, role)
                processed_data = []
                for row in table_data:
                    label, sep, val = row
                    processed_data.append([
                        Paragraph(f"{label}", tbl_lbl_style),
                        Paragraph(f"{sep}", tbl_txt_style),
                        Paragraph(f"{val}", tbl_txt_style)
                    ])

                text_x_start = box_x + IMG_W + (img_margin * 2)
                text_w_avail = BOX_W - IMG_W - (img_margin * 3)
                t = Table(processed_data, colWidths=[45, 5, text_w_avail - 55])
                t.setStyle(tbl_style)
                
                f = Frame(text_x_start, box_y, text_w_avail, BOX_H, leftPadding=0, bottomPadding=5, rightPadding=2, topPadding=5, showBoundary=0)
                f.addFromList([t], canvas)
        canvas.restoreState()

    # --- Page 2+: List Renderer ---
    def draw_list_header(canvas, doc):
        # This ensures borders are drawn on every single page generated by this template
        draw_border_and_header(canvas, doc, is_page_1=False)
        
        canvas.saveState()
        canvas.setFont("Helvetica-Bold", 12)
        canvas.drawString(1.5*cm, HEIGHT - 2.5*cm, "Drawing Room - Additional Items")
        canvas.line(1.5*cm, HEIGHT - 2.5*cm - 5, WIDTH - 1.5*cm, HEIGHT - 2.5*cm - 5)
        canvas.restoreState()

    frame_p1 = Frame(0, 0, WIDTH, HEIGHT, id='f1', showBoundary=0)
    tpl_p1 = PageTemplate(id='Visualizer', frames=[frame_p1], onPage=draw_visualizer_bg)

    frame_p2 = Frame(1.5*cm, 1.5*cm, WIDTH-3*cm, HEIGHT-4.0*cm, id='f2', showBoundary=0)
    tpl_p2 = PageTemplate(id='List', frames=[frame_p2], onPage=draw_list_header)

    doc.addPageTemplates([tpl_p1, tpl_p2])

    # --- Build Story ---
    story = []
    story.append(Spacer(1, 1))

    if others_items:
        story.append(NextPageTemplate('List'))
        story.append(PageBreak())
        
        grid_data = []
        current_row = []
        col_w = (WIDTH - 3*cm - 20) / 2
        
        for item in others_items:
            raw_data = _get_item_pdf_data(item, "Others")
            tbl_rows = []
            for row in raw_data:
                label, sep, val = row
                p_val = Paragraph(val, styles["Normal"])
                p_lbl = Paragraph(f"<b>{label}</b>", styles["Normal"])
                tbl_rows.append([p_lbl, sep, p_val])

            inner_tbl = Table(tbl_rows, colWidths=[65, 5, col_w - 90 - 70]) 
            inner_tbl.setStyle(TableStyle([
                ('VALIGN', (0,0), (-1,-1), 'TOP'),
                ('LEFTPADDING', (0,0), (-1,-1), 0),
                ('RIGHTPADDING', (0,0), (-1,-1), 0),
                ('TOPPADDING', (0,0), (-1,-1), 0),
                ('BOTTOMPADDING', (0,0), (-1,-1), 0),
            ]))

            img_bytes = _get_img_bytes(item)
            img_obj = Spacer(80, 60) 
            if img_bytes:
                try:
                    rdr = ImageReader(BytesIO(img_bytes))
                    iw, ih = rdr.getSize()
                    scale = min(80.0/iw, 60.0/ih)
                    img_obj = RLImage(BytesIO(img_bytes), width=iw*scale, height=ih*scale)
                except: pass
            
            item_data = [[img_obj, inner_tbl]]
            item_tbl = Table(item_data, colWidths=[90, col_w - 90])
            item_tbl.setStyle(TableStyle([
                ('BOX', (0,0), (-1,-1), 0.5, HexColor("#cccccc")),
                ('VALIGN', (0,0), (-1,-1), 'MIDDLE'),
                ('LEFTPADDING', (0,0), (-1,-1), 4),
                ('RIGHTPADDING', (0,0), (-1,-1), 4),
                ('TOPPADDING', (0,0), (-1,-1), 4),
                ('BOTTOMPADDING', (0,0), (-1,-1), 4),
            ]))
            
            current_row.append(item_tbl)
            
            if len(current_row) == 2:
                grid_data.append(current_row)
                current_row = []
        
        if current_row:
            current_row.append(Spacer(1,1))
            grid_data.append(current_row)
            
        main_tbl = Table(grid_data, colWidths=[col_w, col_w])
        main_tbl.setStyle(TableStyle([
            ('VALIGN', (0,0), (-1,-1), 'TOP'),
            ('LEFTPADDING', (0,0), (-1,-1), 0),
            ('RIGHTPADDING', (0,0), (-1, -1), 10),
            ('BOTTOMPADDING', (0,0), (-1,-1), 5),
        ]))
        
        story.append(main_tbl)

    try:
        doc.build(story)
        buffer.seek(0)
        return buffer.read()
    except Exception as e:
        st.error(f"PDF Generation Error: {e}")
        return None

def export_living_room_pdf(cart: list[dict], client_name: str = "", approved_by: str = "") -> Optional[bytes]:
    """
    Generates a Living Room visual layout PDF.
    - Template: living_room_base.png
    - Layout: 3 Items Left, 3 Items Right.
    - Slot Config specific to Living Room roles.
    """
    from datetime import datetime

    # --- Configuration ---
    page_size = landscape(A4)
    WIDTH, HEIGHT = page_size 
    
    # --- LAYOUT SHIFT ---
    Y_SHIFT = 20  

    # --- LAYOUT COORDINATES ---
    BOX_W, BOX_H = 225, 125  
    IMG_W, IMG_H = 95, 95    

    # Slot Definitions for Living Room
    # Right Side: Sofa (Top), TV Unit (Mid), Chest of Tables (Bot)
    # Left Side: Armchair (Top), Side Tables (Mid), Carpets (Bot)
    SLOT_CONFIG = {
        "Sofa":            {"box": (575, 390 - Y_SHIFT), "anchor": (515, 330 - Y_SHIFT), "align": "right"},
        "TV Unit":         {"box": (575, 240 - Y_SHIFT), "anchor": (515, 295 - Y_SHIFT), "align": "right"},
        "Chest of Tables": {"box": (575, 90 - Y_SHIFT),  "anchor": (515, 260 - Y_SHIFT), "align": "right"},
        
        "Armchair":        {"box": (40, 390 - Y_SHIFT),  "anchor": (325, 330 - Y_SHIFT), "align": "left"},
        "Side Tables":     {"box": (40, 240 - Y_SHIFT),  "anchor": (325, 295 - Y_SHIFT), "align": "left"},
        "Carpets":         {"box": (40, 90 - Y_SHIFT),   "anchor": (325, 260 - Y_SHIFT), "align": "left"},
    }
    
    # --- Helpers ---
    def _get_img_bytes(item):
        c = nz_str(item.get("Composition link"))
        l = nz_str(item.get("Link"))
        img = None
        if c: img = get_image_if_cached("", c) or get_image_cached("", c)
        if not img and l: img = get_image_if_cached(l, "") or get_image_cached(l, "")
        return img

    def _draw_elbow_line(canvas, start_x, start_y, box_x, box_y, box_w, box_h, align):
        canvas.saveState()
        canvas.setStrokeColor(black)
        canvas.setLineWidth(1)
        canvas.setDash([3, 2])
        box_mid_y = box_y + (box_h / 2)
        p = canvas.beginPath()
        p.moveTo(start_x, start_y)
        if align == "right":
            end_x = box_x
            mid_x = start_x + 30 if abs(start_y - box_mid_y) > 50 else start_x + 15
            p.lineTo(mid_x, start_y)
            p.lineTo(mid_x, box_mid_y)
            p.lineTo(end_x, box_mid_y)
        else: 
            end_x = box_x + box_w
            mid_x = start_x - 30 if abs(start_y - box_mid_y) > 50 else start_x - 15
            p.lineTo(mid_x, start_y)
            p.lineTo(mid_x, box_mid_y)
            p.lineTo(end_x, box_mid_y)
        canvas.drawPath(p, stroke=1, fill=0)
        canvas.setDash([]) 
        canvas.setFillColor(black)
        canvas.circle(start_x, start_y, 2.5, fill=1, stroke=0)
        canvas.restoreState()

    # --- Header & Border Renderer ---
    def draw_border_and_header(canvas, doc, is_page_1=False):
        canvas.saveState()
        canvas.setStrokeColor(HexColor("#000000"))
        canvas.setLineWidth(1.5)
        canvas.rect(10, 10, WIDTH - 20, HEIGHT - 20)

        if is_page_1:
            date_str = datetime.now().strftime("%d-%m-%Y")
            c_name = client_name.strip() or ""
            app_by = approved_by.strip() or "Altossa"
            s_style = ParagraphStyle('Hdr', fontName='Helvetica', fontSize=10, leading=12)
            row_data = [
                Paragraph(f"<b>Client :</b> {c_name}", s_style),
                Paragraph(f"<b>Approved by :</b> {app_by}", s_style),
                Paragraph(f"<b>Room :</b> Living room", s_style),
                Paragraph(f"<b>Date :</b> {date_str}", s_style)
            ]
            t_width = WIDTH - 40
            col_w = t_width / 4.0
            t = Table([row_data], colWidths=[col_w]*4, rowHeights=[35])
            t.setStyle(TableStyle([
                ('GRID', (0,0), (-1,-1), 1.5, colors.black),
                ('VALIGN', (0,0), (-1,-1), 'MIDDLE'),
                ('LEFTPADDING', (0,0), (-1,-1), 8),
                ('RIGHTPADDING', (0,0), (-1,-1), 8),
            ]))
            w, h = t.wrap(t_width, 100)
            t.drawOn(canvas, 20, HEIGHT - 10 - h - 10)
        canvas.restoreState()

    # --- Filter Content ---
    fixed_slots = {k: [] for k in SLOT_CONFIG.keys()}
    others_items = []
    
    for item in cart:
        # Check Living Room allocation
        raw_alloc = item.get("room_alloc", {})
        lr_data = raw_alloc.get("Living Room")
        
        if not lr_data:
            continue

        lr_qty = 0
        visual_role = "Others"

        if isinstance(lr_data, dict):
            lr_qty = int(lr_data.get("qty", 0))
            visual_role = lr_data.get("visual_role") or "Others"
        else:
            try: lr_qty = int(lr_data)
            except: lr_qty = 0
        
        if lr_qty <= 0:
            continue

        pdf_item = item.copy()
        pdf_item["qty"] = lr_qty

        if visual_role in fixed_slots:
            if not fixed_slots[visual_role]:
                fixed_slots[visual_role].append(pdf_item)
            else:
                others_items.append(pdf_item)
        else:
            others_items.append(pdf_item)

    if fixed_slots.get("Others"):
        others_items.extend(fixed_slots["Others"])
        fixed_slots["Others"] = [] 

    # --- Setup ReportLab ---
    buffer = BytesIO()
    doc = SimpleDocTemplate(
        buffer, pagesize=page_size,
        leftMargin=0, rightMargin=0, topMargin=0, bottomMargin=0
    )
    styles = getSampleStyleSheet()
    tbl_txt_style = ParagraphStyle('Compact', parent=styles['Normal'], fontName='Helvetica', fontSize=6.5, leading=7)
    tbl_lbl_style = ParagraphStyle('CompactBold', parent=styles['Normal'], fontName='Helvetica-Bold', fontSize=6.5, leading=7)
    tbl_style = TableStyle([
        ('VALIGN', (0,0), (-1,-1), 'TOP'),
        ('LEFTPADDING', (0,0), (-1,-1), 0),
        ('RIGHTPADDING', (0,0), (-1,-1), 0),
        ('BOTTOMPADDING', (0,0), (-1,-1), 0),
        ('TOPPADDING', (0,0), (-1,-1), 0),
    ])

    # --- Visualizer Renderer ---
    def draw_visualizer_bg(canvas, doc):
        canvas.saveState()
        if os.path.exists(LIVING_TEMPLATE_PATH):
            try:
                canvas.drawImage(LIVING_TEMPLATE_PATH, 0, -Y_SHIFT, width=WIDTH, height=HEIGHT, preserveAspectRatio=True, anchor='c')
            except: pass
        canvas.restoreState()

        draw_border_and_header(canvas, doc, is_page_1=True)
        
        canvas.saveState()
        for role, items in fixed_slots.items():
            if items:
                item = items[0]
                config = SLOT_CONFIG[role]
                box_x, box_y = config["box"]
                _draw_elbow_line(canvas, config["anchor"][0], config["anchor"][1], box_x, box_y, BOX_W, BOX_H, config["align"])

                canvas.setStrokeColor(HexColor("#333333"))
                canvas.setFillColor(colors.white)
                canvas.setLineWidth(1)
                canvas.roundRect(box_x, box_y, BOX_W, BOX_H, 8, fill=1, stroke=1)

                img_bytes = _get_img_bytes(item)
                img_margin = 8
                if img_bytes:
                    try:
                        rdr = ImageReader(BytesIO(img_bytes))
                        iw, ih = rdr.getSize()
                        scale = min(IMG_W / float(iw), IMG_H / float(ih))
                        w, h = iw * scale, ih * scale
                        img_y_pos = box_y + (BOX_H - h) / 2
                        canvas.drawImage(rdr, box_x + img_margin, img_y_pos, width=w, height=h, mask='auto')
                    except: pass

                table_data = _get_item_pdf_data(item, role)
                processed_data = []
                for row in table_data:
                    label, sep, val = row
                    processed_data.append([
                        Paragraph(f"{label}", tbl_lbl_style),
                        Paragraph(f"{sep}", tbl_txt_style),
                        Paragraph(f"{val}", tbl_txt_style)
                    ])

                text_x_start = box_x + IMG_W + (img_margin * 2)
                text_w_avail = BOX_W - IMG_W - (img_margin * 3)
                t = Table(processed_data, colWidths=[45, 5, text_w_avail - 55])
                t.setStyle(tbl_style)
                
                f = Frame(text_x_start, box_y, text_w_avail, BOX_H, leftPadding=0, bottomPadding=5, rightPadding=2, topPadding=5, showBoundary=0)
                f.addFromList([t], canvas)
        canvas.restoreState()

    # --- List Renderer ---
    def draw_list_header(canvas, doc):
        draw_border_and_header(canvas, doc, is_page_1=False)
        canvas.saveState()
        canvas.setFont("Helvetica-Bold", 12)
        canvas.drawString(1.5*cm, HEIGHT - 2.5*cm, "Living Room - Additional Items")
        canvas.line(1.5*cm, HEIGHT - 2.5*cm - 5, WIDTH - 1.5*cm, HEIGHT - 2.5*cm - 5)
        canvas.restoreState()

    frame_p1 = Frame(0, 0, WIDTH, HEIGHT, id='f1', showBoundary=0)
    tpl_p1 = PageTemplate(id='Visualizer', frames=[frame_p1], onPage=draw_visualizer_bg)

    frame_p2 = Frame(1.5*cm, 1.5*cm, WIDTH-3*cm, HEIGHT-4.0*cm, id='f2', showBoundary=0)
    tpl_p2 = PageTemplate(id='List', frames=[frame_p2], onPage=draw_list_header)

    doc.addPageTemplates([tpl_p1, tpl_p2])

    story = []
    story.append(Spacer(1, 1))

    if others_items:
        story.append(NextPageTemplate('List'))
        story.append(PageBreak())
        
        grid_data = []
        current_row = []
        col_w = (WIDTH - 3*cm - 20) / 2
        
        for item in others_items:
            raw_data = _get_item_pdf_data(item, "Others")
            tbl_rows = []
            for row in raw_data:
                label, sep, val = row
                p_val = Paragraph(val, styles["Normal"])
                p_lbl = Paragraph(f"<b>{label}</b>", styles["Normal"])
                tbl_rows.append([p_lbl, sep, p_val])

            inner_tbl = Table(tbl_rows, colWidths=[65, 5, col_w - 90 - 70]) 
            inner_tbl.setStyle(TableStyle([
                ('VALIGN', (0,0), (-1,-1), 'TOP'),
                ('LEFTPADDING', (0,0), (-1,-1), 0),
                ('RIGHTPADDING', (0,0), (-1,-1), 0),
                ('TOPPADDING', (0,0), (-1,-1), 0),
                ('BOTTOMPADDING', (0,0), (-1,-1), 0),
            ]))

            img_bytes = _get_img_bytes(item)
            img_obj = Spacer(80, 60) 
            if img_bytes:
                try:
                    rdr = ImageReader(BytesIO(img_bytes))
                    iw, ih = rdr.getSize()
                    scale = min(80.0/iw, 60.0/ih)
                    img_obj = RLImage(BytesIO(img_bytes), width=iw*scale, height=ih*scale)
                except: pass
            
            item_data = [[img_obj, inner_tbl]]
            item_tbl = Table(item_data, colWidths=[90, col_w - 90])
            item_tbl.setStyle(TableStyle([
                ('BOX', (0,0), (-1,-1), 0.5, HexColor("#cccccc")),
                ('VALIGN', (0,0), (-1,-1), 'MIDDLE'),
                ('LEFTPADDING', (0,0), (-1,-1), 4),
                ('RIGHTPADDING', (0,0), (-1,-1), 4),
                ('TOPPADDING', (0,0), (-1,-1), 4),
                ('BOTTOMPADDING', (0,0), (-1,-1), 4),
            ]))
            
            current_row.append(item_tbl)
            
            if len(current_row) == 2:
                grid_data.append(current_row)
                current_row = []
        
        if current_row:
            current_row.append(Spacer(1,1))
            grid_data.append(current_row)
            
        main_tbl = Table(grid_data, colWidths=[col_w, col_w])
        main_tbl.setStyle(TableStyle([
            ('VALIGN', (0,0), (-1,-1), 'TOP'),
            ('LEFTPADDING', (0,0), (-1,-1), 0),
            ('RIGHTPADDING', (0,0), (-1, -1), 10),
            ('BOTTOMPADDING', (0,0), (-1,-1), 5),
        ]))
        
        story.append(main_tbl)

    try:
        doc.build(story)
        buffer.seek(0)
        return buffer.read()
    except Exception as e:
        st.error(f"PDF Generation Error: {e}")
        return None

def export_bedroom_pdf(cart: list[dict], room_key: str, client_name: str = "", approved_by: str = "") -> Optional[bytes]:
    """
    Generates a Bedroom visual layout PDF for a specific bedroom (e.g. 'Bedroom 1').
    """
    from datetime import datetime

    # --- Configuration ---
    page_size = landscape(A4)
    WIDTH, HEIGHT = page_size 
    Y_SHIFT = 20  
    BOX_W, BOX_H = 225, 125  
    IMG_W, IMG_H = 95, 95    

    # Bedroom Slots
    SLOT_CONFIG = {
        "Bed":              {"box": (575, 390 - Y_SHIFT), "anchor": (515, 330 - Y_SHIFT), "align": "right"},
        "Wardrobe":         {"box": (575, 240 - Y_SHIFT), "anchor": (515, 295 - Y_SHIFT), "align": "right"},
        "Night Tables":     {"box": (40, 390 - Y_SHIFT),  "anchor": (325, 330 - Y_SHIFT), "align": "left"},
        "Chest of Drawers": {"box": (40, 240 - Y_SHIFT),  "anchor": (325, 295 - Y_SHIFT), "align": "left"},
    }
    
    # --- Helpers ---
    def _get_img_bytes(item):
        c = nz_str(item.get("Composition link"))
        l = nz_str(item.get("Link"))
        img = None
        if c: img = get_image_if_cached("", c) or get_image_cached("", c)
        if not img and l: img = get_image_if_cached(l, "") or get_image_cached(l, "")
        return img

    def _draw_elbow_line(canvas, start_x, start_y, box_x, box_y, box_w, box_h, align):
        canvas.saveState()
        canvas.setStrokeColor(black)
        canvas.setLineWidth(1)
        canvas.setDash([3, 2])
        box_mid_y = box_y + (box_h / 2)
        p = canvas.beginPath()
        p.moveTo(start_x, start_y)
        if align == "right":
            end_x = box_x
            mid_x = start_x + 30 if abs(start_y - box_mid_y) > 50 else start_x + 15
            p.lineTo(mid_x, start_y); p.lineTo(mid_x, box_mid_y); p.lineTo(end_x, box_mid_y)
        else: 
            end_x = box_x + box_w
            mid_x = start_x - 30 if abs(start_y - box_mid_y) > 50 else start_x - 15
            p.lineTo(mid_x, start_y); p.lineTo(mid_x, box_mid_y); p.lineTo(end_x, box_mid_y)
        canvas.drawPath(p, stroke=1, fill=0)
        canvas.setDash([]) 
        canvas.setFillColor(black)
        canvas.circle(start_x, start_y, 2.5, fill=1, stroke=0)
        canvas.restoreState()

    def draw_border_and_header(canvas, doc, is_page_1=False):
        canvas.saveState()
        canvas.setStrokeColor(HexColor("#000000"))
        canvas.setLineWidth(1.5)
        canvas.rect(10, 10, WIDTH - 20, HEIGHT - 20)
        if is_page_1:
            date_str = datetime.now().strftime("%d-%m-%Y")
            c_name = client_name.strip() or ""
            app_by = approved_by.strip() or "Altossa"
            s_style = ParagraphStyle('Hdr', fontName='Helvetica', fontSize=10, leading=12)
            row_data = [
                Paragraph(f"<b>Client :</b> {c_name}", s_style),
                Paragraph(f"<b>Approved by :</b> {app_by}", s_style),
                Paragraph(f"<b>Room :</b> {room_key}", s_style), # [UPDATED]
                Paragraph(f"<b>Date :</b> {date_str}", s_style)
            ]
            t_width = WIDTH - 40
            col_w = t_width / 4.0
            t = Table([row_data], colWidths=[col_w]*4, rowHeights=[35])
            t.setStyle(TableStyle([
                ('GRID', (0,0), (-1,-1), 1.5, colors.black),
                ('VALIGN', (0,0), (-1,-1), 'MIDDLE'),
                ('LEFTPADDING', (0,0), (-1,-1), 8),
                ('RIGHTPADDING', (0,0), (-1,-1), 8),
            ]))
            w, h = t.wrap(t_width, 100)
            t.drawOn(canvas, 20, HEIGHT - 10 - h - 10)
        canvas.restoreState()

    fixed_slots = {k: [] for k in SLOT_CONFIG.keys()}
    others_items = []
    
    for item in cart:
        raw_alloc = item.get("room_alloc", {})
        br_data = raw_alloc.get(room_key) # [UPDATED]
        if not br_data: continue

        br_qty = 0
        visual_role = "Others"
        if isinstance(br_data, dict):
            br_qty = int(br_data.get("qty", 0))
            visual_role = br_data.get("visual_role") or "Others"
        else:
            try: br_qty = int(br_data)
            except: br_qty = 0
        
        if br_qty <= 0: continue

        pdf_item = item.copy()
        pdf_item["qty"] = br_qty

        if visual_role in fixed_slots:
            if not fixed_slots[visual_role]:
                fixed_slots[visual_role].append(pdf_item)
            else:
                others_items.append(pdf_item)
        else:
            others_items.append(pdf_item)

    if fixed_slots.get("Others"):
        others_items.extend(fixed_slots["Others"])
        fixed_slots["Others"] = [] 

    buffer = BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=page_size, leftMargin=0, rightMargin=0, topMargin=0, bottomMargin=0)
    styles = getSampleStyleSheet()
    
    tbl_txt_style = ParagraphStyle('Compact', parent=styles['Normal'], fontName='Helvetica', fontSize=6.5, leading=7)
    tbl_lbl_style = ParagraphStyle('CompactBold', parent=styles['Normal'], fontName='Helvetica-Bold', fontSize=6.5, leading=7)
    tbl_style = TableStyle([
        ('VALIGN', (0,0), (-1,-1), 'TOP'),
        ('LEFTPADDING', (0,0), (-1,-1), 0),
        ('RIGHTPADDING', (0,0), (-1,-1), 0),
        ('BOTTOMPADDING', (0,0), (-1,-1), 0),
        ('TOPPADDING', (0,0), (-1,-1), 0),
    ])

    def draw_visualizer_bg(canvas, doc):
        canvas.saveState()
        if os.path.exists(BEDROOM_TEMPLATE_PATH):
            try:
                canvas.drawImage(BEDROOM_TEMPLATE_PATH, 0, -Y_SHIFT, width=WIDTH, height=HEIGHT, preserveAspectRatio=True, anchor='c')
            except: pass
        canvas.restoreState()
        draw_border_and_header(canvas, doc, is_page_1=True)
        canvas.saveState()
        for role, items in fixed_slots.items():
            if items:
                item = items[0]
                config = SLOT_CONFIG[role]
                box_x, box_y = config["box"]
                _draw_elbow_line(canvas, config["anchor"][0], config["anchor"][1], box_x, box_y, BOX_W, BOX_H, config["align"])
                canvas.setStrokeColor(HexColor("#333333"))
                canvas.setFillColor(colors.white)
                canvas.setLineWidth(1)
                canvas.roundRect(box_x, box_y, BOX_W, BOX_H, 8, fill=1, stroke=1)
                img_bytes = _get_img_bytes(item)
                img_margin = 8
                if img_bytes:
                    try:
                        rdr = ImageReader(BytesIO(img_bytes))
                        iw, ih = rdr.getSize()
                        scale = min(IMG_W / float(iw), IMG_H / float(ih))
                        w, h = iw * scale, ih * scale
                        img_y_pos = box_y + (BOX_H - h) / 2
                        canvas.drawImage(rdr, box_x + img_margin, img_y_pos, width=w, height=h, mask='auto')
                    except: pass
                table_data = _get_item_pdf_data(item, role)
                processed_data = []
                for row in table_data:
                    label, sep, val = row
                    processed_data.append([
                        Paragraph(f"{label}", tbl_lbl_style),
                        Paragraph(f"{sep}", tbl_txt_style),
                        Paragraph(f"{val}", tbl_txt_style)
                    ])
                text_x_start = box_x + IMG_W + (img_margin * 2)
                text_w_avail = BOX_W - IMG_W - (img_margin * 3)
                t = Table(processed_data, colWidths=[45, 5, text_w_avail - 55])
                t.setStyle(tbl_style)
                f = Frame(text_x_start, box_y, text_w_avail, BOX_H, leftPadding=0, bottomPadding=5, rightPadding=2, topPadding=5, showBoundary=0)
                f.addFromList([t], canvas)
        canvas.restoreState()

    def draw_list_header(canvas, doc):
        draw_border_and_header(canvas, doc, is_page_1=False)
        canvas.saveState()
        canvas.setFont("Helvetica-Bold", 12)
        canvas.drawString(1.5*cm, HEIGHT - 2.5*cm, f"{room_key} - Additional Items") # [UPDATED]
        canvas.line(1.5*cm, HEIGHT - 2.5*cm - 5, WIDTH - 1.5*cm, HEIGHT - 2.5*cm - 5)
        canvas.restoreState()

    frame_p1 = Frame(0, 0, WIDTH, HEIGHT, id='f1', showBoundary=0)
    tpl_p1 = PageTemplate(id='Visualizer', frames=[frame_p1], onPage=draw_visualizer_bg)
    frame_p2 = Frame(1.5*cm, 1.5*cm, WIDTH-3*cm, HEIGHT-4.0*cm, id='f2', showBoundary=0)
    tpl_p2 = PageTemplate(id='List', frames=[frame_p2], onPage=draw_list_header)
    doc.addPageTemplates([tpl_p1, tpl_p2])

    story = []
    story.append(Spacer(1, 1))
    if others_items:
        story.append(NextPageTemplate('List'))
        story.append(PageBreak())
        grid_data = []
        current_row = []
        col_w = (WIDTH - 3*cm - 20) / 2
        for item in others_items:
            raw_data = _get_item_pdf_data(item, "Others")
            tbl_rows = []
            for row in raw_data:
                label, sep, val = row
                p_val = Paragraph(val, styles["Normal"])
                p_lbl = Paragraph(f"<b>{label}</b>", styles["Normal"])
                tbl_rows.append([p_lbl, sep, p_val])
            inner_tbl = Table(tbl_rows, colWidths=[65, 5, col_w - 90 - 70]) 
            inner_tbl.setStyle(TableStyle([('VALIGN', (0,0), (-1,-1), 'TOP'),('LEFTPADDING', (0,0), (-1,-1), 0),('RIGHTPADDING', (0,0), (-1,-1), 0),('TOPPADDING', (0,0), (-1,-1), 0),('BOTTOMPADDING', (0,0), (-1,-1), 0)]))
            img_bytes = _get_img_bytes(item)
            img_obj = Spacer(80, 60) 
            if img_bytes:
                try:
                    rdr = ImageReader(BytesIO(img_bytes))
                    iw, ih = rdr.getSize()
                    scale = min(80.0/iw, 60.0/ih)
                    img_obj = RLImage(BytesIO(img_bytes), width=iw*scale, height=ih*scale)
                except: pass
            item_data = [[img_obj, inner_tbl]]
            item_tbl = Table(item_data, colWidths=[90, col_w - 90])
            item_tbl.setStyle(TableStyle([('BOX', (0,0), (-1,-1), 0.5, HexColor("#cccccc")),('VALIGN', (0,0), (-1,-1), 'MIDDLE'),('LEFTPADDING', (0,0), (-1,-1), 4),('RIGHTPADDING', (0,0), (-1,-1), 4),('TOPPADDING', (0,0), (-1,-1), 4),('BOTTOMPADDING', (0,0), (-1,-1), 4)]))
            current_row.append(item_tbl)
            if len(current_row) == 2:
                grid_data.append(current_row); current_row = []
        if current_row:
            current_row.append(Spacer(1,1))
            grid_data.append(current_row)
        main_tbl = Table(grid_data, colWidths=[col_w, col_w])
        main_tbl.setStyle(TableStyle([('VALIGN', (0,0), (-1,-1), 'TOP'),('LEFTPADDING', (0,0), (-1,-1), 0),('RIGHTPADDING', (0,0), (-1, -1), 10),('BOTTOMPADDING', (0,0), (-1,-1), 5)]))
        story.append(main_tbl)

    try:
        doc.build(story)
        buffer.seek(0)
        return buffer.read()
    except Exception as e:
        st.error(f"PDF Generation Error: {e}")
        return None

def export_dining_room_pdf(cart: list[dict], client_name: str = "", approved_by: str = "") -> Optional[bytes]:
    """
    Generates a Dining Room visual layout PDF.
    """
    from datetime import datetime

    page_size = landscape(A4)
    WIDTH, HEIGHT = page_size 
    Y_SHIFT = 20  
    BOX_W, BOX_H = 225, 125  
    IMG_W, IMG_H = 95, 95    

    # Dining Room Slots
    SLOT_CONFIG = {
        "Dining Table": {"box": (575, 390 - Y_SHIFT), "anchor": (515, 330 - Y_SHIFT), "align": "right"},
        "Chairs":       {"box": (40, 390 - Y_SHIFT),  "anchor": (325, 330 - Y_SHIFT), "align": "left"},
    }
    
    # --- Helpers ---
    def _get_img_bytes(item):
        c = nz_str(item.get("Composition link"))
        l = nz_str(item.get("Link"))
        img = None
        if c: img = get_image_if_cached("", c) or get_image_cached("", c)
        if not img and l: img = get_image_if_cached(l, "") or get_image_cached(l, "")
        return img

    def _draw_elbow_line(canvas, start_x, start_y, box_x, box_y, box_w, box_h, align):
        canvas.saveState()
        canvas.setStrokeColor(black)
        canvas.setLineWidth(1)
        canvas.setDash([3, 2])
        box_mid_y = box_y + (box_h / 2)
        p = canvas.beginPath()
        p.moveTo(start_x, start_y)
        if align == "right":
            end_x = box_x
            mid_x = start_x + 30 if abs(start_y - box_mid_y) > 50 else start_x + 15
            p.lineTo(mid_x, start_y); p.lineTo(mid_x, box_mid_y); p.lineTo(end_x, box_mid_y)
        else: 
            end_x = box_x + box_w
            mid_x = start_x - 30 if abs(start_y - box_mid_y) > 50 else start_x - 15
            p.lineTo(mid_x, start_y); p.lineTo(mid_x, box_mid_y); p.lineTo(end_x, box_mid_y)
        canvas.drawPath(p, stroke=1, fill=0)
        canvas.setDash([]) 
        canvas.setFillColor(black)
        canvas.circle(start_x, start_y, 2.5, fill=1, stroke=0)
        canvas.restoreState()

    def draw_border_and_header(canvas, doc, is_page_1=False):
        canvas.saveState()
        canvas.setStrokeColor(HexColor("#000000"))
        canvas.setLineWidth(1.5)
        canvas.rect(10, 10, WIDTH - 20, HEIGHT - 20)
        if is_page_1:
            date_str = datetime.now().strftime("%d-%m-%Y")
            c_name = client_name.strip() or ""
            app_by = approved_by.strip() or "Altossa"
            s_style = ParagraphStyle('Hdr', fontName='Helvetica', fontSize=10, leading=12)
            row_data = [
                Paragraph(f"<b>Client :</b> {c_name}", s_style),
                Paragraph(f"<b>Approved by :</b> {app_by}", s_style),
                Paragraph(f"<b>Room :</b> Dining room", s_style),
                Paragraph(f"<b>Date :</b> {date_str}", s_style)
            ]
            t_width = WIDTH - 40
            col_w = t_width / 4.0
            t = Table([row_data], colWidths=[col_w]*4, rowHeights=[35])
            t.setStyle(TableStyle([
                ('GRID', (0,0), (-1,-1), 1.5, colors.black),
                ('VALIGN', (0,0), (-1,-1), 'MIDDLE'),
                ('LEFTPADDING', (0,0), (-1,-1), 8),
                ('RIGHTPADDING', (0,0), (-1,-1), 8),
            ]))
            w, h = t.wrap(t_width, 100)
            t.drawOn(canvas, 20, HEIGHT - 10 - h - 10)
        canvas.restoreState()

    fixed_slots = {k: [] for k in SLOT_CONFIG.keys()}
    others_items = []
    
    for item in cart:
        raw_alloc = item.get("room_alloc", {})
        din_data = raw_alloc.get("Dining Room")
        if not din_data: continue

        din_qty = 0
        visual_role = "Others"
        if isinstance(din_data, dict):
            din_qty = int(din_data.get("qty", 0))
            visual_role = din_data.get("visual_role") or "Others"
        else:
            try: din_qty = int(din_data)
            except: din_qty = 0
        
        if din_qty <= 0: continue

        pdf_item = item.copy()
        pdf_item["qty"] = din_qty

        if visual_role in fixed_slots:
            if not fixed_slots[visual_role]:
                fixed_slots[visual_role].append(pdf_item)
            else:
                others_items.append(pdf_item)
        else:
            others_items.append(pdf_item)

    if fixed_slots.get("Others"):
        others_items.extend(fixed_slots["Others"])
        fixed_slots["Others"] = [] 

    buffer = BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=page_size, leftMargin=0, rightMargin=0, topMargin=0, bottomMargin=0)
    styles = getSampleStyleSheet()
    
    tbl_txt_style = ParagraphStyle('Compact', parent=styles['Normal'], fontName='Helvetica', fontSize=6.5, leading=7)
    tbl_lbl_style = ParagraphStyle('CompactBold', parent=styles['Normal'], fontName='Helvetica-Bold', fontSize=6.5, leading=7)
    tbl_style = TableStyle([
        ('VALIGN', (0,0), (-1,-1), 'TOP'),
        ('LEFTPADDING', (0,0), (-1,-1), 0),
        ('RIGHTPADDING', (0,0), (-1,-1), 0),
        ('BOTTOMPADDING', (0,0), (-1,-1), 0),
        ('TOPPADDING', (0,0), (-1,-1), 0),
    ])

    def draw_visualizer_bg(canvas, doc):
        canvas.saveState()
        if os.path.exists(DINING_TEMPLATE_PATH):
            try:
                canvas.drawImage(DINING_TEMPLATE_PATH, 0, -Y_SHIFT, width=WIDTH, height=HEIGHT, preserveAspectRatio=True, anchor='c')
            except: pass
        canvas.restoreState()
        draw_border_and_header(canvas, doc, is_page_1=True)
        canvas.saveState()
        for role, items in fixed_slots.items():
            if items:
                item = items[0]
                config = SLOT_CONFIG[role]
                box_x, box_y = config["box"]
                _draw_elbow_line(canvas, config["anchor"][0], config["anchor"][1], box_x, box_y, BOX_W, BOX_H, config["align"])
                canvas.setStrokeColor(HexColor("#333333"))
                canvas.setFillColor(colors.white)
                canvas.setLineWidth(1)
                canvas.roundRect(box_x, box_y, BOX_W, BOX_H, 8, fill=1, stroke=1)
                img_bytes = _get_img_bytes(item)
                img_margin = 8
                if img_bytes:
                    try:
                        rdr = ImageReader(BytesIO(img_bytes))
                        iw, ih = rdr.getSize()
                        scale = min(IMG_W / float(iw), IMG_H / float(ih))
                        w, h = iw * scale, ih * scale
                        img_y_pos = box_y + (BOX_H - h) / 2
                        canvas.drawImage(rdr, box_x + img_margin, img_y_pos, width=w, height=h, mask='auto')
                    except: pass
                table_data = _get_item_pdf_data(item, role)
                processed_data = []
                for row in table_data:
                    label, sep, val = row
                    processed_data.append([
                        Paragraph(f"{label}", tbl_lbl_style),
                        Paragraph(f"{sep}", tbl_txt_style),
                        Paragraph(f"{val}", tbl_txt_style)
                    ])
                text_x_start = box_x + IMG_W + (img_margin * 2)
                text_w_avail = BOX_W - IMG_W - (img_margin * 3)
                t = Table(processed_data, colWidths=[45, 5, text_w_avail - 55])
                t.setStyle(tbl_style)
                f = Frame(text_x_start, box_y, text_w_avail, BOX_H, leftPadding=0, bottomPadding=5, rightPadding=2, topPadding=5, showBoundary=0)
                f.addFromList([t], canvas)
        canvas.restoreState()

    def draw_list_header(canvas, doc):
        draw_border_and_header(canvas, doc, is_page_1=False)
        canvas.saveState()
        canvas.setFont("Helvetica-Bold", 12)
        canvas.drawString(1.5*cm, HEIGHT - 2.5*cm, "Dining Room - Additional Items")
        canvas.line(1.5*cm, HEIGHT - 2.5*cm - 5, WIDTH - 1.5*cm, HEIGHT - 2.5*cm - 5)
        canvas.restoreState()

    frame_p1 = Frame(0, 0, WIDTH, HEIGHT, id='f1', showBoundary=0)
    tpl_p1 = PageTemplate(id='Visualizer', frames=[frame_p1], onPage=draw_visualizer_bg)
    frame_p2 = Frame(1.5*cm, 1.5*cm, WIDTH-3*cm, HEIGHT-4.0*cm, id='f2', showBoundary=0)
    tpl_p2 = PageTemplate(id='List', frames=[frame_p2], onPage=draw_list_header)
    doc.addPageTemplates([tpl_p1, tpl_p2])

    story = []
    story.append(Spacer(1, 1))
    if others_items:
        story.append(NextPageTemplate('List'))
        story.append(PageBreak())
        grid_data = []
        current_row = []
        col_w = (WIDTH - 3*cm - 20) / 2
        for item in others_items:
            raw_data = _get_item_pdf_data(item, "Others")
            tbl_rows = []
            for row in raw_data:
                label, sep, val = row
                p_val = Paragraph(val, styles["Normal"])
                p_lbl = Paragraph(f"<b>{label}</b>", styles["Normal"])
                tbl_rows.append([p_lbl, sep, p_val])
            inner_tbl = Table(tbl_rows, colWidths=[65, 5, col_w - 90 - 70]) 
            inner_tbl.setStyle(TableStyle([('VALIGN', (0,0), (-1,-1), 'TOP'),('LEFTPADDING', (0,0), (-1,-1), 0),('RIGHTPADDING', (0,0), (-1,-1), 0),('TOPPADDING', (0,0), (-1,-1), 0),('BOTTOMPADDING', (0,0), (-1,-1), 0)]))
            img_bytes = _get_img_bytes(item)
            img_obj = Spacer(80, 60) 
            if img_bytes:
                try:
                    rdr = ImageReader(BytesIO(img_bytes))
                    iw, ih = rdr.getSize()
                    scale = min(80.0/iw, 60.0/ih)
                    img_obj = RLImage(BytesIO(img_bytes), width=iw*scale, height=ih*scale)
                except: pass
            item_data = [[img_obj, inner_tbl]]
            item_tbl = Table(item_data, colWidths=[90, col_w - 90])
            item_tbl.setStyle(TableStyle([('BOX', (0,0), (-1,-1), 0.5, HexColor("#cccccc")),('VALIGN', (0,0), (-1,-1), 'MIDDLE'),('LEFTPADDING', (0,0), (-1,-1), 4),('RIGHTPADDING', (0,0), (-1,-1), 4),('TOPPADDING', (0,0), (-1,-1), 4),('BOTTOMPADDING', (0,0), (-1,-1), 4)]))
            current_row.append(item_tbl)
            if len(current_row) == 2:
                grid_data.append(current_row); current_row = []
        if current_row:
            current_row.append(Spacer(1,1))
            grid_data.append(current_row)
        main_tbl = Table(grid_data, colWidths=[col_w, col_w])
        main_tbl.setStyle(TableStyle([('VALIGN', (0,0), (-1,-1), 'TOP'),('LEFTPADDING', (0,0), (-1,-1), 0),('RIGHTPADDING', (0,0), (-1, -1), 10),('BOTTOMPADDING', (0,0), (-1,-1), 5)]))
        story.append(main_tbl)

    try:
        doc.build(story)
        buffer.seek(0)
        return buffer.read()
    except Exception as e:
        st.error(f"PDF Generation Error: {e}")
        return None

def export_others_room_pdf(cart: list[dict], client_name: str = "", approved_by: str = "") -> Optional[bytes]:
    """
    Generates a visual layout PDF for items assigned to 'Others' room or any unassigned overflow.
    - Layout: Grid 4x2 (8 items/page) - Similar to 'Additional Items' pages.
    - No visualizer diagram page (starts directly as a list).
    """
    from datetime import datetime

    # --- Configuration ---
    page_size = landscape(A4)
    WIDTH, HEIGHT = page_size 
    
    # --- Helpers ---
    def _get_img_bytes(item):
        c = nz_str(item.get("Composition link"))
        l = nz_str(item.get("Link"))
        img = None
        if c: img = get_image_if_cached("", c) or get_image_cached("", c)
        if not img and l: img = get_image_if_cached(l, "") or get_image_cached(l, "")
        return img

    def draw_border_and_header(canvas, doc):
        canvas.saveState()
        # Border
        canvas.setStrokeColor(HexColor("#000000"))
        canvas.setLineWidth(1.5)
        canvas.rect(10, 10, WIDTH - 20, HEIGHT - 20)
        
        # Header Box (Simple Title Block)
        date_str = datetime.now().strftime("%d-%m-%Y")
        c_name = client_name.strip() or ""
        app_by = approved_by.strip() or "Altossa"
        
        s_style = ParagraphStyle('Hdr', fontName='Helvetica', fontSize=10, leading=12)
        
        row_data = [
            Paragraph(f"<b>Client :</b> {c_name}", s_style),
            Paragraph(f"<b>Approved by :</b> {app_by}", s_style),
            Paragraph(f"<b>Room :</b> Others / General", s_style),
            Paragraph(f"<b>Date :</b> {date_str}", s_style)
        ]
        
        t_width = WIDTH - 40
        col_w = t_width / 4.0
        
        t = Table([row_data], colWidths=[col_w]*4, rowHeights=[35])
        t.setStyle(TableStyle([
            ('GRID', (0,0), (-1,-1), 1.5, colors.black),
            ('VALIGN', (0,0), (-1,-1), 'MIDDLE'),
            ('LEFTPADDING', (0,0), (-1,-1), 8),
            ('RIGHTPADDING', (0,0), (-1,-1), 8),
        ]))
        
        # Draw header at top inside border
        w, h = t.wrap(t_width, 100)
        t.drawOn(canvas, 20, HEIGHT - 10 - h - 10)
        
        # Title Line
        canvas.setFont("Helvetica-Bold", 12)
        canvas.drawString(1.5*cm, HEIGHT - 2.5*cm, "Other Items / General Selection")
        canvas.line(1.5*cm, HEIGHT - 2.5*cm - 5, WIDTH - 1.5*cm, HEIGHT - 2.5*cm - 5)
        
        canvas.restoreState()

    # --- Filter Content ---
    others_items = []
    
    for item in cart:
        qty_total = int(item.get("qty", 1))
        raw_alloc = item.get("room_alloc", {})
        
        # 1. Calculate total explicitly assigned to ANY known room (including "Others")
        assigned_qty = 0
        explicit_others_qty = 0
        
        if isinstance(raw_alloc, dict):
            for r, v in raw_alloc.items():
                q = _room_qty(v)
                if r == "Others":
                    explicit_others_qty += q
                else:
                    assigned_qty += q
        
        # 2. Calculate unassigned remainder
        # (Total Item Qty) - (Qty assigned to specific rooms like Living/Drawing/etc) - (Qty explicitly assigned to Others)
        # Actually, "Others" IS an assignment. So remainder is what's left after ALL assignments.
        # However, the user wants "Others" PDF to contain:
        #   A) Explicit "Others" assignment
        #   B) Unassigned remainder (implicit Others)
        
        # Total assigned to non-Others rooms
        assigned_to_specific_rooms = assigned_qty 
        
        # Remainder that defaults to "Others"
        remainder = qty_total - assigned_to_specific_rooms - explicit_others_qty
        
        # Safety: remainder shouldn't be negative if data is clean, but clamp it
        if remainder < 0: remainder = 0
        
        final_oth_qty = explicit_others_qty + remainder
            
        if final_oth_qty <= 0:
            continue

        # Create PDF item
        pdf_item = item.copy()
        pdf_item["qty"] = final_oth_qty
        others_items.append(pdf_item)

    if not others_items:
        return None # Nothing to export

    # --- Setup ReportLab ---
    buffer = BytesIO()
    doc = SimpleDocTemplate(
        buffer, pagesize=page_size,
        leftMargin=0, rightMargin=0, topMargin=0, bottomMargin=0
    )
    styles = getSampleStyleSheet()

    # --- Template ---
    def on_page(canvas, doc):
        draw_border_and_header(canvas, doc)

    # Frame starts lower to clear the header block
    frame = Frame(1.5*cm, 1.5*cm, WIDTH-3*cm, HEIGHT-4.0*cm, id='f1', showBoundary=0)
    tpl = PageTemplate(id='OthersList', frames=[frame], onPage=on_page)
    doc.addPageTemplates([tpl])

    # --- Build Story ---
    story = []
    
    # Grid Layout 4x2
    grid_data = []
    current_row = []
    col_w = (WIDTH - 3*cm - 20) / 2
    
    for item in others_items:
        # Text Data
        raw_data = _get_item_pdf_data(item, "General") # "General" as generic label
        tbl_rows = []
        for row in raw_data:
            label, sep, val = row
            p_val = Paragraph(val, styles["Normal"])
            p_lbl = Paragraph(f"<b>{label}</b>", styles["Normal"])
            tbl_rows.append([p_lbl, sep, p_val])

        # Inner Table
        inner_tbl = Table(tbl_rows, colWidths=[65, 5, col_w - 90 - 70]) 
        inner_tbl.setStyle(TableStyle([
            ('VALIGN', (0,0), (-1,-1), 'TOP'),
            ('LEFTPADDING', (0,0), (-1,-1), 0),
            ('RIGHTPADDING', (0,0), (-1,-1), 0),
            ('TOPPADDING', (0,0), (-1,-1), 0),
            ('BOTTOMPADDING', (0,0), (-1,-1), 0),
        ]))

        # Image
        img_bytes = _get_img_bytes(item)
        img_obj = Spacer(80, 60) 
        if img_bytes:
            try:
                rdr = ImageReader(BytesIO(img_bytes))
                iw, ih = rdr.getSize()
                scale = min(80.0/iw, 60.0/ih)
                img_obj = RLImage(BytesIO(img_bytes), width=iw*scale, height=ih*scale)
            except: pass
        
        # Item Container
        item_data = [[img_obj, inner_tbl]]
        item_tbl = Table(item_data, colWidths=[90, col_w - 90])
        item_tbl.setStyle(TableStyle([
            ('BOX', (0,0), (-1,-1), 0.5, HexColor("#cccccc")),
            ('VALIGN', (0,0), (-1,-1), 'MIDDLE'),
            ('LEFTPADDING', (0,0), (-1,-1), 4),
            ('RIGHTPADDING', (0,0), (-1,-1), 4),
            ('TOPPADDING', (0,0), (-1,-1), 4),
            ('BOTTOMPADDING', (0,0), (-1,-1), 4),
        ]))
        
        current_row.append(item_tbl)
        
        if len(current_row) == 2:
            grid_data.append(current_row)
            current_row = []
    
    if current_row:
        current_row.append(Spacer(1,1))
        grid_data.append(current_row)
        
    main_tbl = Table(grid_data, colWidths=[col_w, col_w])
    main_tbl.setStyle(TableStyle([
        ('VALIGN', (0,0), (-1,-1), 'TOP'),
        ('LEFTPADDING', (0,0), (-1,-1), 0),
        ('RIGHTPADDING', (0,0), (-1, -1), 10),
        ('BOTTOMPADDING', (0,0), (-1,-1), 5),
    ]))
    
    story.append(main_tbl)

    try:
        doc.build(story)
        buffer.seek(0)
        return buffer.read()
    except Exception as e:
        st.error(f"PDF Generation Error: {e}")
        return None

def export_all_rooms_pdf(cart: list[dict], client_name: str = "", approved_by: str = "") -> Optional[bytes]:
    """
    Generates a single merged PDF containing visual layouts for ALL rooms that have items.
    Order: Drawing Room -> Living Room -> Dining Room -> Bedrooms (1-3) -> Others.
    """
    try:
        from pypdf import PdfWriter, PdfReader
        from io import BytesIO
    except ImportError:
        st.error("Install `pypdf` to enable merged export: pip install pypdf")
        return None

    writer = PdfWriter()
    has_pages = False

    # Define export sequence
    # (Function, Room Key/Arg)
    tasks = [
        (export_drawing_room_pdf, None),
        (export_living_room_pdf, None),
        (export_dining_room_pdf, None),
        (export_bedroom_pdf, "Bedroom 1"),
        (export_bedroom_pdf, "Bedroom 2"),
        (export_bedroom_pdf, "Bedroom 3"),
        (export_others_room_pdf, None),
    ]

    for func, arg in tasks:
        try:
            if arg:
                pdf_bytes = func(cart, arg, client_name, approved_by)
            else:
                pdf_bytes = func(cart, client_name, approved_by)
            
            if pdf_bytes:
                reader = PdfReader(BytesIO(pdf_bytes))
                for page in reader.pages:
                    writer.add_page(page)
                has_pages = True
        except Exception:
            pass

    if not has_pages:
        return None

    out = BytesIO()
    writer.write(out)
    out.seek(0)
    return out.read()

def export_pro_forma_invoice_pdf(cart: list[dict], client_name: str = "", approved_by: str = "") -> Optional[bytes]:
    """
    Generates a professional Room-based Pro-forma Invoice PDF.
    Includes Price Overrides.
    """
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import A4, landscape
    from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.enums import TA_RIGHT, TA_CENTER
    from reportlab.lib.units import cm
    from datetime import datetime

    buffer = BytesIO()
    # [UPDATED] Use landscape to prevent overlapping text with extra columns
    doc = SimpleDocTemplate(buffer, pagesize=landscape(A4), topMargin=1.5*cm, bottomMargin=1.5*cm, leftMargin=1.0*cm, rightMargin=1.0*cm)
    styles = getSampleStyleSheet()
    
    # Custom Styles
    style_title = ParagraphStyle('InvoiceTitle', parent=styles['Heading1'], alignment=TA_CENTER, fontSize=16, spaceAfter=12)
    style_room = ParagraphStyle('RoomHeader', parent=styles['Heading3'], fontSize=12, spaceBefore=10, spaceAfter=6, textColor=colors.HexColor("#2c3e50"))
    style_normal = ParagraphStyle('NormalSmall', parent=styles["Normal"], fontSize=9, leading=11)
    style_right = ParagraphStyle('RightAlign', parent=style_normal, alignment=TA_RIGHT)
    style_bold_right = ParagraphStyle('RightAlignBold', parent=style_normal, alignment=TA_RIGHT, fontName="Helvetica-Bold")

    elements = []

    # --- 1. Header ---
    elements.append(Paragraph("Pro-forma Invoice / Estimate", style_title))
    
    date_str = datetime.now().strftime("%d-%b-%Y")
    header_data = [
        [Paragraph(f"<b>Client:</b> {client_name}", style_normal), Paragraph(f"<b>Date:</b> {date_str}", style_right)],
        [Paragraph(f"<b>Approved By:</b> {approved_by}", style_normal), ""]
    ]
    t_head = Table(header_data, colWidths=[15*cm, 12*cm]) # Adjusted for landscape
    t_head.setStyle(TableStyle([
        ('VALIGN', (0,0), (-1,-1), 'TOP'),
        ('LINEBELOW', (0,1), (-1,1), 1, colors.black),
        ('BOTTOMPADDING', (0,1), (-1,1), 10),
    ]))
    elements.append(t_head)
    elements.append(Spacer(1, 10))

    # --- 2. Prepare Grouped Data ---
    rooms = {}
    grand_total = 0
    
    for item in cart:
        qty_total = int(item.get("qty", 1))
        price = _get_item_price(item)
        alloc = item.get("room_alloc", {})
        
        assigned_qty = 0
        
        # Assign to specific rooms
        if isinstance(alloc, dict):
            for r, val in alloc.items():
                q = _room_qty(val)
                # Extract visual role (label)
                role = ""
                if isinstance(val, dict): role = val.get("visual_role", "")
                
                if q > 0:
                    if r not in rooms: rooms[r] = []
                    sub = q * price
                    rooms[r].append({"item": item, "qty": q, "price": price, "subtotal": sub, "role": role})
                    assigned_qty += q
                    grand_total += sub
        
        # Assign remainder to Unassigned
        rem = qty_total - assigned_qty
        if rem > 0:
            r = "Unassigned / Others"
            if r not in rooms: rooms[r] = []
            sub = rem * price
            rooms[r].append({"item": item, "qty": rem, "price": price, "subtotal": sub, "role": ""})
            grand_total += sub

    #sort_order = ["Drawing Room", "Living Room", "Bedroom", "Dining Room", "Kitchen"]
    sort_order = ["Drawing Room", "Living Room", "Dining Room", "Kitchen", "Bedroom 1", "Bedroom 2", "Bedroom 3"]
    sorted_keys = [k for k in sort_order if k in rooms] + [k for k in rooms if k not in sort_order]

    # --- 3. Build Tables per Room ---
    # [UPDATED] Columns: Brand, Product Info, Context, Qty, Unit, Total
    # Total width approx 27cm
    col_widths = [3.0*cm, 8.5*cm, 6.0*cm, 2.0*cm, 3.5*cm, 4.0*cm] 
    
    for room in sorted_keys:
        items = rooms[room]
        if not items: continue
        
        elements.append(Paragraph(f"{room}", style_room))
        
        table_data = [[
            Paragraph("<b>Brand</b>", style_normal),
            Paragraph("<b>Product / Details</b>", style_normal),
            Paragraph("<b>Type / Loc / Label</b>", style_normal), # NEW COLUMN
            Paragraph("<b>Qty</b>", style_right),
            Paragraph("<b>Unit Price</b>", style_right),
            Paragraph("<b>Total</b>", style_right)
        ]]
        
        room_total = 0
        
        for entry in items:
            it = entry["item"]
            
            # Product Description
            desc = f"<b>{it.get('Product name','')}</b>"
            if it.get('Name of subtype'): desc += f"<br/>{it['Name of subtype']}"
            if it.get('W'): desc += f"<br/>Dim: {it['W']}x{it['D']}x{it['H']}"
            
            # Context Info (Type, Room, Label)
            ptype = it.get("Product type", "")
            label = entry.get("role", "")
            context_lines = [f"Type: {ptype}"]
            context_lines.append(f"Room: {room}")
            if label: context_lines.append(f"Label: {label}")
            context_str = "<br/>".join(context_lines)

            row = [
                Paragraph(it.get("Brand", ""), style_normal),
                Paragraph(desc, style_normal),
                Paragraph(context_str, style_normal), # NEW DATA
                Paragraph(str(entry["qty"]), style_right),
                Paragraph(f"{entry['price']:,.0f}", style_right),
                Paragraph(f"{entry['subtotal']:,.0f}", style_right)
            ]
            table_data.append(row)
            room_total += entry["subtotal"]
            
        table_data.append([
            "", "", "", "",
            Paragraph("<b>Subtotal:</b>", style_right),
            Paragraph(f"<b>{room_total:,.0f}</b>", style_bold_right)
        ])
        
        t = Table(table_data, colWidths=col_widths, repeatRows=1)
        t.setStyle(TableStyle([
            ('BACKGROUND', (0,0), (-1,0), colors.HexColor("#f2f2f2")),
            ('TEXTCOLOR', (0,0), (-1,0), colors.black),
            ('ALIGN', (3,0), (-1,-1), 'RIGHT'), # Qty onwards right aligned
            ('VALIGN', (0,0), (-1,-1), 'TOP'),
            ('INNERGRID', (0,0), (-1,-2), 0.25, colors.grey),
            ('BOX', (0,0), (-1,-2), 0.25, colors.black),
            ('LINEABOVE', (0,-1), (-1,-1), 1, colors.black),
            ('BOTTOMPADDING', (0,0), (-1,-1), 6),
            ('TOPPADDING', (0,0), (-1,-1), 6),
        ]))
        elements.append(t)
        elements.append(Spacer(1, 15))

    # --- 4. Grand Total ---
    elements.append(Spacer(1, 10))
    gt_style = ParagraphStyle('GT', parent=styles['Heading3'], alignment=TA_RIGHT, fontSize=14)
    gt_data = [[
        "", 
        Paragraph("<b>GRAND TOTAL (INR):</b>", gt_style),
        Paragraph(f"<b>{grand_total:,.0f}</b>", gt_style)
    ]]
    t_gt = Table(gt_data, colWidths=[19.5*cm, 4.5*cm, 4.0*cm])
    t_gt.setStyle(TableStyle([
        ('LINEABOVE', (1,0), (2,0), 1.5, colors.black),
        ('TOPPADDING', (0,0), (-1,-1), 12),
    ]))
    elements.append(t_gt)

    try:
        doc.build(elements)
        buffer.seek(0)
        return buffer.read()
    except Exception as e:
        st.error(f"Invoice PDF Generation Error: {e}")
        return None

# ---------------------------
# Page: Dashboard (catalog admin)
# ---------------------------
if st.session_state.page == "dashboard":

    # Dashboard-specific auth (separate password)
    if not gate_dashboard_auth():
        st.stop()

    render_nav(
        prefix="dash",
        on_back=back,
        on_forward=forward,
        on_home=go_home,
        on_products=lambda: go("products"),
    )

    st.divider()

    st.title("Catalog Dashboard")
    #df = load_data(DATA_PATH)
    df = load_data_from_db(DATABASE_URL)

    st.info("Filter the table; edit cells inline; add or delete rows; click **Save Changes** to persist.")

    # --- Filters ---
    fcols = st.columns(5)
    with fcols[0]:
        brand_opt = ["All"] + sorted([x for x in df["Brand"].dropna().unique() if str(x).strip()])
        pick_brand = st.selectbox("Brand", brand_opt, index=0)
    with fcols[1]:
        type_opt = ["All"] + sorted([x for x in df["Product type"].dropna().unique() if str(x).strip()])
        pick_type = st.selectbox("Product type", type_opt, index=0)
    with fcols[2]:
        name_q = st.text_input("Search product name")
    with fcols[3]:
        sub_q = st.text_input("Search subtype")
    with fcols[4]:
        oc_q = st.text_input("Search Other Categories")

    mask = pd.Series(True, index=df.index)
    if pick_brand != "All":
        mask &= (df["Brand"] == pick_brand)
    if pick_type != "All":
        mask &= (df["Product type"] == pick_type)
    if name_q:
        mask &= df["Product name"].str.contains(name_q, case=False, na=False)
    if sub_q:
        mask &= df["Name of subtype"].str.contains(sub_q, case=False, na=False)
    if oc_q and "Other Categories" in df.columns:
        mask &= df["Other Categories"].str.contains(oc_q, case=False, na=False)

    view = df.loc[mask].copy()

    colcfg = {
        "W": st.column_config.NumberColumn(format="%.2f", step=1.0),
        "D": st.column_config.NumberColumn(format="%.2f", step=1.0),
        "H": st.column_config.NumberColumn(format="%.2f", step=1.0),
        "Fabric Category Price": st.column_config.NumberColumn(step=1.0),
        "Leather Category Price": st.column_config.NumberColumn(step=1.0),
        "Other Price": st.column_config.NumberColumn(step=1.0),
        "Link": st.column_config.LinkColumn(),
        "Composition link": st.column_config.LinkColumn(),
        "Composition parts": st.column_config.TextColumn(),
        "Composition part image URLs": st.column_config.TextColumn(),
    }

    st.markdown("#### Edit Catalog")
    edited = st.data_editor(
        view,
        column_config=colcfg,
        use_container_width=True,
        num_rows="dynamic",
        key="editor_table",
    )

    # --- Add New Row (clean vertical layout) ---
    with st.expander("Add a new product row / add subtype to an existing product"):
        b_choices = sorted([x for x in df["Brand"].dropna().unique() if str(x).strip()])
        t_choices = sorted([x for x in df["Product type"].dropna().unique() if str(x).strip()])

        st.markdown("#### Basics")
        # Brand
        b_sel = st.selectbox(
            "Brand (or choose '<Add new>' to type a new one)",
            ["<Add new>"] + b_choices,
            key="add_brand_sel"
        )
        brand_val = (
            st.text_input("New Brand", value="", key="add_brand_new")
            if b_sel == "<Add new>" else b_sel
        )

        # Product type
        t_sel = st.selectbox(
            "Product type (or choose '<Add new>' to type a new one)",
            ["<Add new>"] + t_choices,
            key="add_type_sel"
        )
        type_val = (
            st.text_input("New Product type", value="", key="add_type_new")
            if t_sel == "<Add new>" else t_sel
        )

        # Optional: pick an existing Product (by Brand + Type) to add a new subtype under it
        existing_prod_label = "<Add new>"
        existing_products_df = pd.DataFrame()
        if brand_val and type_val and brand_val != "<Add new>" and type_val != "<Add new>":
            existing_products_df = (
                df[(df["Brand"] == brand_val) & (df["Product type"] == type_val)]
                [["product_id", "Product name", "Link"]]
                .dropna(subset=["product_id"])
                .drop_duplicates(subset=["product_id"])
                .sort_values(["Product name"])
            )
        prod_options = [existing_prod_label] + [
            f'{r["Product name"]} (id:{int(r["product_id"])})' for _, r in existing_products_df.iterrows()
        ]
        picked_label = st.selectbox(
            "Existing product (optional, to add more subtypes under it)",
            options=prod_options,
            key="add_existing_product_pick"
        )

        # Determine product name + default link depending on selection
        if picked_label != existing_prod_label and not existing_products_df.empty:
            # user chose an existing product → lock the name; auto-fill link
            picked_pid_str = picked_label.split("id:")[-1].rstrip(")")
            try:
                picked_pid = int(picked_pid_str)
            except Exception:
                picked_pid = None

            row_match = existing_products_df[existing_products_df["product_id"] == picked_pid]
            _prefill_name = row_match.iloc[0]["Product name"] if not row_match.empty else ""
            _prefill_link = row_match.iloc[0]["Link"] if not row_match.empty else ""

            pname = st.text_input("Product name (existing)", value=_prefill_name, disabled=True, key="add_pname_existing")
            link_val = st.text_input("Product link (URL)", value=_prefill_link or "", key="add_link_existing")
        else:
            # new product flow
            pname = st.text_input("Product name", key="add_pname_new")
            link_val = st.text_input("Product link (URL)", key="add_link_new")

        comp_link_val = st.text_input("Composition link (image URL)", key="add_comp_link")
        comp_parts_val = st.text_area("Composition parts (pipe | separated labels, in order)", key="add_comp_parts")
        comp_part_urls_val = st.text_area(
            "Composition part image URLs (pipe | separated, match the order of parts)",
            key="add_comp_part_urls"
        )

        st.markdown("#### Dimensions")
        w = st.text_input("W", key="add_w")
        d = st.text_input("D", key="add_d")
        h = st.text_input("H", key="add_h")

        st.markdown("#### Pricing & Categories")
        fcat   = st.text_input("Fabric Category", key="add_fcat")
        fprice = st.text_input("Fabric Category Price", key="add_fprice")
        lcat   = st.text_input("Leather Category", key="add_lcat")
        lprice = st.text_input("Leather Category Price", key="add_lprice")
        oc     = st.text_input("Other Categories", key="add_oc")
        oprice = st.text_input("Other Price", key="add_oprice")

        st.markdown("#### Subtype")
        subtype = st.text_input("Name of subtype", key="add_subtype")

        if st.button("Add Row (product or new subtype)", use_container_width=True):
            new_row = {
                "Brand": (brand_val or "").strip(),
                "Product name": (pname or "").strip(),
                "Link": (link_val or "").strip(),
                "Product type": (type_val or "").strip(),
                "Name of subtype": (subtype or "").strip(),
                "W": _num_or_none(w),
                "D": _num_or_none(d),
                "H": _num_or_none(h),
                "Fabric Category": (fcat or "").strip(),
                "Leather Category": (lcat or "").strip(),
                "Other Categories": (oc or "").strip(),
                "Fabric Category Price": _num_or_none(fprice),
                "Leather Category Price": _num_or_none(lprice),
                "Other Price": _num_or_none(oprice),
                "Composition link": (comp_link_val or "").strip(),
                "Composition parts": (comp_parts_val or "").strip(),
                "Composition part image URLs": (comp_part_urls_val or "").strip(),
            }
            if not new_row["Brand"] or not new_row["Product name"] or not new_row["Product type"]:
                st.error("Brand, Product name, and Product type are required.")
            else:
                # Note: upsert_row_to_db will create the product if new, or attach the new
                # subtype/prices to the chosen existing product automatically.
                upsert_row_to_db(DATABASE_URL, new_row)
                st.success("Row added to database (new product or new subtype).")
                st.rerun()

    # --- Delete Selected Rows (safe defaults) ---
    st.markdown("#### Delete Rows")
    to_delete = safe_multiselect(
        "Select rows to delete (by index from the edited table above)",
        options=list(edited.index),
        default=st.session_state.to_delete_defaults,
        format_func=lambda idx: f"{idx} — {edited.loc[idx, 'Brand']} / {edited.loc[idx, 'Product name']}",
        key="delete_rows_selection",
    )
    st.session_state.to_delete_defaults = to_delete

    if st.button("Delete Selected"):
        if to_delete:
            rows_to_delete = edited.loc[to_delete]
            deleted = delete_rows_in_db(DATABASE_URL, rows_to_delete)
            st.session_state.to_delete_defaults = []
            st.success(f"Deleted {deleted} row(s) from database.")
            st.rerun()
        else:
            st.warning("No rows selected.")

    # --- Update Existing Product & Variant (all attributes) ---
    with st.expander("Update existing product"):
        # unique products list
        prod_list = (
            df[["product_id", "Brand", "Product type", "Product name", "Link", "Composition link"]]
            .dropna(subset=["product_id"])
            .drop_duplicates(subset=["product_id"])
            .sort_values(["Brand", "Product type", "Product name"])
        )

        if prod_list.empty:
            st.info("No products available.")
        else:
            prod_labels = {
                int(r.product_id): f"{r['Brand']} · {r['Product type']} · {r['Product name']} (id:{int(r.product_id)})"
                for _, r in prod_list.iterrows()
            }
            sel_pid = st.selectbox(
                "Pick a product to update",
                options=list(prod_labels.keys()),
                format_func=lambda k: prod_labels[k],
                key="upd_pick_product",
            )

            # Current product fields
            current_prod = prod_list[prod_list["product_id"] == sel_pid].iloc[0]
            u_brand = st.text_input("Brand", value=current_prod["Brand"])
            u_type  = st.text_input("Product type", value=current_prod["Product type"])
            u_name  = st.text_input("Product name", value=current_prod["Product name"])
            u_link  = st.text_input("Link", value=current_prod.get("Link") or "")
            u_comp  = st.text_input("Composition link", value=current_prod.get("Composition link") or "")

            # --- VARIANT picker for this product
            vdf = df[df["product_id"] == sel_pid][["variant_id","Name of subtype","W","D","H"]].dropna(subset=["variant_id"])
            vdf = vdf.drop_duplicates(subset=["variant_id"]).copy()

            def _fmt(v):
                if pd.isna(v): return ""
                try:
                    vf = float(v)
                    return str(int(vf)) if vf.is_integer() else str(vf)
                except Exception:
                    return str(v)

            var_labels = {
                int(r.variant_id): f"{r.get('Name of subtype') or '(no subtype)'} · {_fmt(r['W'])}×{_fmt(r['D'])}×{_fmt(r['H'])} (vid:{int(r.variant_id)})"
                for _, r in vdf.sort_values(["Name of subtype","W","D","H"]).iterrows()
            }
            sel_vid = st.selectbox(
                "Pick a variant to update",
                options=list(var_labels.keys()),
                format_func=lambda k: var_labels[k],
                key="upd_pick_variant",
            )

            # Pull current variant+price values
            vrows = df[(df["product_id"] == sel_pid) & (df["variant_id"] == sel_vid)].copy()
            vfirst = vrows.iloc[0] if not vrows.empty else pd.Series({})

            def _first(series):
                for x in series:
                    if pd.notna(x) and str(x).strip() != "":
                        return x
                return ""

            st.markdown("#### Variant")
            c1, c2, c3, c4 = st.columns(4)
            with c1:
                u_subtype = st.text_input("Name of subtype", value=str(vfirst.get("Name of subtype") or ""))
            with c2:
                u_w = st.text_input("W", value=str(vfirst.get("W") or ""))
            with c3:
                u_d = st.text_input("D", value=str(vfirst.get("D") or ""))
            with c4:
                u_h = st.text_input("H", value=str(vfirst.get("H") or ""))

            st.markdown("#### Pricing")
            fc1, fc2 = st.columns(2)
            with fc1:
                u_fab_cat   = st.text_input("Fabric Category", value=str(_first(vrows["Fabric Category"])) )
            with fc2:
                u_fab_price = st.text_input("Fabric Category Price", value=str(_first(vrows["Fabric Category Price"])) )
            lc1, lc2 = st.columns(2)
            with lc1:
                u_lea_cat   = st.text_input("Leather Category", value=str(_first(vrows["Leather Category"])) )
            with lc2:
                u_lea_price = st.text_input("Leather Category Price", value=str(_first(vrows["Leather Category Price"])) )
            oc1, oc2 = st.columns(2)
            with oc1:
                u_oth_cat   = st.text_input("Other Categories", value=str(_first(vrows["Other Categories"])) )
            with oc2:
                u_oth_price = st.text_input("Other Price", value=str(_first(vrows["Other Price"])) )

            st.markdown("#### Composition (variant-scoped)")
            cc1, cc2 = st.columns(2)
            with cc1:
                u_comp_link  = st.text_input("Composition link (image URL)", value=str(_first(vrows["Composition link"])) )
            with cc2:
                u_comp_parts = st.text_area("Composition parts (pipe | separated)", value=str(_first(vrows["Composition parts"])) )
            u_comp_part_urls = st.text_area(
                "Composition part image URLs (pipe | separated, match order of parts)",
                value=str(_first(vrows["Composition part image URLs"]))
            )

            if st.button("Update Product & Variant", use_container_width=True, key="btn_update_product_variant"):
                try:
                    # product-level (unchanged)
                    update_product_details(
                        DATABASE_URL,
                        int(sel_pid),
                        u_brand, u_type, u_name, u_link, u_comp,
                    )
                    # variant + pricing + compositions
                    update_variant_and_prices(
                        DATABASE_URL,
                        int(sel_vid),
                        u_subtype, u_w, u_d, u_h,
                        u_fab_cat, u_fab_price,
                        u_lea_cat, u_lea_price,
                        u_oth_cat, u_oth_price,
                        u_comp_link, u_comp_parts, u_comp_part_urls,
                    )
                    st.success("Product & variant updated in database.")
                    st.rerun()
                except Exception as e:
                    st.error(str(e))

    if st.button("Save Changes"):
        # Only persist rows that are actually NEW or CHANGED
        saved = 0

        # Columns that matter for catalog persistence
        cols_to_compare = [
            "Brand",
            "Product name",
            "Link",
            "Product type",
            "Name of subtype",
            "W",
            "D",
            "H",
            "Fabric Category",
            "Leather Category",
            "Other Categories",
            "Fabric Category Price",
            "Leather Category Price",
            "Other Price",
            "Composition link",
            "Composition parts",
            "Composition part image URLs",
        ]
        cols_to_compare = [c for c in cols_to_compare if c in edited.columns and c in view.columns]

        def _is_core_empty(row) -> bool:
            """Skip rows that don't even have brand / product / type filled."""
            b = str(row.get("Brand", "")).strip()
            p = str(row.get("Product name", "")).strip()
            t = str(row.get("Product type", "")).strip()
            return not (b and p and t)

        for idx, r in edited.iterrows():
            # Skip totally empty / placeholder rows in the editor
            if _is_core_empty(r):
                continue

            if idx in view.index:
                # Existing row – check if anything actually changed
                orig = view.loc[idx]
                changed = False
                for c in cols_to_compare:
                    if pd.isna(r.get(c)) and pd.isna(orig.get(c)):
                        continue
                    if str(r.get(c)) != str(orig.get(c)):
                        changed = True
                        break
                if not changed:
                    continue  # no diff → skip
            # else: idx not in view.index → this is a brand-new row added in the editor

            # Only here do we touch the DB
            upsert_row_to_db(DATABASE_URL, r.to_dict())
            saved += 1

        if saved > 0:
            st.success(f"Saved {saved} changed/new row(s) to database.")
        else:
            st.info("No changes detected – nothing saved.")
        st.rerun()

# ---------------------------
# Page: Home (Rooms-first landing; product-type gallery with animated tiles)
# ---------------------------
elif st.session_state.page == "home":
    st.markdown("<div style='text-align:center'><h1>Altossa – Product Recommender</h1></div>", unsafe_allow_html=True)
    st.caption("Start by choosing a room. Then pick product types inside that room.")

    rooms = list(st.session_state.room_types.keys())

    # 3x3 compact grid of room buttons — the whole bar is the button (no inner container)
    for i in range(0, len(rooms), 3):
        cols = st.columns(3, gap="large")
        for j, col in enumerate(cols):
            if i + j >= len(rooms):
                continue
            room = rooms[i + j]
            with col:
                # RESET compositions_only when navigating normally
                st.button(
                    room,
                    key=f"open_{room}",
                    use_container_width=True,
                    help=f"Open {room}",
                    on_click=lambda r=room: (
                        st.session_state.update(sel_room=r, sel_type=None, compositions_only=False),
                        go("room")
                    ),
                )

    st.divider()
    st.subheader("…or browse by Product Type")

    # If we arrived via a tile click (?type=...), consume it and jump once.
    _consume_type_query_and_jump()

    # Product Type gallery – each tile uses the FIRST product's own image for that type.
    types = sorted([t for t in df["Product type"].dropna().unique() if t])
    search = st.text_input("Search Product Type")
    if search:
        types = [t for t in types if search.lower() in t.lower()]

    for i in range(0, len(types), 3):
        row_types = types[i:i+3]

        # For this row, figure out the first product's (Link, Composition link) per type
        type_links: dict[str, tuple[str, str]] = {}
        for t in row_types:
            type_links[t] = first_product_links_for_type(df, t)

        # Prefetch just those exact images so tiles paint fast
        pre_pairs = [
            (pl, cl)
            for (pl, cl) in type_links.values()
            if (pl or cl)
        ]
        prefetch_product_images_async(pre_pairs, max_workers=8)

        c1, c2, c3 = st.columns(3, gap="large")
        for j, col in enumerate((c1, c2, c3)):
            if j >= len(row_types):
                continue
            t = row_types[j]
            with col:
                pl, cl = type_links.get(t, ("", ""))
                render_type_tile(t, pl, cl, key=f"type_{i}_{j}")

    # Sofa Compositions shortcut
    st.markdown("---")
    st.subheader("Available Compositions")
    st.caption("Jump straight to sofas that include a composition image and breakdown.")
    st.button(
        "View Compositions",
        use_container_width=True,
        on_click=lambda: (
            st.session_state.update(sel_type="Sofa", sel_room=None, compositions_only=True),
            go("products")
        ),
    )

# ---------------------------
# Page: Room (per-room mapping editor + compact type grid)
# ---------------------------
elif st.session_state.page == "room":
    room = st.session_state.sel_room
    if not room:
        back()

    render_nav(
        prefix="room",
        on_back=back,
        on_forward=forward,
        on_home=go_home,
        on_products=lambda: go("products"),
    )

    st.markdown(f"<h2 style='text-align:center'>{room}</h2>", unsafe_allow_html=True)

    # Per-room mapping editor (moved here, scoped to the selected room)
    with st.expander(f"Configure Product Types for {room}"):
        all_types = sorted([t for t in df["Product type"].dropna().unique() if t])
        current = st.session_state.room_types.get(room, [])
        in_options = [t for t in current if t in all_types]
        custom = [t for t in current if t not in all_types]
        sel = safe_multiselect(
            f"{room} product types",
            options=all_types + custom,
            default=in_options + custom,
            key=f"room_cfg_{room}",
        )
        if sel != current:
            st.session_state.room_types[room] = sel
            st.success("Updated mapping for this room (session).")

    # Compact 3x3 grid of product types for this room
    types_here = st.session_state.room_types.get(room, [])
    if not types_here:
        st.info("No product types added to this room yet.")
    else:
        for i in range(0, len(types_here), 3):
            cols = st.columns(3, gap="large")
            for j, col in enumerate(cols):
                if i + j >= len(types_here):
                    continue
                ptype = types_here[i + j]
                with col:
                    st.button(
                        ptype,
                        use_container_width=True,
                        on_click=lambda r=room, t=ptype: (st.session_state.update(sel_room=r, sel_type=t), go("products")),
                    )

# ---------------------------
# Page: Products (within a type)
# ---------------------------
elif st.session_state.page == "products":

    render_nav(
        prefix="prod",
        on_back=back,
        on_forward=forward,
        on_home=go_home,
        on_products=lambda: go("products"),
    )

    sel_type = st.session_state.sel_type
    sel_room = st.session_state.sel_room
    only_comp = bool(st.session_state.get("compositions_only"))

    # NEW: restore scroll position per product type when returning here
    inject_scroll_restoration(f"products_{sel_type or ''}")

    st.subheader(
        f"{sel_type}"
        + (f" · {sel_room}" if sel_room else "")
        + (" · compositions only" if only_comp else "")
    )

    # base filter by Product type
    df_t = df[df["Product type"].str.lower() == str(sel_type).lower()].copy()

    # when compositions-only is ON, keep only product_ids that have composition info
    if only_comp:
        comp_pids = set(
            df_t.loc[df_t["Composition link"].ne(""), "product_id"]
            .dropna()
            .astype(int)
            .unique()
        )
        df_t = df_t[df_t["product_id"].isin(comp_pids)]

    if df_t.empty:
        st.info("No products yet in this category.")
    
    else:
        # NEW: persistent filters per Product type (explicit state management)
        filter_store = st.session_state.product_filters.get(sel_type, {})
        
        # Keys for widgets
        brand_key = f"prod_brand_filter_{sel_type or 'all'}"
        price_key = f"prod_price_filter_{sel_type or 'all'}"
        search_key = f"prod_search_{sel_type or 'all'}"

        # 1. Brand Filter
        brand_options = ["All"] + sorted(df_t["Brand"].dropna().unique())
        
        # Restore index from store
        saved_brand = filter_store.get("brand", "All")
        try:
            brand_idx = brand_options.index(saved_brand)
        except ValueError:
            brand_idx = 0
            
        brand_filter = st.selectbox(
            "Brand",
            brand_options,
            index=brand_idx,
            key=brand_key,
        )
        
        # Update store if changed
        if brand_filter != saved_brand:
            filter_store["brand"] = brand_filter
            st.session_state.product_filters[sel_type] = filter_store

        if brand_filter != "All":
            df_t = df_t[df_t["Brand"] == brand_filter]

        # 2. Price Range Filter
        def _get_row_min_max_price(row):
            """Return min and max final price for a single row."""
            prices = []
            brand = row.get("Brand", "")
            for col in ["Fabric Category Price", "Leather Category Price", "Other Price"]:
                val = calculate_final_price(brand, row.get(col))
                if val is not None:
                    prices.append(val)
            if not prices:
                return None, None
            return min(prices), max(prices)

        min_prices = []
        max_prices = []
        
        for _, r in df_t.iterrows():
            mn, mx = _get_row_min_max_price(r)
            if mn is not None:
                min_prices.append(mn)
                max_prices.append(mx)
        
        if min_prices:
            global_min = int(min(min_prices))
            global_max = int(max(max_prices))
            
            if global_max > global_min:
                step = 1000
                slider_min = (global_min // step) * step
                slider_max = ((global_max // step) + 1) * step
                
                # Restore value from store
                default_val = (slider_min, slider_max)
                saved_price = filter_store.get("price", default_val)
                
                # Clamp to current bounds to prevent errors if filtered data shrinks range
                s_start = max(slider_min, min(saved_price[0], slider_max))
                s_end = max(slider_min, min(saved_price[1], slider_max))
                if s_start > s_end: s_start = s_end
                
                sel_price_range = st.slider(
                    "Price Range (₹)",
                    min_value=slider_min,
                    max_value=slider_max,
                    value=(s_start, s_end),
                    step=step,
                    key=price_key
                )
                
                # Update store
                if sel_price_range != saved_price:
                    filter_store["price"] = sel_price_range
                    st.session_state.product_filters[sel_type] = filter_store

                # Filter logic
                def _row_matches_price(row, p_min, p_max):
                    mn, mx = _get_row_min_max_price(row)
                    if mn is None: return False
                    return (mn <= p_max) and (mx >= p_min)

                mask = df_t.apply(lambda r: _row_matches_price(r, sel_price_range[0], sel_price_range[1]), axis=1)
                df_t = df_t[mask]

        # 3. Search Filter
        saved_search = filter_store.get("search", "")
        q = st.text_input(
            "Search by product name",
            value=saved_search,
            key=search_key,
        )
        
        if q != saved_search:
            filter_store["search"] = q
            st.session_state.product_filters[sel_type] = filter_store

        if q:
            df_t = df_t[df_t["Product name"].str.contains(q, case=False, na=False)]

        prods = list(df_t.groupby("product_id", sort=False))

        for i in range(0, len(prods), 3):
            # prefetch images in parallel for this row of 3 cards
            slice_items = prods[i:i+3]
            pairs = []
            for pid, g0 in slice_items:
                r0 = g0.iloc[0]
                pairs.append(
                    (
                        r0.get("Link", "") or "",
                        r0.get("Composition link", "") or "",
                    )
                )
            prefetch_product_images(pairs, max_workers=8)

            c1, c2, c3 = st.columns(3)
            for j, col in enumerate([c1, c2, c3]):
                if i + j >= len(prods):
                    continue
                pid, g0 = prods[i + j]
                row0 = g0.iloc[0]
                name = row0["Product name"]
                brand = row0["Brand"]
                with col.container(border=True):
                    st.markdown(f"**{name}**")
                    st.caption(f"Brand: {brand}")
                    if row0.get("Link"):
                        st.markdown(f"[Open product page]({row0['Link']})")

                    img = get_image_cached(
                        row0.get("Link", ""),
                        row0.get("Composition link", ""),
                    )
                    if img:
                        st.image(img)

                    st.button(
                        "Select",
                        key=f"sel_{pid}",
                        on_click=lambda _pid=int(pid), _nm=str(name): on_select_product(
                            _pid, _nm
                        ),
                    )

# ---------------------------
# Page: Sizes (W×D×H) for a product
# ---------------------------
elif st.session_state.page == "sizes":

    render_nav(
        prefix="sizes",
        on_back=back,
        on_forward=forward,
        on_home=go_home,
        on_products=lambda: go("products"),
    )

    sel_type = st.session_state.sel_type
    sel_pid = st.session_state.sel_product_id
    sel_prod = st.session_state.sel_product
    st.subheader(f"{sel_prod} - choose size")

    #g = df[df["product_id"] == int(sel_pid)].copy()
    # Base rows for this product
    g = df[df["product_id"] == int(sel_pid)].copy()

    # If we came from "Available Compositions", keep only the variants that
    # themselves have a composition link. This prevents plain (non-composition)
    # sizes from appearing in this flow.
    if bool(st.session_state.get("compositions_only")):
        g = g[g["Composition link"].astype(str).str.strip() != ""].copy()


    # Build available sizes (dedupe by W, D, H, and Name of subtype)
    cards = _dedupe_size_cards(g)

    if cards:
        # Small formatter to show ints without trailing .0, and keep strings/ranges as-is
        def _fmt(v):
            if pd.isna(v):
                return ""
            try:
                vf = float(v)
                return str(int(vf)) if vf.is_integer() else str(vf)
            except Exception:
                return str(v)

        for i in range(0, len(cards), 3):
            c1, c2, c3 = st.columns(3)
            for j, col in enumerate((c1, c2, c3)):
                if i + j >= len(cards):
                    continue
                s = cards[i + j]
                with col.container(border=True):
                    st.markdown(f"**{_fmt(s['W'])} × {_fmt(s['D'])} × {_fmt(s['H'])}**")
                    if s["subtype"]:
                        st.caption(s["subtype"])

                    # --- Composition image preview (only if this variant has a composition link)
                    comp_preview = str(s["row"].get("Composition link") or "").strip()
                    if comp_preview:
                        # Use the same "big hero" image approach we use elsewhere
                        hero_url = _drive_hero_url(comp_preview, width=1600)
                        comp_img = get_image_cached("", hero_url)   # synchronous fetch; returns bytes or None
                        if comp_img:
                            st.image(comp_img, use_container_width=True)

                    st.button(
                        "Choose this size",
                        key=f"size_{s['idx']}",
                        on_click=lambda rr=s["row"]: (
                            st.session_state.update(sel_size_idx=int(rr.name)),
                            go("materials")
                        ),
                    )
    else:
        # Auto-skip sizes when no dimensions exist
        if g.empty:
            st.warning("No data rows for this product.")
        else:
            st.session_state.sel_size_idx = None  # no size context
            go("materials")
            st.rerun()

# ---------------------------
# Page: Materials & Pricing
# ---------------------------
elif st.session_state.page == "materials":

    render_nav(
        prefix="mat",
        on_back=back,
        on_forward=forward,
        on_home=go_home,
        on_products=lambda: go("products"),
    )

    sel_type = st.session_state.sel_type
    sel_pid = st.session_state.sel_product_id
    sel_prod = st.session_state.sel_product

    # Base subset: all rows for this product
    base = df[df["product_id"] == int(sel_pid)].copy()

    # If a specific variant (size card) was chosen, filter by full variant key:
    # Brand, Product type, Product name, W, D, H, and Name of subtype.
    chosen_idx = st.session_state.get("sel_size_idx")
    if chosen_idx is not None and chosen_idx in df.index:
        selected_row = df.loc[chosen_idx]
        g = _filter_variant_rows(df, selected_row)
    else:
        g = base

    has_cat = (g["Fabric Category"].ne("") | g["Leather Category"].ne("")).any()
    has_other = g["Other Price"].notna().any()

    brand = g["Brand"].iloc[0] if not g.empty else ""
    link  = g["Link"].iloc[0] if not g.empty else ""

    head_l, head_r = st.columns([0.75, 0.25])
    with head_l:
        st.subheader(f"{sel_prod}")
        if link:
            st.markdown(f"[Product page]({link})")

    with head_r:
        # Show "View compositions" ONLY if this variant has any composition labels or URLs
        has_compositions = False
        if not g.empty and ("Composition parts" in g.columns or "Composition part image URLs" in g.columns):
            for _, rr in g.iterrows():
                parts_nonempty = bool(str(rr.get("Composition parts") or "").strip())
                urls_nonempty  = bool(str(rr.get("Composition part image URLs") or "").strip())
                if parts_nonempty or urls_nonempty:
                    has_compositions = True
                    break

        if has_compositions:
            with st.popover("View compositions", use_container_width=True):
                st.markdown("#### Composition breakdown")
                render_compositions_for_variant(g, link)
        else:
            st.empty()  # nothing to show when no compositions

    # ---- Prefetch images (PRIORITIZE composition hero if available)
    comp_link = g["Composition link"].iloc[0] if ("Composition link" in g.columns and not g.empty) else ""

    # ---- Choose image strictly from the SELECTED variant (if any).
    # If the selected row has a composition link, use it; otherwise fall back to product hero.

    # collect composition part image URLs (first non-empty row)
    comp_urls_str = ""
    for _, rr in g.iterrows():
        uu = str(rr.get("Composition part image URLs") or "").strip()
        if uu:
            comp_urls_str = uu
            break

    pairs = []
    if comp_link:
        pairs.append(("", comp_link))      # fetch composition hero first (force composition branch)
    if link:
        pairs.append((link, ""))           # fetch product hero as fallback
    pairs += [("", u.strip()) for u in comp_urls_str.split("|") if u.strip()]
    prefetch_product_images_async(pairs, max_workers=8)

    # ---- Display hero with preference: Composition link if present, else Product hero
    hero_slot = st.empty()

    def _placeholder():
        hero_slot.markdown(
            "<div style='height:260px; background:rgba(255,255,255,0.04); "
            "border-radius:12px; display:flex; align-items:center; justify-content:center;'>"
            "Loading product image…</div>",
            unsafe_allow_html=True,
        )

    shown = False
    if comp_link:
        # Try composition hero from cache
        img = get_image_if_cached("", comp_link)
        if img:
            hero_slot.image(img, use_container_width=True)
            shown = True
        else:
            _placeholder()
            # bounded short poll + gentle rerun to reveal when background fetch completes
            import time as _t
            key = f"rearm_comp_{hash(comp_link)}"
            last = st.session_state.get(key, 0.0)
            for _ in range(4):  # up to 1.0s
                _t.sleep(0.25)
                img = get_image_if_cached("", comp_link)
                if img:
                    hero_slot.image(img, use_container_width=True)
                    shown = True
                    break
            if not shown and (_t.time() - float(last) > 2.0):
                st.session_state[key] = _t.time()
                st.rerun()
    else:
        # No composition → show product hero (cached) or placeholder + nudge
        img = get_image_if_cached(link, "")
        if img:
            hero_slot.image(img, use_container_width=True)
            shown = True
        else:
            _placeholder()
            import time as _t
            key = f"rearm_prod_{hash(link)}"
            last = st.session_state.get(key, 0.0)
            for _ in range(4):  # up to 1.0s
                _t.sleep(0.25)
                img = get_image_if_cached(link, "")
                if img:
                    hero_slot.image(img, use_container_width=True)
                    shown = True
                    break
            if not shown and (_t.time() - float(last) > 2.0):
                st.session_state[key] = _t.time()
                st.rerun()

    # Case A: Fabric/Leather categories
    if has_cat:
        st.markdown("### Select a material")
        left, right = st.columns(2)

        fab = g[g["Fabric Category"].ne("") & g["Fabric Category Price"].notna()] \
                .sort_values(["Fabric Category", "Fabric Category Price"], ascending=[True, False])
        fab = fab.groupby("Fabric Category", as_index=False).first()

        for _, r in fab.iterrows():
            btn = left.button(
                f"Fabric – {r['Fabric Category']} · ₹{calculate_final_price(brand, r['Fabric Category Price']):,}",
                key=f"fab_{_}"
            )
            if btn:
                st.session_state.cart.append({
                    "Brand": brand,
                    "Product type": sel_type,
                    "Product name": sel_prod,
                    "Name of subtype": r.get("Name of subtype", ""),
                    "W": r["W"], "D": r["D"], "H": r["H"],
                    "Material": r["Fabric Category"],
                    "Other Category": "",
                    "Base Price": float(r["Fabric Category Price"]),
                    "qty": 1,
                    "Link": link,
                    "Composition link": r.get("Composition link", ""),
                    "Composition parts": r.get("Composition parts", ""),
                })
                st.success("Added to cart")
                st.rerun()

        lea = g[g["Leather Category"].ne("") & g["Leather Category Price"].notna()] \
                .sort_values(["Leather Category", "Leather Category Price"], ascending=[True, False])
        lea = lea.groupby("Leather Category", as_index=False).first()

        for _, r in lea.iterrows():
            btn = right.button(
                f"Leather – {r['Leather Category']} · ₹{calculate_final_price(brand, r['Leather Category Price']):,}",
                key=f"lea_{_}"
            )
            if btn:
                st.session_state.cart.append({
                    "Brand": brand,
                    "Product type": sel_type,
                    "Product name": sel_prod,
                    "Name of subtype": r.get("Name of subtype", ""),
                    "W": r["W"], "D": r["D"], "H": r["H"],
                    "Material": r["Leather Category"],
                    "Other Category": "",
                    "Base Price": float(r["Leather Category Price"]),
                    "qty": 1,
                    "Link": link,
                    "Composition link": r.get("Composition link", ""),
                    "Composition parts": r.get("Composition parts", ""),
                })
                st.success("Added to cart")
                st.rerun()

        if fab.empty and lea.empty:
            st.info("No valid fabric/leather prices found for this selection.")

    # Case B: No fabric/leather, use Other Price (with optional Other Categories)
    elif has_other:
        st.markdown("### Select a configuration")
        if "Other Categories" in g.columns and (g["Other Categories"].ne("")).any():
            other = g[g["Other Categories"].ne("") & g["Other Price"].notna()] \
                    .sort_values(["Other Categories", "Other Price"], ascending=[True, False])
            grp_cols = ["Other Categories", "W", "D", "H", "Name of subtype"]
            other = other.groupby(grp_cols, as_index=False).first()

            for i, r in other.iterrows():
                def _fmt_dim(v):
                    if pd.isna(v):
                        return ""
                    try:
                        vf = float(v)
                        return str(int(vf)) if float(vf).is_integer() else str(vf)
                    except Exception:
                        # Non-numeric values (ranges like "260-330", "260/330", etc.) should be returned as-is
                        return str(v)

                dims = " × ".join(
                    _fmt_dim(v) for v in [r["W"], r["D"], r["H"]] if not pd.isna(v)
                )
                label = r["Other Categories"]
                if dims:
                    label += f" · {dims}"
                if r.get("Name of subtype"):
                    label += f" · {r['Name of subtype']}"

                # --- NEW: safe final price text (handles NaN / missing)
                raw_price = r["Other Price"]
                final_price = calculate_final_price(brand, raw_price)
                price_text = f"₹{final_price:,}" if final_price is not None else "—"

                btn = st.button(f"{label} – {price_text}", key=f"oc_{i}")
                if btn:
                    base_price = float(raw_price) if pd.notna(raw_price) else 0.0
                    st.session_state.cart.append({
                        "Brand": brand,
                        "Product type": sel_type,
                        "Product name": sel_prod,
                        "Name of subtype": r.get("Name of subtype", ""),
                        "W": r["W"], "D": r["D"], "H": r["H"],
                        "Material": "",
                        "Other Category": r["Other Categories"],
                        "Base Price": base_price,
                        "qty": 1,
                        "Link": link,
                        "Composition link": r.get("Composition link", ""),
                        "Composition parts": r.get("Composition parts", ""),
                    })
                    st.success("Added to cart")
                    st.rerun()

        else:
            # No named "Other Categories" – show UNIQUE configurations by (W, D, H, subtype)
            gx = (
                g.sort_values(by=["W", "D", "H", "Other Price"], ascending=[True, True, True, False])
                 .drop_duplicates(subset=["Other Categories", "Name of subtype", "W", "D", "H"])
            )

            def _fmt_dim_any(v):
                if pd.isna(v):
                    return ""
                try:
                    vf = float(v)
                    return str(int(vf)) if vf.is_integer() else str(vf)
                except Exception:
                    # Non-numeric strings like "25/30/40/50" should be shown as-is
                    return str(v)

            for i, r in gx.iterrows():
                dims = " × ".join(
                    _fmt_dim_any(v)
                    for v in [r["W"], r["D"], r["H"]]
                    if not pd.isna(v)
                )

                # When Other Categories is empty, label using subtype + dimensions
                label_parts = []
                if r.get("Name of subtype"):
                    label_parts.append(str(r["Name of subtype"]))
                if dims:
                    label_parts.append(dims)
                label = " · ".join([p for p in label_parts if p])

                raw_price = r["Other Price"]
                final_price = calculate_final_price(brand, raw_price)
                price_text = f"₹{final_price:,}" if final_price is not None else "—"

                btn = st.button(f"{label or 'Configuration'} – {price_text}", key=f"oth_{i}")
                if btn:
                    base_price = float(raw_price) if pd.notna(raw_price) else 0.0
                    st.session_state.cart.append({
                        "Brand": brand,
                        "Product type": sel_type,
                        "Product name": sel_prod,
                        "Name of subtype": r.get("Name of subtype", ""),
                        "W": r["W"], "D": r["D"], "H": r["H"],
                        "Material": "",
                        "Other Category": r.get("Other Categories", ""),
                        "Base Price": base_price,
                        "qty": 1,
                        "Link": link
                    })
                    st.success("Added to cart")
                    st.rerun()

    # Case C: No categories, no sizes -> top 2 prices
    else:
        st.markdown("### Select price")
        top = g.sort_values(by="Other Price", ascending=False).head(2)
        if top.empty:
            st.warning("No price data available for this product.")
        for i, r in top.iterrows():
            price = r["Other Price"]
            oc_label = r.get("Other Categories", "")
            final_price = calculate_final_price(brand, price)
            price_txt = f"₹{final_price:,}" if final_price is not None else "—"
            prefix = f"({oc_label}) " if oc_label else ""
            text = f"{prefix}{price_txt}"
            btn = st.button(text, key=f"nop_{i}")
            if btn:
                st.session_state.cart.append({
                    "Brand": brand,
                    "Product type": sel_type,
                    "Product name": sel_prod,
                    "Name of subtype": r.get("Name of subtype", ""),
                    "W": r["W"], "D": r["D"], "H": r["H"],
                    "Material": "",
                    "Other Category": oc_label,
                    "Base Price": float(price) if pd.notna(price) else 0.0,
                    "qty": 1,
                    "Link": link
                })
                st.success("Added to cart")
                st.rerun()

# ---------------------------
# Page: Review & Notes (selected items)
# ---------------------------
elif st.session_state.page == "review":
    # Back button that first flushes uploader state into cart
    def _back_from_review():
        _flush_review_uploads_to_cart()
        back()

    def _forward_from_review():
        _flush_review_uploads_to_cart()
        forward()

    def _home_from_review():
        _flush_review_uploads_to_cart()
        go_home()

    render_nav(
        prefix="rev",
        on_back=_back_from_review,
        on_forward=_forward_from_review,
        on_home=_home_from_review,
        on_products=lambda: (_flush_review_uploads_to_cart(), go("products")),
    )

    st.title("Selection Review")

    # --- Client Info (Shared) ---
    # We read from sidebar state to initialize these inputs.
    # We do NOT write back to 'client_name_sidebar' here to avoid StreamlitAPIException.
    # If the user edits these fields on the main page, they are used for this page's exports.
    with st.container(border=True):
        c1, c2 = st.columns(2)
        
        # Get current values from sidebar (or empty default)
        sb_client = st.session_state.get("client_name_sidebar", "")
        sb_approved = st.session_state.get("approved_by_sidebar", "")
        
        # Main page inputs initialized with sidebar values
        client_name = c1.text_input("Client name", value=sb_client, key="client_name_main")
        approved_by = c2.text_input("Approved by", value=sb_approved, key="approved_by_main")
        
        # Note: removed the manual assignment back to session_state keys 
        # to prevent "cannot be modified after instantiation" error.

    cart = st.session_state.cart
    if not cart:
        st.info("No products selected yet.")
        st.stop()

    # --- TABS: Detailed View vs Invoice View ---
    tab_detail, tab_invoice = st.tabs(["Detailed Review", "Pro-forma Invoice & Pricing"])
    
    # CSS to enlarge tab text
    st.markdown("""
    <style>
    [data-baseweb="tab-list"] button {
        font-size: 50px !important;
        padding: 12px 24px !important;
        font-weight: 600 !important;
    }
    </style>
    """, unsafe_allow_html=True)

    # ==================================================
    # TAB 1: DETAILED REVIEW (Existing Logic)
    # ==================================================
    with tab_detail:
        # --- Filters
        f1, f2, f3 = st.columns(3)
        with f1:
            brand_opts = ["All brands"] + sorted({it.get("Brand") for it in cart if it.get("Brand")})
            sel_brand = st.selectbox("Brand", brand_opts, key="rev_filter_brand")
        with f2:
            type_opts = ["All types"] + sorted({it.get("Product type") for it in cart if it.get("Product type")})
            sel_type = st.selectbox("Product type", type_opts, key="rev_filter_type")
        with f3:
            rooms_found = set()
            for it in cart:
                alloc = it.get("room_alloc") or {}
                if isinstance(alloc, dict):
                    for room_name, qty in alloc.items():
                        if _room_qty(qty) > 0 and room_name: rooms_found.add(room_name)
            room_opts = ["All rooms"] + sorted(rooms_found)
            sel_room = st.selectbox("Room", room_opts, key="rev_filter_room")

        # --- Filter Items
        filtered_items = []
        for idx, it in enumerate(cart):
            if sel_brand != "All brands" and it.get("Brand") != sel_brand: continue
            if sel_type != "All types" and it.get("Product type") != sel_type: continue
            if sel_room != "All rooms":
                has_room = False
                alloc = it.get("room_alloc") or {}
                if isinstance(alloc, dict):
                    for r, q in alloc.items():
                        if r == sel_room and _room_qty(q) > 0: has_room = True
                if not has_room: continue
            filtered_items.append((idx, it))

        if not filtered_items:
            st.info("No products match filters.")
        else:
            for idx, it in filtered_items:
                uid = _ensure_item_uid(idx)
                with st.container(border=True):
                    left, right = st.columns([0.3, 0.7])
                    
                    # Image
                    with left:
                        comp = (it.get("Composition link") or "").strip()
                        link = (it.get("Link") or "").strip()
                        img = get_image_cached(link, comp)
                        if img: st.image(img, use_container_width=True)
                        else: st.caption("No image")

                    # Details
                    with right:
                        # Title: Product Name (Brand)
                        st.markdown(f"**{it.get('Product name')}** ({it.get('Brand')})")
                        
                        # Line 1: Type & Subtype
                        type_str = it.get("Product type", "")
                        sub_str = it.get("Name of subtype", "")
                        if sub_str:
                            st.text(f"{type_str} • {sub_str}")
                        else:
                            st.text(type_str)

                        # Line 2: Specs (Dimensions & Material)
                        specs = []
                        if it.get("W"): specs.append(f"{it['W']}x{it['D']}x{it['H']}")
                        if it.get("Material"): specs.append(it['Material'])
                        if specs:
                            st.caption(" · ".join(specs))

                        # Line 3: Room & Visual Role
                        rooms_txt = _room_allocation_text(it)
                        # If rooms_txt contains visual roles implicitly or explicitly, show it.
                        # Otherwise, we can reconstruct specific visual roles if needed, 
                        # but _room_allocation_text usually covers the room distribution.
                        # To explicitly show visual role if stored:
                        roles = []
                        alloc = it.get("room_alloc", {})
                        if isinstance(alloc, dict):
                            for r, v in alloc.items():
                                if isinstance(v, dict) and v.get("visual_role"):
                                    roles.append(f"{r}: {v['visual_role']}")
                        
                        if roles:
                            st.caption(f"{rooms_txt}  |  {', '.join(roles)}")
                        else:
                            st.caption(f"{rooms_txt}")
                        
                        # Price (Read-only here, edit in Invoice tab)
                        u_price = _get_item_price(it)
                        total = u_price * int(it.get("qty", 1))
                        st.markdown(f"**Qty:** {it.get('qty')} | **Unit:** ₹{u_price:,.0f} | **Total:** ₹{total:,.0f}")
                        
                        # Notes
                        c_key = f"note_{uid}"
                        curr_note = it.get("internal_comment", "")
                        new_note = st.text_area("Notes", value=curr_note, key=c_key, height=60, label_visibility="collapsed", placeholder="Add client notes...")
                        if new_note != curr_note:
                            st.session_state.cart[idx]["internal_comment"] = new_note
                        
                        # Uploads
                        u_key = _uploader_key_for(uid)
                        _ensure_uploader_cleared(u_key)
                        ups = st.file_uploader("Attach files", type=["jpg","png"], key=u_key, accept_multiple_files=True, label_visibility="collapsed")
                        if ups: _sync_uploader_assets(idx, u_key)
                        
                        assets = _get_assets_for(uid)
                        if assets:
                            st.caption(f"{len(assets)} files attached")

    # ==================================================
    # TAB 2: PRO-FORMA INVOICE (Room Based + Price Edit)
    # ==================================================
    with tab_invoice:
        st.info("ℹ️ **Price Override:** Editing the Unit Price here updates the product globally for this session. It does not affect the database.")
        
        # 1. Group Data by Room
        rooms_data = {} # {Room: [ (cart_idx, qty_in_room) ]}
        
        for idx, it in enumerate(cart):
            qty_total = int(it.get("qty", 1))
            alloc = it.get("room_alloc", {})
            assigned = 0
            
            # Allocated rooms
            if isinstance(alloc, dict):
                for r, v in alloc.items():
                    q = _room_qty(v)
                    if q > 0:
                        if r not in rooms_data: rooms_data[r] = []
                        rooms_data[r].append((idx, q))
                        assigned += q
            
            # Unassigned
            rem = qty_total - assigned
            if rem > 0:
                r = "Unassigned"
                if r not in rooms_data: rooms_data[r] = []
                rooms_data[r].append((idx, rem))

        # 2. Sort Rooms
        sort_keys = ["Drawing Room", "Living Room", "Bedroom 1", "Bedroom 2", "Bedroom 3", "Dining Room", "Kitchen"]
        all_keys = sorted(rooms_data.keys(), key=lambda k: sort_keys.index(k) if k in sort_keys else 999)

        grand_total = 0

        # 3. Render Tables
        for room in all_keys:
            st.subheader(f"{room}")
            
            # [UPDATED] Header row - Added column for Context info (Type/Room/Label)
            h1, h2, h3, h4, h5, h6 = st.columns([0.2, 0.2, 0.2, 0.1, 0.15, 0.15])
            h1.markdown("**Product**")
            h2.markdown("**Context**") # NEW HEADER
            h3.markdown("**Details**")
            h4.markdown("**Qty**")
            h5.markdown("**Unit Price (₹)**")
            h6.markdown("**Total (₹)**")
            st.divider()
            
            room_subtotal = 0
            
            items_list = rooms_data[room]
            for (c_idx, r_qty) in items_list:
                item = cart[c_idx]
                uid = _ensure_item_uid(c_idx)
                
                # [UPDATED] Fetch specific visual role label if available
                visual_label = "—"
                alloc = item.get("room_alloc", {})
                if isinstance(alloc, dict):
                    room_data = alloc.get(room)
                    if isinstance(room_data, dict):
                        visual_label = room_data.get("visual_role", "—")

                with st.container():
                    c1, c2, c3, c4, c5, c6 = st.columns([0.2, 0.2, 0.2, 0.1, 0.15, 0.15])
                    
                    # 1. Product (Brand + Name)
                    c1.caption(item.get("Brand", ""))
                    c1.markdown(f"**{item.get('Product name')}**")
                    
                    # 2. [NEW] Context (Type, Room, Label)
                    c2.caption(f"Type: {item.get('Product type', '')}")
                    c2.text(f"Room: {room}")
                    c2.markdown(f"Label: **{visual_label}**")

                    # 3. Details (Subtype + Dims)
                    desc = []
                    if item.get("Name of subtype"): desc.append(item["Name of subtype"])
                    if item.get("W"): desc.append(f"{item['W']}x{item['D']}x{item['H']}")
                    if desc:
                        for d in desc:
                            c3.text(d)
                    else:
                        c3.text("—")
                    
                    # 4. Qty
                    c4.text(str(r_qty))
                    
                    # 5. Price Editor
                    current_price = _get_item_price(item)
                    new_price = c5.number_input(
                        "Price", 
                        value=float(current_price), 
                        step=100.0, 
                        label_visibility="collapsed",
                        key=f"price_edit_{room}_{uid}"
                    )
                    
                    if new_price != current_price:
                        cart[c_idx]["override_price"] = new_price
                        st.rerun()
                        
                    # 6. Total
                    line_total = new_price * r_qty
                    c6.markdown(f"**{line_total:,.0f}**")
                    
                    room_subtotal += line_total
                
                # [UPDATED] Separator line between products
                st.divider()
            
            # Room Footer
            st.markdown(f"""
            <div style="text-align: right; padding: 10px; background-color: rgba(0,0,0,0.05); border-radius: 5px;">
                <b>{room} Total: ₹ {room_subtotal:,.0f}</b>
            </div>
            """, unsafe_allow_html=True)
            grand_total += room_subtotal
            st.markdown(" ")

        # 4. Grand Total & Invoice Export
        st.markdown("---")
        c_left, c_right = st.columns([0.6, 0.4])
        with c_right:
            st.metric("GRAND TOTAL", f"₹ {grand_total:,.0f}")
            
            if client_name:
                inv_pdf = export_pro_forma_invoice_pdf(cart, client_name, approved_by)
                if inv_pdf:
                    st.download_button(
                        "📄 Download Invoice PDF", 
                        data=inv_pdf, 
                        file_name=f"Invoice_{client_name}_{datetime.now().strftime('%Y%m%d')}.pdf", 
                        mime="application/pdf", 
                        use_container_width=True,
                        type="primary"
                    )
            else:
                st.warning("Enter Client Name above to download Invoice PDF.")

    # --- Footer Navigation ---
    st.divider()
    
    # [UPDATED] Visual Layout Exports (Unified Dropdown UI)
    with st.expander("Visual Layout Exports", expanded=True):
        st.caption("Generate and download visual layout PDFs for each room.")
        
        # Helper to render row: [Generate Button] [Download Button (Disabled until generated)]
        def _render_export_row(label, key_suffix, generator_func, *args):
            c1, c2 = st.columns([0.4, 0.6])
            
            # State key for this specific export
            pdf_key = f"pdf_bytes_{key_suffix}"
            
            with c1:
                if st.button(f"Generate {label}", key=f"gen_{key_suffix}", use_container_width=True):
                    with st.spinner(f"Generating {label}..."):
                        pdf_data = generator_func(*args)
                        if pdf_data:
                            st.session_state[pdf_key] = pdf_data
                        else:
                            st.session_state[pdf_key] = None
                            st.warning(f"No items found for {label}")
            
            with c2:
                data = st.session_state.get(pdf_key)
                if data:
                    st.download_button(
                        f"⬇ Download {label} PDF",
                        data=data,
                        file_name=f"{label.replace(' ', '_')}_Visual.pdf",
                        mime="application/pdf",
                        key=f"dl_{key_suffix}",
                        use_container_width=True,
                        type="primary"
                    )
                else:
                    st.button(
                        f"⬇ Download {label} PDF", 
                        disabled=True, 
                        key=f"dl_dis_{key_suffix}", 
                        use_container_width=True
                    )

        # 1. Export All (Merged)
        _render_export_row("All Rooms Merged", "all", export_all_rooms_pdf, st.session_state.cart, client_name, approved_by)
        st.markdown("---")

        # 2. Individual Rooms
        _render_export_row("Drawing Room", "dr", export_drawing_room_pdf, st.session_state.cart, client_name, approved_by)
        _render_export_row("Living Room", "lr", export_living_room_pdf, st.session_state.cart, client_name, approved_by)
        _render_export_row("Dining Room", "din", export_dining_room_pdf, st.session_state.cart, client_name, approved_by)
        _render_export_row("Bedroom 1", "br1", export_bedroom_pdf, st.session_state.cart, "Bedroom 1", client_name, approved_by)
        _render_export_row("Bedroom 2", "br2", export_bedroom_pdf, st.session_state.cart, "Bedroom 2", client_name, approved_by)
        _render_export_row("Bedroom 3", "br3", export_bedroom_pdf, st.session_state.cart, "Bedroom 3", client_name, approved_by)
        _render_export_row("Others / General", "oth", export_others_room_pdf, st.session_state.cart, client_name, approved_by)

    st.button("⬅ Back to Products", on_click=_back_from_review, key="rev_back_bottom")